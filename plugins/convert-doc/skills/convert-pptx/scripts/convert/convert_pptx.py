#!/usr/bin/env python3
"""
convert_pptx.py - Markdown to PowerPoint (PPTX) converter

Usage:
  python convert_pptx.py <input.md> [output.pptx]
    [--title TITLE] [--subtitle SUB] [--aspect 16:9|4:3]
    [--primary-color "#003879"] [--max-body-chars 2400]

Dependencies: python-pptx, Pillow, requests
"""

import argparse
import base64
import io
import json
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, List, Optional
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# --- Design constants (edit here for a re-skinning) ---
PRIMARY_DEFAULT = "#003879"   # primary navy
ACCENT = RGBColor(0x1D, 0x6F, 0xD1)
TEXT = RGBColor(0x1F, 0x2D, 0x3D)
CODE_BG = RGBColor(0xF5, 0xF6, 0xF8)
CODE_TEXT = RGBColor(0x1F, 0x2D, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BODY_FONT = "Yu Gothic UI"
HEADING_FONT = "Yu Gothic UI"
CODE_FONT = "Consolas"

TITLE_BAND_HEIGHT_IN = 0.9
CONTENT_PADDING_IN = 0.5
MERMAID_MAX_WIDTH_IN = 11.5
MERMAID_MAX_HEIGHT_IN = 5.5
IMAGE_MAX_WIDTH_IN = 11.5
IMAGE_MAX_HEIGHT_IN = 5.5

MERMAID_PNG_ENDPOINT = "https://mermaid.ink/img/{}?type=png"

# ----- Pygments syntax highlight palette (aligned with the HTML friendly theme) -----
_SYNTAX_PALETTE = {
    "keyword":   RGBColor(0x00, 0x7B, 0x83),   # Token.Keyword*
    "builtin":   RGBColor(0x5C, 0x35, 0x66),   # Token.Name.Builtin*
    "func":      RGBColor(0x00, 0x4E, 0xB0),   # Token.Name.Function / Token.Name.Class
    "tag":       RGBColor(0x00, 0x7B, 0x83),   # Token.Name.Tag
    "attr":      RGBColor(0xAA, 0x55, 0x00),   # Token.Name.Attribute
    "string":    RGBColor(0x4E, 0x95, 0x2A),   # Token.Literal.String*
    "number":    RGBColor(0xAA, 0x55, 0x00),   # Token.Literal.Number*
    "operator":  RGBColor(0x33, 0x33, 0x33),
    "comment":   RGBColor(0x80, 0x80, 0x80),   # Token.Comment*
    "error":     RGBColor(0xB7, 0x25, 0x25),
    "heading":   RGBColor(0x00, 0x4E, 0xB0),   # Token.Generic.Heading*
}


# ===================== Markdown parsing =====================

@dataclass
class Block:
    """A Markdown block element."""
    kind: str
    text: str = ""
    level: int = 0
    lang: str = ""
    rows: List[List[str]] = field(default_factory=list)
    src: str = ""
    alt: str = ""


HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
IMAGE_RE = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$")
UL_RE = re.compile(r"^(\s*)[-*+]\s+(.+)$")
OL_RE = re.compile(r"^(\s*)\d+\.\s+(.+)$")
HR_RE = re.compile(r"^\s*(-{3,}|\*{3,}|_{3,})\s*$")
TABLE_DIVIDER_RE = re.compile(r"^\s*\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)+\|?\s*$")


def parse_markdown(md_text: str) -> List[Block]:
    """Parse the Markdown text into a flat list of Block objects.
    Minimal block-level parser: covers headings, paragraphs, bullet/numbered lists,
    fenced code (including mermaid), tables, images (standalone), and horizontal rules.
    """
    lines = md_text.splitlines()
    blocks: List[Block] = []
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]

        # Fenced code
        fence = re.match(r"^```(\w*)\s*$", line)
        if fence:
            lang = fence.group(1)
            i += 1
            buf = []
            while i < n and not re.match(r"^```\s*$", lines[i]):
                buf.append(lines[i])
                i += 1
            if i < n:
                i += 1  # consume closing fence
            kind = "mermaid" if lang == "mermaid" else "code"
            blocks.append(Block(kind=kind, text="\n".join(buf), lang=lang))
            continue

        # Heading
        m = HEADING_RE.match(line)
        if m:
            blocks.append(Block(kind="heading", text=m.group(2), level=len(m.group(1))))
            i += 1
            continue

        # Standalone image
        m = IMAGE_RE.match(line)
        if m:
            blocks.append(Block(kind="image", alt=m.group(1), src=m.group(2)))
            i += 1
            continue

        # Horizontal rule
        if HR_RE.match(line):
            blocks.append(Block(kind="hr"))
            i += 1
            continue

        # Table: header | divider | body rows
        if "|" in line and i + 1 < n and TABLE_DIVIDER_RE.match(lines[i + 1]):
            header = _split_table_row(line)
            i += 2
            rows = [header]
            while i < n and "|" in lines[i] and lines[i].strip():
                rows.append(_split_table_row(lines[i]))
                i += 1
            blocks.append(Block(kind="table", rows=rows))
            continue

        # Lists (group consecutive items)
        m_ul = UL_RE.match(line)
        m_ol = OL_RE.match(line)
        if m_ul or m_ol:
            kind = "ul" if m_ul else "ol"
            items = []
            while i < n:
                mm_ul = UL_RE.match(lines[i])
                mm_ol = OL_RE.match(lines[i])
                if mm_ul and kind == "ul":
                    indent = len(mm_ul.group(1))
                    items.append((indent // 2, mm_ul.group(2)))
                    i += 1
                elif mm_ol and kind == "ol":
                    indent = len(mm_ol.group(1))
                    items.append((indent // 2, mm_ol.group(2)))
                    i += 1
                else:
                    break
            blocks.append(Block(kind=kind, rows=[[str(lvl), txt] for lvl, txt in items]))
            continue

        # Blank line -> paragraph separator
        if not line.strip():
            i += 1
            continue

        # Paragraph (accumulate consecutive non-blank lines)
        para_lines = [line]
        i += 1
        while (
            i < n
            and lines[i].strip()
            and not HEADING_RE.match(lines[i])
            and not re.match(r"^```", lines[i])
            and not IMAGE_RE.match(lines[i])
            and not HR_RE.match(lines[i])
            and not UL_RE.match(lines[i])
            and not OL_RE.match(lines[i])
            and not ("|" in lines[i] and i + 1 < n and TABLE_DIVIDER_RE.match(lines[i + 1]))
        ):
            para_lines.append(lines[i])
            i += 1
        blocks.append(Block(kind="paragraph", text=" ".join(s.strip() for s in para_lines)))

    return blocks


def _split_table_row(line: str) -> List[str]:
    """Split a Markdown table row into cells."""
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


# ===================== Slide segmentation =====================

@dataclass
class SlideSpec:
    title: Optional[str]
    blocks: List[Block]


def split_into_slides(blocks: List[Block]) -> (Optional[str], Optional[str], List[SlideSpec]):
    """Produce (doc_title, doc_subtitle, slides) from a flat block list.
    - First H1 -> doc_title (becomes title slide)
    - The paragraph immediately after the first H1 (if any) -> subtitle
    - Each H2 starts a new slide; its text becomes the slide title
    - If there's no H2 at all, a single content slide holds everything after the H1
    """
    doc_title: Optional[str] = None
    doc_subtitle: Optional[str] = None
    slides: List[SlideSpec] = []

    idx = 0
    n = len(blocks)

    # Find first H1
    while idx < n and not (blocks[idx].kind == "heading" and blocks[idx].level == 1):
        idx += 1
    if idx < n:
        doc_title = blocks[idx].text
        idx += 1
        # Optional subtitle: the first paragraph before any H2
        j = idx
        while j < n:
            b = blocks[j]
            if b.kind == "heading" and b.level == 2:
                break
            if b.kind == "paragraph" and doc_subtitle is None:
                doc_subtitle = b.text
                break
            j += 1
    else:
        idx = 0  # no H1 at all; still try to build slides

    # Segment remaining blocks by H2
    current: Optional[SlideSpec] = None
    any_h2 = False
    for b in blocks[idx:]:
        if b.kind == "heading" and b.level == 2:
            any_h2 = True
            current = SlideSpec(title=b.text, blocks=[])
            slides.append(current)
            continue
        if current is None:
            # Before the first H2; skip the subtitle paragraph we already captured
            if b.kind == "paragraph" and b.text == doc_subtitle:
                continue
            current = SlideSpec(title=None, blocks=[])
            slides.append(current)
        current.blocks.append(b)

    if not any_h2 and not slides:
        # Rare case: only H1 (or nothing). Create an empty slide to avoid zero-slide decks.
        slides.append(SlideSpec(title=None, blocks=[]))

    return doc_title, doc_subtitle, slides


# ===================== Rendering helpers =====================

def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def fetch_mermaid_png(code: str) -> Optional[bytes]:
    """Download a mermaid diagram as PNG via mermaid.ink. Returns None on failure."""
    import requests
    try:
        payload = {"code": code, "mermaid": {"theme": "default"}}
        encoded = base64.urlsafe_b64encode(json.dumps(payload).encode()).decode().rstrip("=")
        url = MERMAID_PNG_ENDPOINT.format(encoded)
        resp = requests.get(url, timeout=20, headers={"User-Agent": "convert-pptx/1.0"})
        if resp.status_code == 200 and resp.content.startswith(b"\x89PNG"):
            return resp.content
        print(f"Warning: mermaid.ink returned status={resp.status_code}", file=sys.stderr)
    except Exception as e:
        print(f"Warning: mermaid.ink fetch failed: {e}", file=sys.stderr)
    return None


def fit_size_inches(px_w: int, px_h: int, max_w_in: float, max_h_in: float) -> (Inches, Inches):
    """Fit a (pixel) size within (max_w, max_h) inches preserving aspect ratio."""
    # Assume 96 DPI for raster images coming from mermaid.ink
    w_in = px_w / 96
    h_in = px_h / 96
    if w_in <= 0 or h_in <= 0:
        return Inches(max_w_in), Inches(max_h_in)
    ratio = min(max_w_in / w_in, max_h_in / h_in, 1.0)
    return Inches(w_in * ratio), Inches(h_in * ratio)


def measure_image_bytes(data: bytes) -> (int, int):
    from PIL import Image
    with Image.open(io.BytesIO(data)) as img:
        return img.size


def _lexer_for(lang: str):
    """Return a pygments lexer for the given hint, falling back to plain text."""
    from pygments.lexers import get_lexer_by_name
    from pygments.lexers.special import TextLexer
    from pygments.util import ClassNotFound

    if lang:
        try:
            return get_lexer_by_name(lang, stripnl=False, ensurenl=False)
        except ClassNotFound:
            pass
    return TextLexer(stripnl=False, ensurenl=False)


def _token_color(tok_type):
    """Resolve a pygments token to an RGB color by walking up the token hierarchy."""
    from pygments.token import Token

    mapping = [
        (Token.Keyword,             "keyword"),
        (Token.Name.Function,       "func"),
        (Token.Name.Class,          "func"),
        (Token.Name.Decorator,      "func"),
        (Token.Name.Builtin,        "builtin"),
        (Token.Name.Builtin.Pseudo, "builtin"),
        (Token.Name.Tag,            "tag"),
        (Token.Name.Attribute,      "attr"),
        (Token.Literal.String,      "string"),
        (Token.Literal.Number,      "number"),
        (Token.Operator,            "operator"),
        (Token.Punctuation,         "operator"),
        (Token.Comment,             "comment"),
        (Token.Error,               "error"),
        (Token.Generic.Heading,     "heading"),
        (Token.Generic.Subheading,  "heading"),
    ]
    t = tok_type
    while t is not None:
        for base, key in mapping:
            if t in base:
                return _SYNTAX_PALETTE[key]
        parent = getattr(t, "parent", None)
        if parent is None or parent is t:
            break
        t = parent
    return None


def _lex_code_lines(code_text: str, lang: str):
    """Yield token-list-per-line for the given code string.
    Each item is a list of (token_type, text_fragment) tuples.  Empty lines
    yield an empty list so the paragraph structure matches the source.
    """
    from pygments import lex

    tokens = lex(code_text, _lexer_for(lang))
    line: list = []
    for tok, text in tokens:
        while "\n" in text:
            head, _, text = text.partition("\n")
            if head:
                line.append((tok, head))
            yield line
            line = []
        if text:
            line.append((tok, text))
    if line:
        yield line


NBSP = " "

def _normalize_leading(line_tokens):
    """Convert leading-whitespace-only runs to NBSP so PowerPoint preserves indentation.

    PowerPoint collapses/hides runs of ASCII spaces during layout, so we replace
    the spaces inside leading-whitespace-only runs with NBSP. Tabs are expanded
    to 4 spaces everywhere. Non-leading spaces are left as normal spaces so word
    wrap behaves naturally.
    """
    normalized = []
    indent_done = False
    for tok, text in line_tokens:
        if not indent_done and text and text.strip() == "":
            normalized.append((tok, text.replace("	", "    ").replace(" ", NBSP)))
        else:
            indent_done = True
            normalized.append((tok, text.replace("	", "    ")))
    return normalized
def _block_has_content(block: "Block") -> bool:
    """Return True if a block would render any visible shape on a slide."""
    if block.kind == "paragraph":
        return bool(block.text and block.text.strip())
    if block.kind in ("ul", "ol"):
        return bool(block.rows) and any((r[1] or "").strip() for r in block.rows)
    if block.kind == "heading":
        return bool(block.text and block.text.strip())
    if block.kind == "table":
        return bool(block.rows) and any(any((c or "").strip() for c in row) for row in block.rows)
    if block.kind == "code":
        return bool(block.text and block.text.strip())
    if block.kind == "mermaid":
        return bool(block.text and block.text.strip())
    if block.kind == "image":
        return bool(block.src)
    if block.kind == "hr":
        return False  # horizontal rule on its own does not justify a slide
    return False


def strip_inline_markdown(text: str) -> str:
    """Minimal inline cleanup for display in text frames."""
    # GFM task list markers -> unicode ballot boxes. Do this *before* we strip
    # links so the brackets are not confused with link syntax.
    text = re.sub(r"^\s*\[[ xX]\]\s+", lambda m: ("☑ " if "x" in m.group(0).lower() else "☐ "), text)
    # Links: [text](url) -> text
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
    # Inline code
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # Bold / italic markers (we keep content; runs are not split here)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    # Strikethrough
    text = re.sub(r"~~([^~]+)~~", r"\1", text)
    return text


# ===================== PPTX builder =====================

class Deck:
    def __init__(self, *, aspect: str, primary: RGBColor):
        self.prs = Presentation()
        if aspect == "4:3":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
        else:  # 16:9
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.primary = primary

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self.blank)
        return slide

    def _add_title_band(self, slide, text: str):
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(TITLE_BAND_HEIGHT_IN),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = self.primary
        band.line.fill.background()
        tf = band.text_frame
        tf.margin_left = Inches(0.35)
        tf.margin_right = Inches(0.35)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = text
        p = tf.paragraphs[0]
        p.font.name = HEADING_FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

    def add_title_slide(self, title: str, subtitle: Optional[str]):
        slide = self._new_slide()
        # Full-bleed primary block on the left
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.4), self.prs.slide_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.primary
        bar.line.fill.background()

        title_tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.3),
            self.prs.slide_width - Inches(1.5), Inches(2.0),
        )
        tf = title_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = HEADING_FONT
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.primary

        if subtitle:
            sub_tb = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.4),
                self.prs.slide_width - Inches(1.5), Inches(1.5),
            )
            tf = sub_tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = strip_inline_markdown(subtitle)
            p.font.name = BODY_FONT
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT

    def add_content_slide(
        self,
        spec: SlideSpec,
        base_dir: Path,
        max_body_chars: int,
    ):
        # Drop blocks that would render as nothing (empty paragraphs, blank hr
        # lists produced by upstream quirks) so we don't emit empty slides.
        nonempty_blocks = [b for b in spec.blocks if _block_has_content(b)]
        if not nonempty_blocks and not spec.title:
            # Truly empty slide -> skip entirely.
            return
        chunks = self._chunk_blocks(nonempty_blocks, max_body_chars)
        # Drop empty chunks defensively. _chunk_blocks returns [[]] for empty
        # input; a slide whose only purpose was the title band is still useful,
        # but a chunk with no blocks AND no title should not create a slide.
        chunks = [c for c in chunks if c or spec.title]
        if not chunks:
            return
        for idx, chunk in enumerate(chunks):
            slide = self._new_slide()
            if spec.title:
                title = spec.title if idx == 0 else f"{spec.title} ({idx + 1})"
                self._add_title_band(slide, title)
                top_cursor = Inches(TITLE_BAND_HEIGHT_IN + 0.2)
            else:
                top_cursor = Inches(0.3)
            self._render_blocks(slide, chunk, top_cursor, base_dir)

    @staticmethod
    def _estimate_block_height_in(b: Block) -> float:
        """Estimate the rendered height of a block in inches.

        The numbers intentionally mirror the actual geometry used by the
        _render_* methods plus the 0.1in gap added after each block, so the
        chunker's budget stays in sync with what _render_blocks will actually
        consume. A small over-estimate is preferable to under-estimate because
        it avoids overflow cutoffs that silently drop shapes.
        """
        gap = 0.1
        if b.kind == "heading":
            return 0.45 + gap
        if b.kind == "paragraph":
            approx_lines = max(1, len(b.text) // 60 + 1)
            return 0.35 * approx_lines + gap
        if b.kind in ("ul", "ol"):
            n = max(1, len(b.rows))
            return 0.4 * n + 0.1 + gap
        if b.kind == "code":
            n_lines = max(1, b.text.count("\n") + 1)
            return max(0.5, 0.28 * n_lines + 0.2) + gap
        if b.kind == "table":
            n_rows = max(1, len(b.rows))
            return 0.4 * n_rows + gap
        if b.kind == "mermaid":
            return MERMAID_MAX_HEIGHT_IN + gap
        if b.kind == "image":
            return IMAGE_MAX_HEIGHT_IN + gap
        if b.kind == "hr":
            return 0.15 + gap
        return 0.3 + gap

    @classmethod
    def _chunk_blocks(cls, blocks: List[Block], max_chars: int) -> List[List[Block]]:
        """Chunk the block list into slide-sized groups based on estimated height.

        The `max_chars` parameter is preserved for API compatibility but is now
        used only as a safety cap: the primary budget is the slide's remaining
        vertical space (slide_height - title band - top/bottom padding - a
        small safety margin). This keeps the chunker's splits aligned with the
        renderer's vertical layout, so shapes aren't silently dropped by the
        bottom-limit cutoff in _render_blocks.

        After packing, an orphan-heading pass moves trailing headings
        (e.g. `### 7.1 ...`) to the next chunk so a heading never sits alone at
        the tail of a slide while its content starts on the next one.
        """
        # Vertical budget per slide (inches). Assumes 7.5in slide height.
        budget_in = 7.5 - TITLE_BAND_HEIGHT_IN - CONTENT_PADDING_IN - 0.6

        chunks: List[List[Block]] = []
        current: List[Block] = []
        height_used = 0.0
        chars_used = 0

        for b in blocks:
            h = cls._estimate_block_height_in(b)
            chars = max(1, len(b.text or "")) + sum(len(x[1] or "") + 4 for x in (b.rows or []))

            would_overflow_height = current and (height_used + h > budget_in)
            would_overflow_chars = current and (chars_used + chars > max_chars)
            if would_overflow_height or would_overflow_chars:
                chunks.append(current)
                current = []
                height_used = 0.0
                chars_used = 0

            current.append(b)
            height_used += h
            chars_used += chars

        if current:
            chunks.append(current)

        # Orphan-heading fix: a heading at the very end of a chunk (with its
        # actual content spilling into the next chunk) reads as a dangling
        # title. Move it forward to sit with the content it introduces.
        for i in range(len(chunks) - 1):
            moved: List[Block] = []
            while chunks[i] and chunks[i][-1].kind == "heading":
                moved.append(chunks[i].pop())
            if moved:
                chunks[i + 1] = list(reversed(moved)) + chunks[i + 1]

        # Drop any chunks that became empty after the orphan migration.
        chunks = [c for c in chunks if c]
        if not chunks:
            chunks.append([])
        return chunks

    def _render_blocks(self, slide, blocks: List[Block], start_top: Emu, base_dir: Path):
        left = Inches(CONTENT_PADDING_IN)
        width = self.prs.slide_width - Inches(CONTENT_PADDING_IN * 2)
        bottom_limit = self.prs.slide_height - Inches(0.3)
        cursor = start_top

        for b in blocks:
            if cursor >= bottom_limit:
                break

            if b.kind == "paragraph":
                height = self._render_text(
                    slide, left, cursor, width, strip_inline_markdown(b.text), Pt(16)
                )
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "heading":
                # H3+ rendered as bold accent line
                size = max(14, 22 - (b.level - 3) * 2)
                height = self._render_text(
                    slide, left, cursor, width,
                    strip_inline_markdown(b.text),
                    Pt(size), bold=True, color=self.primary,
                )
                cursor = cursor + height + Inches(0.1)

            elif b.kind in ("ul", "ol"):
                items = b.rows
                tb = slide.shapes.add_textbox(left, cursor, width, Inches(0.4 * len(items) + 0.1))
                tf = tb.text_frame
                tf.word_wrap = True
                for i, (lvl_str, txt) in enumerate(items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    bullet = "- " if b.kind == "ul" else f"{i + 1}. "
                    p.text = f"{'  ' * int(lvl_str)}{bullet}{strip_inline_markdown(txt)}"
                    p.font.name = BODY_FONT
                    p.font.size = Pt(15)
                    p.font.color.rgb = TEXT
                cursor = cursor + tb.height + Inches(0.1)

            elif b.kind == "code":
                height = self._render_code(slide, left, cursor, width, b.text, b.lang)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "mermaid":
                height = self._render_mermaid(slide, left, cursor, width, b.text)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "image":
                height = self._render_image(slide, left, cursor, width, b.src, b.alt, base_dir)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "table":
                height = self._render_table(slide, left, cursor, width, b.rows)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "hr":
                ln = slide.shapes.add_connector(1, left, cursor, left + width, cursor)
                ln.line.color.rgb = RGBColor(0xC9, 0xD0, 0xD8)
                cursor = cursor + Inches(0.15)

    def _render_text(
        self, slide, left, top, width, text, size: Pt, *, bold=False, color: RGBColor = TEXT
    ):
        # Rough height estimate based on char count
        approx_lines = max(1, int(len(text) / 60) + 1)
        height = Inches(0.35 * approx_lines)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = BODY_FONT
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        return height

    def _render_code(self, slide, left, top, width, code_text: str, lang: str = ""):
        """Render a fenced code block with Pygments syntax highlight.

        Each line becomes one paragraph; each Pygments token becomes a run so
        the whole block keeps a uniform monospace font while coloring keywords,
        strings, comments, etc.  Leading whitespace on each line is converted
        to NBSP so PowerPoint preserves indentation.
        """
        line_groups = list(_lex_code_lines(code_text, lang))
        if not line_groups:
            line_groups = [[]]

        height = Inches(max(0.5, 0.28 * len(line_groups) + 0.2))
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = CODE_BG
        box.line.color.rgb = RGBColor(0xDD, 0xE1, 0xE8)
        tf = box.text_frame
        # Code should not reflow: it preserves the author's line breaks.
        tf.word_wrap = False
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)

        for i, line_tokens in enumerate(line_groups):
            line_tokens = _normalize_leading(line_tokens)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Baseline paragraph formatting (applies if a run doesn't override).
            p.font.name = CODE_FONT
            p.font.size = Pt(11)
            p.font.color.rgb = CODE_TEXT

            if not line_tokens:
                # Preserve blank lines with a NBSP run so the line height stays.
                run = p.add_run()
                run.text = NBSP
                run.font.name = CODE_FONT
                run.font.size = Pt(11)
                run.font.color.rgb = CODE_TEXT
                continue

            for tok, txt in line_tokens:
                if not txt:
                    continue
                run = p.add_run()
                run.text = txt
                run.font.name = CODE_FONT
                run.font.size = Pt(11)
                color = _token_color(tok) or CODE_TEXT
                run.font.color.rgb = color
        return height

    def _render_mermaid(self, slide, left, top, width, code: str):
        png = fetch_mermaid_png(code)
        if not png:
            return self._render_code(slide, left, top, width, code)
        try:
            w, h = measure_image_bytes(png)
            img_w, img_h = fit_size_inches(w, h, MERMAID_MAX_WIDTH_IN, MERMAID_MAX_HEIGHT_IN)
        except Exception:
            img_w, img_h = Inches(MERMAID_MAX_WIDTH_IN), Inches(MERMAID_MAX_HEIGHT_IN)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=img_w, height=img_h)
        return img_h

    def _render_image(self, slide, left, top, width, src: str, alt: str, base_dir: Path):
        data = self._load_image_bytes(src, base_dir)
        if data is None:
            return self._render_text(
                slide, left, top, width, f"[画像が見つかりません: {alt or src}]", Pt(13)
            )
        try:
            w, h = measure_image_bytes(data)
            img_w, img_h = fit_size_inches(w, h, IMAGE_MAX_WIDTH_IN, IMAGE_MAX_HEIGHT_IN)
        except Exception:
            img_w, img_h = Inches(IMAGE_MAX_WIDTH_IN), Inches(IMAGE_MAX_HEIGHT_IN)
        slide.shapes.add_picture(io.BytesIO(data), left, top, width=img_w, height=img_h)
        return img_h

    @staticmethod
    def _load_image_bytes(src: str, base_dir: Path) -> Optional[bytes]:
        parsed = urlparse(src)
        if parsed.scheme in ("http", "https"):
            try:
                import requests
                r = requests.get(src, timeout=20, headers={"User-Agent": "convert-pptx/1.0"})
                if r.status_code == 200:
                    return r.content
            except Exception:
                return None
            return None
        # local path
        p = Path(src)
        if not p.is_absolute():
            p = base_dir / src
        if p.exists():
            return p.read_bytes()
        return None

    def _render_table(self, slide, left, top, width, rows: List[List[str]]):
        if not rows:
            return Inches(0.3)
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        row_h = Inches(0.4)
        height = Inches(0.4 * n_rows)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table
        for r, row in enumerate(rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.text = strip_inline_markdown(row[c]) if c < len(row) else ""
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = BODY_FONT
                        run.font.size = Pt(12)
                        if r == 0:
                            run.font.bold = True
                            run.font.color.rgb = WHITE
                        else:
                            run.font.color.rgb = TEXT
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.primary
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 1 else RGBColor(0xF5, 0xF6, 0xF8)
        return height

    def save(self, path: Path):
        self.prs.save(str(path))


# ===================== Entry point =====================

def main() -> None:
    parser = argparse.ArgumentParser(description="Convert Markdown to PowerPoint (PPTX)")
    parser.add_argument("input", help="Input .md file path")
    parser.add_argument("output", nargs="?", help="Output .pptx file path")
    parser.add_argument("--title", help="Title slide heading (default: first H1 in the Markdown)")
    parser.add_argument("--subtitle", help="Title slide subtitle")
    parser.add_argument("--aspect", default="16:9", choices=["16:9", "4:3"])
    parser.add_argument("--primary-color", default=PRIMARY_DEFAULT, help="Primary color #RRGGBB")
    parser.add_argument("--max-body-chars", type=int, default=2400,
                        help="Soft budget for a single slide's body text")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output) if args.output else input_path.with_suffix(".pptx")
    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    md_text = input_path.read_text(encoding="utf-8")
    blocks = parse_markdown(md_text)
    doc_title, doc_subtitle, slides = split_into_slides(blocks)

    if args.title:
        doc_title = args.title
    if args.subtitle:
        doc_subtitle = args.subtitle

    primary = hex_to_rgb(args.primary_color)
    deck = Deck(aspect=args.aspect, primary=primary)
    if doc_title:
        deck.add_title_slide(doc_title, doc_subtitle)
    for spec in slides:
        deck.add_content_slide(spec, base_dir=input_path.parent, max_body_chars=args.max_body_chars)

    deck.save(output_path)
    print(f"Generated: {output_path}")


if __name__ == "__main__":
    main()
