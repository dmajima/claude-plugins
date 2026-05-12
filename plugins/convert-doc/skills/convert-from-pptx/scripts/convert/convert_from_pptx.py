#!/usr/bin/env python3
"""convert_from_pptx.py - PPTX を Markdown に変換するスクリプト.

使い方:
    python convert_from_pptx.py <入力PPTX> [<出力MD>] [--images-dir DIR]
                                [--no-mermaid] [--include-notes]
                                [--include-hidden] [--no-first-slide-as-title]
                                [--max-image-size BYTES]
"""

from __future__ import annotations

import argparse
import re
import sys
import zipfile
from pathlib import Path
from typing import Iterable, List, Optional

try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:
    print(f"Error: python-pptx is not installed: {exc}", file=sys.stderr)
    sys.exit(2)

try:
    from lxml import etree
except ImportError as exc:
    print(f"Error: lxml is not installed: {exc}", file=sys.stderr)
    sys.exit(2)


# ----------------------------------------------------------------------------- #
# 定数
# ----------------------------------------------------------------------------- #

MONOSPACE_FONTS = {
    "Consolas",
    "Courier",
    "Courier New",
    "Menlo",
    "Monaco",
    "Cascadia Code",
    "Cascadia Mono",
    "Fira Code",
    "Source Code Pro",
    "Lucida Console",
    "MS Gothic",
    "ＭＳ ゴシック",
    "MS Mincho",
    "ＭＳ 明朝",
}

PPTX_MAGIC = b"PK\x03\x04"
DEFAULT_MAX_IMAGE_BYTES = 5 * 1024 * 1024  # 5 MiB
MAX_TOTAL_UNCOMPRESSED_BYTES = 256 * 1024 * 1024  # 256 MiB (ZIP bomb 防御)
MAX_COMPRESSION_RATIO = 200  # compress_size に対する展開サイズ倍率上限

ALLOWED_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp", "emf", "wmf"}

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _hardened_xml_parser() -> "etree.XMLParser":
    """XXE / DTD / 巨大ツリー攻撃を遮断した lxml パーサを返す."""
    return etree.XMLParser(
        resolve_entities=False,
        no_network=True,
        load_dtd=False,
        huge_tree=False,
    )


# ----------------------------------------------------------------------------- #
# ユーティリティ
# ----------------------------------------------------------------------------- #


def _safe_text(value: Optional[str]) -> str:
    if not value:
        return ""
    return value.replace("\r", "").strip()


def _escape_md_pipe(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", "<br>")


def _is_monospace_font(font_name: Optional[str]) -> bool:
    if not font_name:
        return False
    return font_name in MONOSPACE_FONTS


def _mermaid_escape_label(text: str) -> str:
    text = (text or "").replace('"', "'").replace("\n", "<br/>")
    # Mermaid のメタ文字をエスケープ（後段 HTML 化時の XSS 対策含む）
    text = text.replace("<", "&lt;").replace(">", "&gt;")
    if len(text) > 80:
        text = text[:77] + "..."
    return text or " "


def _escape_md_link_text(text: str) -> str:
    """Markdown のリンクテキスト・alt 部分のメタ文字をエスケープ."""
    return (
        (text or "")
        .replace("\\", "\\\\")
        .replace("[", "\\[")
        .replace("]", "\\]")
        .replace("(", "\\(")
        .replace(")", "\\)")
        .replace("\n", " ")
    )


def _validate_pptx(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(f"Input file not found: {path}")
    if path.suffix.lower() not in (".pptx", ".pptm"):
        raise ValueError(f"Input file is not PPTX/PPTM (extension): {path}")
    with open(path, "rb") as fh:
        magic = fh.read(4)
    if magic != PPTX_MAGIC:
        raise ValueError(f"Input file is not a valid PPTX (zip magic): {path}")
    try:
        with zipfile.ZipFile(path, "r") as zf:
            # ZIP bomb 検査: 総展開サイズ・圧縮率
            total_uncompressed = 0
            for info in zf.infolist():
                total_uncompressed += info.file_size
                if total_uncompressed > MAX_TOTAL_UNCOMPRESSED_BYTES:
                    raise ValueError(
                        f"Input PPTX exceeds uncompressed size limit "
                        f"({total_uncompressed} > {MAX_TOTAL_UNCOMPRESSED_BYTES}): {path}"
                    )
                if info.compress_size > 0:
                    ratio = info.file_size / info.compress_size
                    if ratio > MAX_COMPRESSION_RATIO:
                        raise ValueError(
                            f"Suspicious compression ratio ({ratio:.1f}x) "
                            f"for entry '{info.filename}' in {path}"
                        )
            content_types = zf.read("[Content_Types].xml").decode("utf-8", errors="ignore")
    except KeyError as exc:
        raise ValueError(f"Input file missing [Content_Types].xml: {path}") from exc
    if "presentationml" not in content_types.lower():
        raise ValueError(f"Input file is not a PresentationML document: {path}")


def _resolve_images_dir(output_md: Path, images_dir_opt: Optional[str]) -> Path:
    base = output_md.parent.resolve()
    if images_dir_opt:
        candidate = Path(images_dir_opt)
        if not candidate.is_absolute():
            candidate = output_md.parent / candidate
        candidate = candidate.resolve()
        try:
            candidate.relative_to(base)
        except ValueError as exc:
            raise ValueError(
                f"images-dir must be under output MD directory (path traversal blocked): {candidate}"
            ) from exc
        return candidate
    return (base / f"{output_md.stem}_images").resolve()


# ----------------------------------------------------------------------------- #
# Mermaid 生成（コネクタ駆動）
# ----------------------------------------------------------------------------- #


class FlowExtractor:
    """同一スライド内の図形群とコネクタからフロー図 Mermaid を生成する."""

    def __init__(self, shapes_meta: List[dict], connectors: List[dict]):
        self.shapes_meta = shapes_meta
        self.connectors = connectors
        self.used_node_ids: set = set()

    def build_mermaid(self) -> Optional[str]:
        self.used_node_ids = set()
        if not self.connectors or len(self.shapes_meta) < 2:
            return None
        id_to_label = {
            meta["shape_id"]: (meta["text"] or meta["name"] or f"shape{meta['shape_id']}")
            for meta in self.shapes_meta
        }
        edges = []
        for connector in self.connectors:
            begin = connector.get("begin")
            end = connector.get("end")
            if begin is None or end is None:
                continue
            if begin not in id_to_label or end not in id_to_label:
                continue
            edges.append((begin, end))
            self.used_node_ids.add(begin)
            self.used_node_ids.add(end)
        if not edges:
            return None
        direction = self._infer_direction()
        lines = [f"flowchart {direction}"]
        for nid in sorted(self.used_node_ids):
            meta = next((m for m in self.shapes_meta if m["shape_id"] == nid), None)
            label = _mermaid_escape_label(id_to_label.get(nid, f"shape{nid}"))
            open_b, close_b = self._bracket_for_shape(meta)
            lines.append(f'    N{nid}{open_b}"{label}"{close_b}')
        for src, dst in edges:
            lines.append(f"    N{src} --> N{dst}")
        return "\n".join(lines)

    def _infer_direction(self) -> str:
        xs = [m["left"] for m in self.shapes_meta if m.get("left") is not None]
        ys = [m["top"] for m in self.shapes_meta if m.get("top") is not None]
        if not xs or not ys:
            return "TD"
        x_range = max(xs) - min(xs)
        y_range = max(ys) - min(ys)
        if x_range > y_range * 1.5:
            return "LR"
        return "TD"

    @staticmethod
    def _bracket_for_shape(meta: Optional[dict]) -> tuple:
        if not meta:
            return ("[", "]")
        auto = (meta.get("auto_shape_type") or "").upper()
        if "OVAL" in auto or "ELLIPSE" in auto or "FLOWCHART_TERMINATOR" in auto:
            return ("((", "))")
        if "DIAMOND" in auto or "FLOWCHART_DECISION" in auto:
            return ("{", "}")
        if "PARALLELOGRAM" in auto or "FLOWCHART_DATA" in auto:
            return ("[/", "/]")
        return ("[", "]")


# ----------------------------------------------------------------------------- #
# 変換本体
# ----------------------------------------------------------------------------- #


class PPTXMarkdownConverter:
    """python-pptx で読んだプレゼンを Markdown に転記する."""

    def __init__(self, args: argparse.Namespace) -> None:
        self.input_path = Path(args.input).resolve()
        self.output_path = (
            Path(args.output).resolve() if args.output else self.input_path.with_suffix(".md")
        )
        self.images_dir = _resolve_images_dir(self.output_path, args.images_dir)
        self.no_mermaid = bool(args.no_mermaid)
        self.include_notes = bool(args.include_notes)
        self.include_hidden = bool(args.include_hidden)
        self.first_slide_as_title = not bool(args.no_first_slide_as_title)
        self.max_image_bytes = int(args.max_image_size)
        self._image_seq: dict[int, int] = {}

    # ---------- エントリ ----------

    def convert(self) -> Path:
        _validate_pptx(self.input_path)
        presentation = Presentation(str(self.input_path))
        self.images_dir.mkdir(parents=True, exist_ok=True)

        markdown_chunks: List[str] = []
        emitted_slide_no = 0
        for slide in presentation.slides:
            if self._is_hidden(slide) and not self.include_hidden:
                continue
            emitted_slide_no += 1
            chunk = self._convert_slide(slide, emitted_slide_no)
            if chunk:
                markdown_chunks.append(chunk)

        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        body = "\n\n".join(markdown_chunks)
        if not body.endswith("\n"):
            body += "\n"
        with open(self.output_path, "w", encoding="utf-8", newline="\n") as fh:
            fh.write(body)
        return self.output_path

    # ---------- スライド単位 ----------

    @staticmethod
    def _is_hidden(slide) -> bool:
        try:
            return slide.element.get("show") == "0"
        except Exception:
            return False

    def _convert_slide(self, slide, slide_no: int) -> str:
        title = self._extract_title(slide)
        if slide_no == 1 and self.first_slide_as_title:
            heading = f"# {title or 'タイトル'}"
        else:
            heading = f"## {title or f'スライド{slide_no}'}"

        shapes_meta: List[dict] = []
        connectors: List[dict] = []
        collected: List[dict] = []  # 各要素 {"kind": str, "shape_id": Optional[int], "text": str}

        for shape in slide.shapes:
            self._collect_shape(shape, slide_no, shapes_meta, connectors, collected)

        mermaid_md: Optional[str] = None
        used_ids: set = set()
        if not self.no_mermaid:
            extractor = FlowExtractor(shapes_meta, connectors)
            mermaid_md = extractor.build_mermaid()
            used_ids = extractor.used_node_ids

        body_blocks: List[str] = []
        for item in collected:
            if (
                item["kind"] == "text"
                and item["shape_id"] is not None
                and item["shape_id"] in used_ids
            ):
                continue
            if item["text"]:
                body_blocks.append(item["text"])

        if mermaid_md:
            body_blocks.append(f"```mermaid\n{mermaid_md}\n```")

        if self.include_notes:
            notes = self._extract_notes(slide)
            if notes:
                body_blocks.append(notes)

        parts = [heading] + [b for b in body_blocks if b]
        return "\n\n".join(parts)

    @staticmethod
    def _extract_title(slide) -> Optional[str]:
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                return _safe_text(title_shape.text_frame.text)
        except Exception:
            return None
        return None

    # ---------- シェイプ単位 ----------

    def _collect_shape(
        self,
        shape,
        slide_no: int,
        shapes_meta: list,
        connectors: list,
        collected: list,
    ) -> None:
        if self._is_title_shape(shape):
            return

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                self._collect_shape(sub, slide_no, shapes_meta, connectors, collected)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            block = self._handle_picture(shape, slide_no)
            if block:
                collected.append({"kind": "image", "shape_id": shape.shape_id, "text": block})
            return

        if getattr(shape, "has_table", False) and shape.has_table:
            block = self._convert_table(shape.table)
            if block:
                collected.append({"kind": "table", "shape_id": shape.shape_id, "text": block})
            return

        if getattr(shape, "has_chart", False) and shape.has_chart:
            block = self._summarize_chart(shape.chart)
            collected.append({"kind": "chart", "shape_id": shape.shape_id, "text": block})
            return

        if shape.shape_type == MSO_SHAPE_TYPE.LINE or self._is_connector(shape):
            self._collect_connector(shape, connectors)
            return

        if self._is_smartart(shape):
            mermaid = None
            if not self.no_mermaid:
                mermaid = self._smartart_to_mermaid(shape)
            if mermaid:
                collected.append(
                    {"kind": "mermaid", "shape_id": shape.shape_id, "text": f"```mermaid\n{mermaid}\n```"}
                )
            else:
                collected.append(
                    {"kind": "smartart", "shape_id": shape.shape_id, "text": self._smartart_fallback_text(shape)}
                )
            return

        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            shapes_meta.append(self._shape_meta(shape))
            text = self._convert_text_frame(shape.text_frame)
            if text:
                collected.append({"kind": "text", "shape_id": shape.shape_id, "text": text})
            return

    @staticmethod
    def _is_title_shape(shape) -> bool:
        try:
            placeholder = shape.placeholder_format
            if placeholder is None:
                return False
            if placeholder.idx == 0:
                return True
            placeholder_type = placeholder.type
            return placeholder_type is not None and "TITLE" in str(placeholder_type)
        except Exception:
            return False

    @staticmethod
    def _is_connector(shape) -> bool:
        try:
            return shape.element.tag.endswith("}cxnSp")
        except Exception:
            return False

    @staticmethod
    def _is_smartart(shape) -> bool:
        try:
            for child in shape.element.iter():
                if child.tag.endswith("}graphicData"):
                    uri = child.get("uri", "")
                    if "diagram" in uri:
                        return True
        except Exception:
            return False
        return False

    @staticmethod
    def _shape_meta(shape) -> dict:
        text = ""
        try:
            if shape.has_text_frame:
                text = _safe_text(shape.text_frame.text)
        except Exception:
            text = ""
        auto = ""
        try:
            auto = str(shape.auto_shape_type) if getattr(shape, "auto_shape_type", None) else ""
        except Exception:
            auto = ""
        return {
            "shape_id": shape.shape_id,
            "name": shape.name,
            "text": text,
            "auto_shape_type": auto,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        }

    def _collect_connector(self, shape, connectors: list) -> None:
        try:
            begin = None
            end = None
            for element in shape.element.iter():
                tag = element.tag
                if tag.endswith("}stCxn") and element.get("id"):
                    begin = int(element.get("id"))
                elif tag.endswith("}endCxn") and element.get("id"):
                    end = int(element.get("id"))
            if begin is not None and end is not None:
                connectors.append({"begin": begin, "end": end})
        except Exception:
            return

    # ---------- テキストフレーム ----------

    def _convert_text_frame(self, text_frame) -> str:
        rendered_paragraphs: List[str] = []
        code_buffer: List[str] = []

        def flush_code():
            if code_buffer:
                rendered_paragraphs.append("```\n" + "\n".join(code_buffer) + "\n```")
                code_buffer.clear()

        for paragraph in text_frame.paragraphs:
            text = self._render_paragraph_runs(paragraph)
            if not text.strip():
                flush_code()
                continue

            level = paragraph.level or 0
            if self._is_monospace_paragraph(paragraph):
                # 装飾を取り除いた素のテキストを残す
                raw = "".join((run.text or "") for run in paragraph.runs)
                code_buffer.append(raw)
                continue

            flush_code()

            if level > 0:
                indent = "  " * level
                rendered_paragraphs.append(f"{indent}- {text}")
            elif self._has_bullet(paragraph):
                rendered_paragraphs.append(f"- {text}")
            else:
                rendered_paragraphs.append(text)

        flush_code()
        return "\n".join(rendered_paragraphs)

    @staticmethod
    def _has_bullet(paragraph) -> bool:
        try:
            p_pr = paragraph._pPr
            if p_pr is None:
                return False
            for child in p_pr.iter():
                tag = child.tag
                if tag.endswith("}buChar") or tag.endswith("}buAutoNum"):
                    return True
                if tag.endswith("}buNone"):
                    return False
        except Exception:
            return False
        return False

    @staticmethod
    def _render_paragraph_runs(paragraph) -> str:
        rendered_parts: List[str] = []
        for run in paragraph.runs:
            text = run.text or ""
            if not text:
                continue
            bold = bool(run.font.bold)
            italic = bool(run.font.italic)
            strike = False
            try:
                r_pr = run._r.find(f"{{{NS_A}}}rPr")
                if r_pr is not None:
                    strike_attr = r_pr.get("strike")
                    if strike_attr and strike_attr != "noStrike":
                        strike = True
            except Exception:
                strike = False
            wrapped = text
            if strike:
                wrapped = f"~~{wrapped}~~"
            if italic:
                wrapped = f"*{wrapped}*"
            if bold:
                wrapped = f"**{wrapped}**"
            rendered_parts.append(wrapped)
        return "".join(rendered_parts)

    @staticmethod
    def _is_monospace_paragraph(paragraph) -> bool:
        runs = list(paragraph.runs)
        if not runs:
            return False
        for run in runs:
            try:
                name = run.font.name
            except Exception:
                name = None
            if not _is_monospace_font(name):
                return False
        return True

    # ---------- 表 ----------

    def _convert_table(self, table) -> str:
        rows = list(table.rows)
        if not rows:
            return ""
        header_cells = [self._cell_text(cell) for cell in rows[0].cells]
        col_count = len(header_cells)
        if col_count == 0:
            return ""
        lines = [
            "| " + " | ".join(_escape_md_pipe(c) for c in header_cells) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]
        for row in rows[1:]:
            cells = [self._cell_text(cell) for cell in row.cells]
            while len(cells) < col_count:
                cells.append("")
            lines.append(
                "| " + " | ".join(_escape_md_pipe(c) for c in cells[:col_count]) + " |"
            )
        return "\n".join(lines)

    @staticmethod
    def _cell_text(cell) -> str:
        if cell is None:
            return ""
        return _safe_text(cell.text or "")

    # ---------- 画像 ----------

    def _handle_picture(self, shape, slide_no: int) -> Optional[str]:
        try:
            blob = shape.image.blob
        except Exception as exc:
            print(
                f"Warning: failed to extract image on slide {slide_no}: {exc}",
                file=sys.stderr,
            )
            return None

        size = len(blob)
        if size > self.max_image_bytes:
            print(
                f"Warning: image too large ({size} bytes > {self.max_image_bytes}); "
                f"slide {slide_no} shape '{shape.name}'",
                file=sys.stderr,
            )
            return (
                f"> 画像（サイズ超過のためバイナリ非保存）: slide={slide_no}, "
                f"name={shape.name}, size={size}B"
            )

        raw_ext = (shape.image.ext or "png").lower()
        ext = raw_ext if raw_ext in ALLOWED_IMAGE_EXTS else "bin"
        if ext != raw_ext:
            print(
                f"Warning: image extension '{raw_ext}' is not in allowlist; "
                f"normalized to 'bin' on slide {slide_no}",
                file=sys.stderr,
            )
        self._image_seq.setdefault(slide_no, 0)
        self._image_seq[slide_no] += 1
        img_no = self._image_seq[slide_no]
        filename = f"slide{slide_no}_img{img_no}.{ext}"
        out_path = self.images_dir / filename
        try:
            out_path.write_bytes(blob)
        except OSError as exc:
            print(f"Warning: failed to write image: {exc}", file=sys.stderr)
            return None

        alt = _escape_md_link_text(self._image_alt(shape, img_no))
        try:
            rel = out_path.relative_to(self.output_path.parent)
        except ValueError:
            rel = out_path
        return f"![{alt}]({rel.as_posix()})"

    @staticmethod
    def _image_alt(shape, img_no: int) -> str:
        try:
            descr = shape._element.xpath("string(.//@descr)")
            if descr:
                return descr
        except Exception:
            pass
        return shape.name or f"image{img_no}"

    # ---------- チャート ----------

    @staticmethod
    def _summarize_chart(chart) -> str:
        try:
            type_name = str(chart.chart_type)
        except Exception:
            type_name = "unknown"
        series_names: List[str] = []
        try:
            for series in chart.series:
                try:
                    series_names.append(series.name)
                except Exception:
                    continue
        except Exception:
            series_names = []
        return f"> チャート: {type_name} 系列={series_names}"

    # ---------- SmartArt ----------

    def _smartart_to_mermaid(self, shape) -> Optional[str]:
        try:
            ns = {"dgm": NS_DGM, "r": NS_R, "a": NS_A}
            rel_ids = shape.element.findall(".//dgm:relIds", namespaces=ns)
            if not rel_ids:
                return None
            rid_data = rel_ids[0].get(f"{{{NS_R}}}dm")
            if not rid_data:
                return None
            slide_part = shape.part
            rel = slide_part.rels.get(rid_data)
            if rel is None:
                return None
            data_xml = etree.fromstring(rel.target_part.blob, parser=_hardened_xml_parser())
            points: dict[str, str] = {}
            for point in data_xml.findall(".//dgm:pt", namespaces=ns):
                pid = point.get("modelId")
                if not pid:
                    continue
                texts = point.findall(".//a:t", namespaces=ns)
                value = " ".join((t.text or "") for t in texts).strip()
                points[pid] = value
            edges: List[tuple] = []
            for connection in data_xml.findall(".//dgm:cxn", namespaces=ns):
                src = connection.get("srcId")
                dst = connection.get("destId")
                ctype = connection.get("type", "parOf")
                if ctype not in ("parOf", "presOf"):
                    continue
                if src in points and dst in points:
                    edges.append((src, dst))
            if not edges:
                return None

            def safe_id(pid: str) -> str:
                cleaned = re.sub(r"[^A-Za-z0-9]", "", pid)
                return f"S{cleaned[:20]}" if cleaned else f"S{abs(hash(pid)) % 100000}"

            used = {src for src, _ in edges} | {dst for _, dst in edges}
            lines = ["flowchart TD"]
            for nid in sorted(used):
                label = _mermaid_escape_label(points.get(nid, "") or nid)
                lines.append(f'    {safe_id(nid)}["{label}"]')
            for src, dst in edges:
                lines.append(f"    {safe_id(src)} --> {safe_id(dst)}")
            return "\n".join(lines)
        except Exception as exc:
            print(f"Warning: failed to parse SmartArt: {exc}", file=sys.stderr)
            return None

    @staticmethod
    def _smartart_fallback_text(shape) -> str:
        try:
            text_nodes = shape.element.findall(f".//{{{NS_A}}}t")
            lines: List[str] = []
            for node in text_nodes:
                if node.text and node.text.strip():
                    lines.append(f"- {node.text.strip()}")
            if lines:
                return "> SmartArt（テキスト抽出）:\n" + "\n".join(lines)
        except Exception:
            pass
        return "> SmartArt（解析失敗・テキスト抽出不可）"

    # ---------- スピーカーノート ----------

    @staticmethod
    def _extract_notes(slide) -> Optional[str]:
        try:
            if not slide.has_notes_slide:
                return None
            text = _safe_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            return None
        if not text:
            return None
        lines = ["> [!NOTE]"]
        for line in text.split("\n"):
            lines.append(f"> {line}")
        return "\n".join(lines)


# ----------------------------------------------------------------------------- #
# エントリポイント
# ----------------------------------------------------------------------------- #


def parse_args(argv: Iterable[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert PPTX to Markdown.")
    parser.add_argument("input", help="Input PPTX file path")
    parser.add_argument(
        "output",
        nargs="?",
        default=None,
        help="Output MD file path (default: <input>.md)",
    )
    parser.add_argument(
        "--images-dir",
        default=None,
        help="Image extraction directory (default: <basename>_images/ next to output MD)",
    )
    parser.add_argument(
        "--no-mermaid",
        action="store_true",
        help="Disable flowchart/SmartArt Mermaid conversion",
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes",
    )
    parser.add_argument(
        "--include-hidden",
        action="store_true",
        help="Include hidden slides",
    )
    parser.add_argument(
        "--no-first-slide-as-title",
        action="store_true",
        help="Treat first slide as H2 instead of H1",
    )
    parser.add_argument(
        "--max-image-size",
        type=int,
        default=DEFAULT_MAX_IMAGE_BYTES,
        help="Maximum bytes per image (default 5 MiB)",
    )
    return parser.parse_args(list(argv))


def main(argv: Optional[Iterable[str]] = None) -> int:
    args = parse_args(argv if argv is not None else sys.argv[1:])
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        return 1
    try:
        converter = PPTXMarkdownConverter(args)
        output_path = converter.convert()
        print(f"Wrote: {output_path}")
        print(f"Images dir: {converter.images_dir}")
        return 0
    except FileNotFoundError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1
    except ValueError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1
    except Exception as exc:  # noqa: BLE001
        print(f"Error: unexpected failure: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    sys.exit(main())
