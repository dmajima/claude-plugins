#!/usr/bin/env python3
"""convert_from_pptx.py - PPTX を Markdown に変換するスクリプト.

使い方:
    python convert_from_pptx.py <入力PPTX> [<出力MD>] [--images-dir DIR]
                                [--no-mermaid] [--include-notes]
                                [--include-hidden] [--no-first-slide-as-title]
                                [--max-image-size BYTES]
"""

from __future__ import annotations

import argparse
import re
import sys
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, List, Optional

try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

# XML 攻撃対策（XXE / Billion Laughs / DTD / external entity, CWE-611 / CWE-776）:
# - 本スクリプト独自の lxml 解析は `_hardened_xml_parser()` で
#   resolve_entities=False / no_network=True / load_dtd=False / huge_tree=False を適用.
# - python-pptx 内部の lxml は直接ハードニングしない代わりに、`_validate_pptx` の
#   ZIP bomb 検査（MAX_TOTAL_UNCOMPRESSED_BYTES / MAX_COMPRESSION_RATIO）と
#   上限定数（MAX_SLIDES / MAX_SHAPES_PER_SLIDE / MAX_GROUP_DEPTH /
#   MAX_TEXT_PER_SHAPE / MAX_TOTAL_IMAGE_BYTES / MAX_IMAGE_COUNT_PER_PPTX）で
#   DoS / 巨大エンティティ展開の影響範囲を限定する.
# - 旧コード `defusedxml.lxml.monkey_patch_lxml()` は defusedxml 0.7 で API が
#   削除されたため撤去した（呼び出すと AttributeError → fail-close で起動不能）.

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:
    print(f"Error: python-pptx is not installed: {exc}", file=sys.stderr, flush=True)
    sys.stderr.flush()
    sys.exit(2)

try:
    from lxml import etree
except ImportError as exc:
    print(f"Error: lxml is not installed: {exc}", file=sys.stderr, flush=True)
    sys.stderr.flush()
    sys.exit(2)


# ----------------------------------------------------------------------------- #
# 定数
# ----------------------------------------------------------------------------- #

MONOSPACE_FONTS = {
    "Consolas",
    "Courier",
    "Courier New",
    "Menlo",
    "Monaco",
    "Cascadia Code",
    "Cascadia Mono",
    "Fira Code",
    "Source Code Pro",
    "Lucida Console",
    "MS Gothic",
    "ＭＳ ゴシック",
    "MS Mincho",
    "ＭＳ 明朝",
}

PPTX_MAGIC = b"PK\x03\x04"
DEFAULT_MAX_IMAGE_BYTES = 5 * 1024 * 1024  # 5 MiB
MAX_TOTAL_UNCOMPRESSED_BYTES = 256 * 1024 * 1024  # 256 MiB (ZIP bomb 防御)
MAX_COMPRESSION_RATIO = 200  # compress_size に対する展開サイズ倍率上限

# DoS 防御: shape 数 / スライド数 / 再帰深度 / テキスト長 / 画像総量の上限
MAX_SLIDES = 1000
MAX_SHAPES_PER_SLIDE = 5000
MAX_GROUP_DEPTH = 20
MAX_TEXT_PER_SHAPE = 1_000_000  # 1 MB 相当
MAX_TOTAL_IMAGE_BYTES = 256 * 1024 * 1024  # 256 MiB (画像総量)
MAX_IMAGE_COUNT_PER_PPTX = 1000
MAX_ID_STRING_LEN = 32  # 数値属性文字列の長さ上限 (CWE-754, int_max_str_digits 対策)

# 双方向制御文字 (Bidi / homograph 偽装対策、CWE-1007)
_BIDI_CONTROL_CHARS = "‎‏‪‫‬‭‮⁦⁧⁨⁩"
_BIDI_CONTROL_RE = re.compile(f"[{_BIDI_CONTROL_CHARS}]")

ALLOWED_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp", "emf", "wmf"}

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _hardened_xml_parser() -> "etree.XMLParser":
    """XXE / DTD / 巨大ツリー攻撃を遮断した lxml パーサを返す."""
    return etree.XMLParser(
        resolve_entities=False,
        no_network=True,
        load_dtd=False,
        huge_tree=False,
    )


# ----------------------------------------------------------------------------- #
# ユーティリティ
# ----------------------------------------------------------------------------- #


def _safe_text(value: Optional[str]) -> str:
    """テキストを安全な形に正規化する.

    - `\\r` 除去 + 前後空白 strip
    - 双方向制御文字（Bidi override 等）の除去（homograph 偽装対策・CWE-1007）
    - サロゲートペア破損対策（UTF-8 round trip）
    - 長さ上限を超える場合は切り詰め（DoS 防御）
    """
    if not value:
        return ""
    text = value.replace("\r", "").strip()
    text = _BIDI_CONTROL_RE.sub("", text)
    # サロゲート文字の round trip 防御（JSON dump の UnicodeEncodeError 回避）
    text = text.encode("utf-8", errors="replace").decode("utf-8", errors="replace")
    if len(text) > MAX_TEXT_PER_SHAPE:
        text = text[:MAX_TEXT_PER_SHAPE]
    return text


def _escape_md_pipe(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", "<br>")


def _is_monospace_font(font_name: Optional[str]) -> bool:
    if not font_name:
        return False
    return font_name in MONOSPACE_FONTS


def _mermaid_escape_label(text: str) -> str:
    """Mermaid ラベルの HTML エスケープ（HR-G: XSS / SVG onclick / 構文破壊対策の強化版）.

    エスケープ対象: `&` / `<` / `>` / `"` / `'` / `` ` `` / `\\`.
    `&` を最優先で `&amp;` にエスケープし、後段の `&lt;` 等を二重デコードさせない.
    改行は Unicode 非文字 SENTINEL で退避し、HTML エスケープ後に `<br/>` へ復元する.
    Mermaid ラベルはノード描画時に SVG/HTML として展開されうるため、危険記号を
    包括的に HTML エンティティ化することで後段 convert-html での XSS を防ぐ.
    """
    text = text or ""
    # `&` を先にエスケープ（後段の `&lt;` 等を二重デコードさせないため）
    text = text.replace("&", "&amp;")
    # 改行は Unicode 非文字 SENTINEL で退避（入力に出現しないため衝突しない）
    SENTINEL = "￿BR￾"
    text = text.replace("\n", SENTINEL)
    # HTML / Mermaid メタ文字を包括的にエスケープ（XSS / SVG onclick / 構文破壊対策）
    text = text.replace("<", "&lt;").replace(">", "&gt;")
    text = text.replace('"', "&quot;").replace("'", "&#39;")
    text = text.replace("`", "&#96;").replace("\\", "&#92;")
    # SENTINEL を Mermaid の <br/> に復元
    text = text.replace(SENTINEL, "<br/>")
    if len(text) > 80:
        text = text[:77] + "..."
    return text or " "


def _escape_md_link_text(text: str) -> str:
    """Markdown のリンクテキスト・alt 部分のメタ文字をエスケープ."""
    return (
        (text or "")
        .replace("\\", "\\\\")
        .replace("[", "\\[")
        .replace("]", "\\]")
        .replace("(", "\\(")
        .replace(")", "\\)")
        .replace("\n", " ")
    )


def _validate_pptx(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(f"Input file not found: {path}")
    if path.suffix.lower() not in (".pptx", ".pptm"):
        raise ValueError(f"Input file is not PPTX/PPTM (extension): {path}")
    with open(path, "rb") as fh:
        magic = fh.read(4)
    if magic != PPTX_MAGIC:
        raise ValueError(f"Input file is not a valid PPTX (zip magic): {path}")
    try:
        with zipfile.ZipFile(path, "r") as zf:
            # ZIP bomb 検査: 総展開サイズ・圧縮率
            total_uncompressed = 0
            for info in zf.infolist():
                total_uncompressed += info.file_size
                if total_uncompressed > MAX_TOTAL_UNCOMPRESSED_BYTES:
                    raise ValueError(
                        f"Input PPTX exceeds uncompressed size limit "
                        f"({total_uncompressed} > {MAX_TOTAL_UNCOMPRESSED_BYTES}): {path}"
                    )
                if info.compress_size > 0:
                    ratio = info.file_size / info.compress_size
                    if ratio > MAX_COMPRESSION_RATIO:
                        raise ValueError(
                            f"Suspicious compression ratio ({ratio:.1f}x) "
                            f"for entry '{info.filename}' in {path}"
                        )
            content_types = zf.read("[Content_Types].xml").decode("utf-8", errors="ignore")
    except KeyError as exc:
        raise ValueError(f"Input file missing [Content_Types].xml: {path}") from exc
    if "presentationml" not in content_types.lower():
        raise ValueError(f"Input file is not a PresentationML document: {path}")


def _enforce_under(
    base: Path,
    candidate: Path,
    label: str,
    msg_suffix: str = "workspace root",
) -> Path:
    """candidate が base 配下に解決されることを検証し、絶対パスを返す.

    パストラバーサル攻撃の防止用（CWE-22 / CWE-73）。出力先パス引数すべてに
    共通適用される。candidate は絶対 / 相対のいずれでも受け付け、base からの
    相対解決後に base 配下に収まることを `relative_to` で確認する。

    シンボリックリンクは `resolve()` が追跡するため、リンク先が base 外なら拒否する
    （CWE-59 / CWE-367 部分対策。完全な TOCTOU 対策は書き込み時の is_symlink チェック
    を併用する）。
    """
    base_resolved = base.resolve(strict=False)
    if candidate.is_absolute():
        cand_resolved = candidate.resolve(strict=False)
    else:
        cand_resolved = (base / candidate).resolve(strict=False)
    try:
        cand_resolved.relative_to(base_resolved)
    except ValueError as exc:
        raise ValueError(
            f"{label} must be under {msg_suffix} (path traversal blocked): {cand_resolved}"
        ) from exc
    return cand_resolved


def _check_safe_output(path: Path, force: bool = False) -> None:
    """書込前の安全チェック（CWE-59 / CWE-377 対策）.

    - 既存ファイルがシンボリックリンクなら拒否（TOCTOU 対策）
    - 既存ファイル + `force=False` なら拒否（無確認上書き防止）
    """
    if path.exists():
        if path.is_symlink():
            raise ValueError(f"Refusing to write through symlink: {path}")
        if not force:
            raise ValueError(f"Output already exists (use --force to overwrite): {path}")


def _apply_safe_perm(path: Path) -> None:
    """書込後のファイル権限を 0o600 に絞る（POSIX のみ・CWE-732 対策）."""
    try:
        path.chmod(0o600)
    except (OSError, NotImplementedError):
        # Windows / 非 POSIX 環境では chmod が無効。続行する
        pass


def _safe_write_bytes(path: Path, data: bytes, force: bool = False) -> None:
    """ファイルにバイト書き込みする際の安全策（_check_safe_output + _apply_safe_perm を統合）."""
    _check_safe_output(path, force)
    path.write_bytes(data)
    _apply_safe_perm(path)


def _safe_int_from_str(value: Optional[str], max_len: int = MAX_ID_STRING_LEN) -> Optional[int]:
    """XML 由来の数値文字列を安全に int に変換する.

    長すぎる文字列（DoS 攻撃）は None を返す。CWE-754 / int_max_str_digits 対策。
    """
    if value is None:
        return None
    if len(value) > max_len:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _resolve_images_dir(output_md: Path, images_dir_opt: Optional[str]) -> Path:
    base = output_md.parent.resolve()
    if images_dir_opt:
        return _enforce_under(
            base, Path(images_dir_opt), "images-dir", "output MD directory"
        )
    return (base / f"{output_md.stem}_images").resolve()


# ----------------------------------------------------------------------------- #
# Mermaid 生成（コネクタ駆動）
# ----------------------------------------------------------------------------- #


class FlowExtractor:
    """同一スライド内の図形群とコネクタからフロー図 Mermaid を生成する."""

    def __init__(self, shapes_meta: List[dict], connectors: List[dict]):
        self.shapes_meta = shapes_meta
        self.connectors = connectors
        self.used_node_ids: set = set()

    def build_mermaid(self) -> Optional[str]:
        self.used_node_ids = set()
        if not self.connectors or len(self.shapes_meta) < 2:
            return None
        id_to_label = {
            meta["shape_id"]: (meta["text"] or meta["name"] or f"shape{meta['shape_id']}")
            for meta in self.shapes_meta
        }
        edges = []
        for connector in self.connectors:
            begin = connector.get("begin")
            end = connector.get("end")
            if begin is None or end is None:
                continue
            if begin not in id_to_label or end not in id_to_label:
                continue
            edges.append((begin, end))
            self.used_node_ids.add(begin)
            self.used_node_ids.add(end)
        if not edges:
            return None
        direction = self._infer_direction()
        lines = [f"flowchart {direction}"]
        # O(N) ルックアップを O(1) 化するため shape_id → meta の辞書を事前構築
        id_to_meta = {m["shape_id"]: m for m in self.shapes_meta}
        for nid in sorted(self.used_node_ids):
            meta = id_to_meta.get(nid)
            label = _mermaid_escape_label(id_to_label.get(nid, f"shape{nid}"))
            open_b, close_b = self._bracket_for_shape(meta)
            lines.append(f'    N{nid}{open_b}"{label}"{close_b}')
        for src, dst in edges:
            lines.append(f"    N{src} --> N{dst}")
        return "\n".join(lines)

    def _infer_direction(self) -> str:
        xs = [m["left"] for m in self.shapes_meta if m.get("left") is not None]
        ys = [m["top"] for m in self.shapes_meta if m.get("top") is not None]
        if not xs or not ys:
            return "TD"
        x_range = max(xs) - min(xs)
        y_range = max(ys) - min(ys)
        if x_range > y_range * 1.5:
            return "LR"
        return "TD"

    @staticmethod
    def _bracket_for_shape(meta: Optional[dict]) -> tuple:
        if not meta:
            return ("[", "]")
        auto = (meta.get("auto_shape_type") or "").upper()
        if "OVAL" in auto or "ELLIPSE" in auto or "FLOWCHART_TERMINATOR" in auto:
            return ("((", "))")
        if "DIAMOND" in auto or "FLOWCHART_DECISION" in auto:
            return ("{", "}")
        if "PARALLELOGRAM" in auto or "FLOWCHART_DATA" in auto:
            return ("[/", "/]")
        return ("[", "]")


# ----------------------------------------------------------------------------- #
# 変換本体
# ----------------------------------------------------------------------------- #


@dataclass
class SlideContext:
    """1 スライド処理中のローカル状態（Suggestion 12: 状態変数の集約）.

    PPTXMarkdownConverter のインスタンス変数 `_current_*` 群を集約することで、
    スライド処理の純関数化とテスタビリティ向上の準備とする.
    既存の `self._current_*` 参照は property 経由でこの dataclass のフィールドに
    委譲されるため、既存呼び出しは無傷で、新規実装は `self._slide_ctx` を直接
    参照できる. 将来的に `self._current_*` を撤去して `self._slide_ctx` に完全
    移行することを想定する.
    """
    title_shape_id: Optional[int] = None
    is_section_cover: bool = False
    repeated_texts: set = field(default_factory=set)
    max_font_pt: Optional[float] = None
    median_font_pt: Optional[float] = None
    decoration_shape_ids: set = field(default_factory=set)

    def reset(self) -> None:
        """スライド処理の境界でローカル状態を初期化する."""
        self.title_shape_id = None
        self.is_section_cover = False
        self.repeated_texts = set()
        self.max_font_pt = None
        self.median_font_pt = None
        self.decoration_shape_ids = set()


class PPTXMarkdownConverter:
    """python-pptx で読んだプレゼンを Markdown に転記する."""

    def __init__(self, args: argparse.Namespace) -> None:
        # HR-δ: 入力 PPTX のシンボリックリンクを拒否（CWE-59 / CWE-367 TOCTOU 対策）.
        # verify_md.py の _check_input_path と対称化.
        input_arg = Path(args.input)
        if input_arg.is_symlink():
            raise ValueError(
                f"Input PPTX is a symbolic link (refused for safety): {input_arg}"
            )
        self.input_path = input_arg.resolve()
        # workspace_root: 全副次出力のパストラバーサル検証 (CWE-22 / CWE-73) で共通利用するベース.
        # --workspace-root が明示された場合は output_path 本体もそこに強制（HR-B 修正）.
        # 未指定なら output_path.parent を workspace_root とする（後方互換）.
        workspace_arg = getattr(args, "workspace_root", None)
        default_output = (
            Path(args.output) if args.output else self.input_path.with_suffix(".md")
        )
        if workspace_arg:
            self._workspace_root: Path = Path(workspace_arg).resolve()
            self.output_path = _enforce_under(
                self._workspace_root,
                default_output,
                "--output (output MD path)",
                "workspace root (--workspace-root)",
            )
        else:
            self.output_path = default_output.resolve()
            self._workspace_root = self.output_path.parent.resolve()
        self.images_dir = _resolve_images_dir(self.output_path, args.images_dir)
        self.no_mermaid = bool(args.no_mermaid)
        self.include_notes = bool(args.include_notes)
        self.include_hidden = bool(args.include_hidden)
        self.first_slide_as_title = not bool(args.no_first_slide_as_title)
        self.max_image_bytes = int(args.max_image_size)
        self._image_seq: dict[int, int] = {}
        self._template_texts: set[str] = set()
        # SlideContext: 1 スライド処理中のローカル状態を集約.
        # 既存の self._current_* は property 経由でこの _slide_ctx を参照する（後方互換）.
        # property setter 経由で代入を受けるため、_slide_ctx を最初に初期化する必要がある.
        self._slide_ctx: SlideContext = SlideContext()
        self._current_title_shape_id = None
        self._current_is_section_cover = False
        self._current_repeated_texts = set()
        self._slide_width_emu: Optional[int] = None
        self._slide_height_emu: Optional[int] = None
        self._current_max_font_pt = None
        self._current_median_font_pt = None
        self._current_decoration_shape_ids = set()
        # Presentation インスタンスを 1 回だけ読み込んで複数 export メソッドで再利用する
        # ためのキャッシュ。_load_presentation() 経由でアクセスする
        self._presentation_cache: Optional["Presentation"] = None
        # マスタ/レイアウト由来のテンプレ装飾テキスト集合のキャッシュ（複数 export メソッド共有）
        self._template_texts_cache: Optional[set] = None
        # 既存ファイルの上書き許可フラグ（CWE-377 対策）
        self._force: bool = bool(getattr(args, "force", False))
        # 画像の総書込バイト数と枚数（MAX_TOTAL_IMAGE_BYTES / MAX_IMAGE_COUNT_PER_PPTX）
        self._image_total_bytes: int = 0
        self._image_total_count: int = 0
        self.structured_json_path = (
            _enforce_under(
                self._workspace_root,
                Path(args.structured_json),
                "--structured-json",
                "workspace root",
            )
            if getattr(args, "structured_json", None)
            else None
        )
        self.json_only = bool(getattr(args, "json_only", False))
        self.per_slide_json_dir = (
            _enforce_under(
                self._workspace_root,
                Path(args.per_slide_json),
                "--per-slide-json",
                "workspace root",
            )
            if getattr(args, "per_slide_json", None)
            else None
        )
        self.compact_view_dir = (
            _enforce_under(
                self._workspace_root,
                Path(args.compact_view),
                "--compact-view",
                "workspace root",
            )
            if getattr(args, "compact_view", None)
            else None
        )

    # ---------- SlideContext property（後方互換、Suggestion 12）----------
    # 既存の `self._current_*` 参照を SlideContext のフィールドに委譲する.
    # 既存呼び出しは無傷で、内部的にはスライド状態を 1 つの dataclass に集約.

    @property
    def _current_title_shape_id(self) -> Optional[int]:
        return self._slide_ctx.title_shape_id

    @_current_title_shape_id.setter
    def _current_title_shape_id(self, value: Optional[int]) -> None:
        self._slide_ctx.title_shape_id = value

    @property
    def _current_is_section_cover(self) -> bool:
        return self._slide_ctx.is_section_cover

    @_current_is_section_cover.setter
    def _current_is_section_cover(self, value: bool) -> None:
        self._slide_ctx.is_section_cover = bool(value)

    @property
    def _current_repeated_texts(self) -> set:
        return self._slide_ctx.repeated_texts

    @_current_repeated_texts.setter
    def _current_repeated_texts(self, value: set) -> None:
        self._slide_ctx.repeated_texts = value

    @property
    def _current_max_font_pt(self) -> Optional[float]:
        return self._slide_ctx.max_font_pt

    @_current_max_font_pt.setter
    def _current_max_font_pt(self, value: Optional[float]) -> None:
        self._slide_ctx.max_font_pt = value

    @property
    def _current_median_font_pt(self) -> Optional[float]:
        return self._slide_ctx.median_font_pt

    @_current_median_font_pt.setter
    def _current_median_font_pt(self, value: Optional[float]) -> None:
        self._slide_ctx.median_font_pt = value

    @property
    def _current_decoration_shape_ids(self) -> set:
        return self._slide_ctx.decoration_shape_ids

    @_current_decoration_shape_ids.setter
    def _current_decoration_shape_ids(self, value: set) -> None:
        self._slide_ctx.decoration_shape_ids = value

    # ---------- エントリ ----------

    def _load_presentation(self) -> "Presentation":
        """PPTX を 1 回だけ検証・読込し、キャッシュを返す.

        複数の export メソッド（convert / export_structured_json /
        export_per_slide_json / export_compact_view）が同一の Presentation
        インスタンスを再利用することで、PPTX を毎回開き直す重複初期化を解消する.

        注意: python-pptx 内部の lxml パーサは `_hardened_xml_parser` を経由しない.
        悪意 PPTX の Billion Laughs / 巨大 entity 展開リスクへの保護は、
        `_validate_pptx` の ZIP bomb 検査（MAX_TOTAL_UNCOMPRESSED_BYTES /
        MAX_COMPRESSION_RATIO）と上限定数（MAX_SLIDES / MAX_SHAPES_PER_SLIDE /
        MAX_GROUP_DEPTH / MAX_TEXT_PER_SHAPE / MAX_TOTAL_IMAGE_BYTES /
        MAX_IMAGE_COUNT_PER_PPTX）の併用によって与える.
        過去には `defusedxml.lxml.monkey_patch_lxml()` の起動時呼び出しを推奨していたが、
        defusedxml 0.7 で当該 API が削除されたため依存自体を廃止した.
        """
        if self._presentation_cache is None:
            _validate_pptx(self.input_path)
            self._presentation_cache = Presentation(str(self.input_path))
            try:
                self._slide_width_emu = self._presentation_cache.slide_width
                self._slide_height_emu = self._presentation_cache.slide_height
            except Exception:
                self._slide_width_emu = None
                self._slide_height_emu = None
            # スライド数の上限チェック（DoS 防御・CWE-400）
            try:
                slide_count = len(list(self._presentation_cache.slides))
                if slide_count > MAX_SLIDES:
                    raise ValueError(
                        f"slide count {slide_count} exceeds MAX_SLIDES ({MAX_SLIDES})"
                    )
            except ValueError:
                raise
            except Exception:
                pass
        return self._presentation_cache

    def _get_template_texts(self) -> set:
        """マスタ/レイアウト由来のテンプレ装飾テキスト集合を取得（キャッシュ付き）.

        複数の export メソッドが同じデータを使うため、初回呼び出しで構築して
        キャッシュする（MR-2: 重複走査の解消）。
        """
        if self._template_texts_cache is None:
            presentation = self._load_presentation()
            self._template_texts_cache = self._collect_template_texts(presentation)
        return self._template_texts_cache

    def convert(self) -> Path:
        presentation = self._load_presentation()
        self.images_dir.mkdir(parents=True, exist_ok=True)

        self._template_texts = self._get_template_texts()

        markdown_chunks: List[str] = []
        emitted_slide_no = 0
        for slide in presentation.slides:
            if self._is_hidden(slide) and not self.include_hidden:
                continue
            emitted_slide_no += 1
            chunk = self._convert_slide(slide, emitted_slide_no)
            if chunk:
                markdown_chunks.append(chunk)

        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        body = "\n\n".join(markdown_chunks)
        if not body.endswith("\n"):
            body += "\n"
        _check_safe_output(self.output_path, self._force)
        with open(self.output_path, "w", encoding="utf-8", newline="\n") as fh:
            fh.write(body)
        _apply_safe_perm(self.output_path)
        return self.output_path

    def export_structured_json(self) -> Path:
        """PPTX の全 shape を機械的に dump した構造化 JSON を出力する。

        装飾フィルタや並べ替えは適用せず、LLM が後段で文脈的に解釈できる
        生の情報を保持する。スキーマは references/json-schema.md を参照。
        """
        import json

        presentation = self._load_presentation()
        self.images_dir.mkdir(parents=True, exist_ok=True)

        # マスタ/レイアウト由来のテキストは "装飾候補" として参考情報のみ提供（除外はしない）
        template_texts = list(self._get_template_texts())

        slides_data = []
        emitted_slide_no = 0
        # 画像は MD 出力時と同様に抽出（pictures をスキャンして PNG を吐く）
        for slide in presentation.slides:
            if self._is_hidden(slide) and not self.include_hidden:
                continue
            emitted_slide_no += 1
            slides_data.append(self._slide_to_dict(slide, emitted_slide_no))

        document = {
            "metadata": {
                "input_path": str(self.input_path),
                "slide_count": len(slides_data),
                "slide_width_emu": self._slide_width_emu,
                "slide_height_emu": self._slide_height_emu,
                "images_dir": str(self.images_dir),
                "template_decoration_texts": template_texts,
                "schema_version": "1.0",
            },
            "slides": slides_data,
        }

        if self.structured_json_path is None:
            raise ValueError("structured_json_path is required for export_structured_json()")
        self.structured_json_path.parent.mkdir(parents=True, exist_ok=True)
        _check_safe_output(self.structured_json_path, self._force)
        with open(self.structured_json_path, "w", encoding="utf-8", newline="\n") as fh:
            json.dump(document, fh, ensure_ascii=False, indent=2)
        _apply_safe_perm(self.structured_json_path)
        return self.structured_json_path

    def export_per_slide_json(self) -> int:
        """大規模 PPTX 向け: スライド単位 JSON とメタデータ JSON を出力する。

        出力構成:
            <per-slide-dir>/metadata.json     -- 全体メタ（slide_count, 寸法, template_decoration_texts）
            <per-slide-dir>/slide-01.json     -- スライドごとの shape/connector データ
            <per-slide-dir>/slide-02.json
            ...

        Phase 2 で Claude がスライド単位に Read することで、コンテキスト ウィンドウへの
        負荷を分散できる（サブエージェント並列分担にも適する）。
        """
        import json

        presentation = self._load_presentation()
        self.images_dir.mkdir(parents=True, exist_ok=True)

        template_texts = list(self._get_template_texts())

        if self.per_slide_json_dir is None:
            raise ValueError("per_slide_json_dir is required for export_per_slide_json()")
        self.per_slide_json_dir.mkdir(parents=True, exist_ok=True)

        slide_summaries: List[dict] = []
        emitted = 0
        for slide in presentation.slides:
            if self._is_hidden(slide) and not self.include_hidden:
                continue
            emitted += 1
            slide_data = self._slide_to_dict(slide, emitted)
            slide_path = self.per_slide_json_dir / f"slide-{emitted:02d}.json"
            _check_safe_output(slide_path, self._force)
            with open(slide_path, "w", encoding="utf-8", newline="\n") as fh:
                json.dump(slide_data, fh, ensure_ascii=False, indent=2)
            _apply_safe_perm(slide_path)
            slide_summaries.append({
                "slide_no": emitted,
                "layout_name": slide_data["layout_name"],
                "is_section_cover_layout": slide_data["is_section_cover_layout"],
                "shape_count": len(slide_data["shapes"]),
                "connector_count": len(slide_data["connectors"]),
                "has_notes": bool(slide_data["notes"]),
                "file": slide_path.name,
            })

        metadata = {
            "input_path": str(self.input_path),
            "slide_count": emitted,
            "slide_width_emu": self._slide_width_emu,
            "slide_height_emu": self._slide_height_emu,
            "images_dir": str(self.images_dir),
            "template_decoration_texts": template_texts,
            "schema_version": "1.0",
            "slides_index": slide_summaries,
        }
        metadata_path = self.per_slide_json_dir / "metadata.json"
        _check_safe_output(metadata_path, self._force)
        with open(metadata_path, "w", encoding="utf-8", newline="\n") as fh:
            json.dump(metadata, fh, ensure_ascii=False, indent=2)
        _apply_safe_perm(metadata_path)
        return emitted

    def export_compact_view(self) -> int:
        """大規模 PPTX 向け: 1 スライド 1 ファイルの人間/LLM 可読な簡潔ビューを出力する。

        各 shape は 1 行で「pos (top,left) / size (h,w) / フォント / プレースホルダ / フラグ / テキスト」を表示する。
        Claude が Read で読み込み Phase 2 解釈する際の標準フォーマット。
        """
        presentation = self._load_presentation()
        self.images_dir.mkdir(parents=True, exist_ok=True)

        template_texts = self._get_template_texts()

        if self.compact_view_dir is None:
            raise ValueError("compact_view_dir is required for export_compact_view()")
        self.compact_view_dir.mkdir(parents=True, exist_ok=True)

        emitted = 0
        for slide in presentation.slides:
            if self._is_hidden(slide) and not self.include_hidden:
                continue
            emitted += 1
            view_text = self._render_compact_slide_view(slide, emitted, template_texts)
            view_path = self.compact_view_dir / f"slide-{emitted:02d}.txt"
            _check_safe_output(view_path, self._force)
            with open(view_path, "w", encoding="utf-8", newline="\n") as fh:
                fh.write(view_text)
            _apply_safe_perm(view_path)
        return emitted

    def _render_compact_slide_view(self, slide, slide_no: int, template_texts: set) -> str:
        """1 スライドを人間/LLM 可読な簡潔ビュー文字列にする。"""
        try:
            layout_name = slide.slide_layout.name or ""
        except Exception:
            layout_name = ""
        is_cover = self._is_section_cover_layout(slide)

        # shape を再帰展開して位置情報付きで集める
        shapes_data: List[dict] = []
        connectors_data: List[dict] = []
        self._walk_shape_to_dict(slide.shapes, slide_no, shapes_data, connectors_data, parent_path=[])

        lines: List[str] = []
        lines.append(f"=== Slide {slide_no} ===")
        lines.append(f"layout: {layout_name!r}")
        lines.append(f"is_section_cover: {is_cover}")
        lines.append(f"shape_count: {len(shapes_data)}, connector_count: {len(connectors_data)}")
        lines.append("")

        # 視覚順 (top, left)
        sorted_shapes = sorted(
            shapes_data,
            key=lambda s: (
                (s["geometry"].get("top_ratio") or 0),
                (s["geometry"].get("left_ratio") or 0),
            ),
        )

        for idx, shape in enumerate(sorted_shapes):
            g = shape["geometry"]
            top_r = g.get("top_ratio") or 0
            left_r = g.get("left_ratio") or 0
            h_r = g.get("height_ratio") or 0
            w_r = g.get("width_ratio") or 0
            ph = shape.get("placeholder")
            ph_str = f"ph={ph['type']}" if ph else "ph=-"
            auto = shape.get("auto_shape_type") or ""
            auto_str = f"auto={auto}" if auto else ""
            text = (shape.get("text") or "").replace("\n", "\\n")
            if len(text) > 200:
                text = text[:197] + "..."
            font_sz = shape.get("font_size_max_pt")
            color = shape.get("font_color")
            flags = []
            if shape.get("is_grayish_color"):
                flags.append("gray")
            if shape.get("text") and shape["text"].strip() in template_texts:
                flags.append("TEMPLATE")
            if shape.get("kind") and shape["kind"] != "TEXT_FRAME":
                flags.append(shape["kind"])
            if shape.get("group_path"):
                flags.append(f"grp={len(shape['group_path'])}")
            flag_str = " ".join(flags)
            table_info = ""
            if shape.get("table"):
                t = shape["table"]
                table_info = f" TABLE({t['row_count']}x{t['col_count']})"
            image_info = " IMAGE" if shape.get("image") else ""
            lines.append(
                f"  [{idx:02d}] id={shape['shape_id']:>4} "
                f"pos=({top_r:.2f},{left_r:.2f}) "
                f"size=({h_r:.2f},{w_r:.2f}) "
                f"font={font_sz} color={color} "
                f"{ph_str} {auto_str} {flag_str}{table_info}{image_info}"
            )
            lines.append(f"        text={text!r}")
            lines.append("")

        if connectors_data:
            lines.append("--- connectors ---")
            for c in connectors_data:
                lines.append(f"  {c.get('begin_shape_id')} -> {c.get('end_shape_id')}")
            lines.append("")

        try:
            if slide.has_notes_slide and slide.notes_slide is not None:
                notes = _safe_text(slide.notes_slide.notes_text_frame.text)
                if notes:
                    lines.append("--- notes ---")
                    lines.append(notes)
        except Exception:
            pass

        return "\n".join(lines) + "\n"

    def _slide_to_dict(self, slide, slide_no: int) -> dict:
        """1 スライドを JSON 形式の dict にする。"""
        try:
            layout_name = slide.slide_layout.name or ""
        except Exception:
            layout_name = ""

        shapes_list: List[dict] = []
        connectors_list: List[dict] = []
        self._walk_shape_to_dict(slide.shapes, slide_no, shapes_list, connectors_list, parent_path=[])

        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide is not None:
                notes_text = _safe_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            notes_text = ""

        return {
            "slide_no": slide_no,
            "layout_name": layout_name,
            "is_section_cover_layout": self._is_section_cover_layout(slide),
            "shapes": shapes_list,
            "connectors": connectors_list,
            "notes": notes_text,
        }

    def _walk_shape_to_dict(self, shape_iter, slide_no: int, shapes_out: list, connectors_out: list, parent_path: list, depth: int = 0) -> None:
        """shape ツリーを再帰的に走査し、JSON 用辞書を構築する.

        DoS 防御:
        - グループのネスト深度を MAX_GROUP_DEPTH で制限（CWE-674）
        - 1 スライドあたりの shape 数を MAX_SHAPES_PER_SLIDE で制限（CWE-400）
        - グループ shape 自身が装飾扱いの場合は子の再帰走査もスキップ（誤フィルタ防止）
        """
        if depth > MAX_GROUP_DEPTH:
            raise ValueError(
                f"Group nesting exceeds MAX_GROUP_DEPTH ({MAX_GROUP_DEPTH}) on slide {slide_no}"
            )
        for shape in shape_iter:
            if len(shapes_out) >= MAX_SHAPES_PER_SLIDE:
                raise ValueError(
                    f"shape count exceeds MAX_SHAPES_PER_SLIDE ({MAX_SHAPES_PER_SLIDE}) on slide {slide_no}"
                )
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_id = shape.shape_id
                    # グループ自身が装飾扱いなら子もスキップ（HR-2: 誤フィルタ防止）
                    if group_id in self._current_decoration_shape_ids:
                        continue
                    group_name = shape.name
                    new_path = parent_path + [{"shape_id": group_id, "name": group_name}]
                    self._walk_shape_to_dict(
                        shape.shapes, slide_no, shapes_out, connectors_out, new_path, depth=depth + 1
                    )
                    continue
            except ValueError:
                raise
            except Exception:
                pass

            # コネクタは別配列にも格納する
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.LINE or self._is_connector(shape):
                    info = self._extract_connector_info(shape)
                    if info:
                        connectors_out.append(info)
            except Exception:
                pass

            entry = self._shape_to_dict(shape, slide_no, parent_path)
            if entry is not None:
                shapes_out.append(entry)

    def _shape_to_dict(self, shape, slide_no: int, parent_path: list) -> Optional[dict]:
        """1 shape を JSON 用辞書化."""
        try:
            shape_id = shape.shape_id
            name = shape.name
        except Exception:
            return None

        try:
            top = int(shape.top) if shape.top is not None else None
        except Exception:
            top = None
        try:
            left = int(shape.left) if shape.left is not None else None
        except Exception:
            left = None
        try:
            width = int(shape.width) if shape.width is not None else None
        except Exception:
            width = None
        try:
            height = int(shape.height) if shape.height is not None else None
        except Exception:
            height = None

        slide_w = self._slide_width_emu
        slide_h = self._slide_height_emu
        top_ratio = (top / slide_h) if (top is not None and slide_h) else None
        left_ratio = (left / slide_w) if (left is not None and slide_w) else None
        width_ratio = (width / slide_w) if (width is not None and slide_w) else None
        height_ratio = (height / slide_h) if (height is not None and slide_h) else None

        # 種別判定
        kind = "OTHER"
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            kind = "PICTURE"
        elif getattr(shape, "has_table", False):
            try:
                if shape.has_table:
                    kind = "TABLE"
            except Exception:
                pass
        elif shape_type == MSO_SHAPE_TYPE.LINE or self._is_connector(shape):
            kind = "CONNECTOR"
        elif self._is_smartart(shape):
            kind = "SMARTART"
        elif getattr(shape, "has_chart", False):
            try:
                if shape.has_chart:
                    kind = "CHART"
            except Exception:
                pass
        elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            kind = "TEXT_FRAME"

        # placeholder 情報
        ph_info: Optional[dict] = None
        try:
            ph = shape.placeholder_format
            if ph is not None:
                ph_info = {
                    "idx": ph.idx,
                    "type": str(ph.type) if ph.type is not None else None,
                }
        except Exception:
            ph_info = None

        # auto_shape_type
        auto_type = ""
        try:
            auto_type = str(shape.auto_shape_type) if getattr(shape, "auto_shape_type", None) else ""
        except Exception:
            auto_type = ""

        # テキスト
        text = ""
        paragraphs_data: List[dict] = []
        try:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = _safe_text(shape.text_frame.text)
                for para in shape.text_frame.paragraphs:
                    runs_data = []
                    for run in para.runs:
                        run_text = run.text or ""
                        if not run_text:
                            continue
                        rf = {"text": run_text}
                        try:
                            sz = run.font.size
                            if sz is not None:
                                rf["font_size_pt"] = float(sz.pt)
                        except Exception:
                            pass
                        try:
                            rf["bold"] = bool(run.font.bold)
                        except Exception:
                            pass
                        try:
                            color = run.font.color
                            if color is not None:
                                rgb = color.rgb
                                if rgb is not None:
                                    rf["color"] = str(rgb).upper()
                        except Exception:
                            pass
                        runs_data.append(rf)
                    paragraphs_data.append({
                        "level": int(para.level or 0),
                        "runs": runs_data,
                    })
        except Exception:
            pass

        # テーブル
        table_data: Optional[dict] = None
        try:
            if getattr(shape, "has_table", False) and shape.has_table:
                rows_out = []
                for row in shape.table.rows:
                    cells_out = []
                    for cell in row.cells:
                        cells_out.append(_safe_text(cell.text or ""))
                    rows_out.append(cells_out)
                table_data = {
                    "rows": rows_out,
                    "row_count": len(rows_out),
                    "col_count": len(rows_out[0]) if rows_out else 0,
                }
        except Exception:
            table_data = None

        # 画像
        image_data: Optional[dict] = None
        if kind == "PICTURE":
            block = self._handle_picture(shape, slide_no)
            if block:
                # block は "![alt](images_dir/file.png)" 形式の Markdown 文字列
                image_data = {"markdown_link": block}

        font_size_max = self._extract_max_font_size_pt(shape)
        font_color = self._extract_dominant_font_color(shape)

        return {
            "shape_id": shape_id,
            "name": name,
            "kind": kind,
            "auto_shape_type": auto_type,
            "placeholder": ph_info,
            "group_path": parent_path,
            "geometry": {
                "top_emu": top,
                "left_emu": left,
                "width_emu": width,
                "height_emu": height,
                "top_ratio": top_ratio,
                "left_ratio": left_ratio,
                "width_ratio": width_ratio,
                "height_ratio": height_ratio,
            },
            "text": text,
            "paragraphs": paragraphs_data,
            "font_size_max_pt": font_size_max,
            "font_color": font_color,
            "is_grayish_color": self._is_grayish_color(font_color),
            "table": table_data,
            "image": image_data,
        }

    def _extract_connector_info(self, shape) -> Optional[dict]:
        """コネクタの接続元/接続先 shape_id を返す."""
        try:
            begin = None
            end = None
            for element in shape.element.iter():
                tag = element.tag
                if tag.endswith("}stCxn") and element.get("id"):
                    begin = _safe_int_from_str(element.get("id"))
                elif tag.endswith("}endCxn") and element.get("id"):
                    end = _safe_int_from_str(element.get("id"))
            return {"begin_shape_id": begin, "end_shape_id": end, "connector_shape_id": shape.shape_id}
        except Exception:
            return None

    # ---------- スライド単位 ----------

    @staticmethod
    def _is_hidden(slide) -> bool:
        try:
            return slide.element.get("show") == "0"
        except Exception:
            return False

    def _convert_slide(self, slide, slide_no: int) -> str:
        # スライド単位のフォント・装飾統計を事前計算（視覚的役割推定の基礎データ）
        self._current_is_section_cover = self._is_section_cover_layout(slide)
        self._compute_slide_visual_stats(slide)

        title, title_shape_id = self._extract_title(slide)
        self._current_title_shape_id = title_shape_id
        # スライド内で同一テキストが複数回出現する shape を装飾とみなす集合
        self._current_repeated_texts = self._collect_repeated_slide_texts(slide)
        if slide_no == 1 and self.first_slide_as_title:
            heading = f"# {title or 'タイトル'}"
        else:
            heading = f"## {title or f'スライド{slide_no}'}"

        shapes_meta: List[dict] = []
        connectors: List[dict] = []
        collected: List[dict] = []  # 各要素 {"kind": str, "shape_id": Optional[int], "text": str}

        for shape in slide.shapes:
            self._collect_shape(shape, slide_no, shapes_meta, connectors, collected)

        mermaid_md: Optional[str] = None
        used_ids: set = set()
        if not self.no_mermaid:
            extractor = FlowExtractor(shapes_meta, connectors)
            mermaid_md = extractor.build_mermaid()
            used_ids = extractor.used_node_ids

        # 視覚順 (top → left) でコンテンツブロックを並べ替える。
        # PPTX 内部の Z 順序ではなく、読者が画面で目で追う順序に近づける。
        collected.sort(key=lambda it: (it.get("top") or 0, it.get("left") or 0))

        # Mermaid 化されたノードのテキストは本文側で除外（行マージ前に行うこと）
        title_text_norm = (title or "").strip()
        filtered: List[dict] = []
        for item in collected:
            if (
                item["kind"] == "text"
                and item.get("shape_id") is not None
                and item["shape_id"] in used_ids
            ):
                continue
            if title_text_norm and item["kind"] == "text":
                stripped = (item.get("text") or "").lstrip("- ").strip()
                if stripped == title_text_norm:
                    continue
            filtered.append(item)

        # 同一行 (top の差が小さい) の連続する text 要素は左→右で結合する。
        # PowerPoint で「番号 + 章名」のように水平に並ぶ shape を 1 行として扱う。
        filtered = self._merge_horizontal_text_rows(filtered)

        body_blocks: List[str] = []
        for item in filtered:
            text_block = item.get("text") or ""
            if not text_block:
                continue
            body_blocks.append(text_block)

        if mermaid_md:
            body_blocks.append(f"```mermaid\n{mermaid_md}\n```")

        if self.include_notes:
            notes = self._extract_notes(slide)
            if notes:
                body_blocks.append(notes)

        # タイトルも本文も空ならスライド見出しごと出力を抑制（テンプレ装飾のみで本体が無いスライド）
        if not title and not any(b.strip() for b in body_blocks):
            return ""

        parts = [heading] + [b for b in body_blocks if b]
        return "\n\n".join(parts)

    def _extract_title(self, slide) -> tuple:
        """タイトル文字列とその shape_id を返す。

        優先順位:
        1. slide.shapes.title （CENTER_TITLE / TITLE placeholder）
        2. スライド最上部の短文 placeholder / textbox（テキストボックスで描画されたタイトル対応）
        """
        try:
            title_shape = slide.shapes.title
            if (
                title_shape is not None
                and getattr(title_shape, "has_text_frame", False)
                and title_shape.has_text_frame
                and title_shape.text_frame.text.strip()
            ):
                return _safe_text(title_shape.text_frame.text), title_shape.shape_id
        except Exception:
            pass

        candidate = self._guess_title_shape(slide)
        if candidate is not None:
            return _safe_text(candidate.text_frame.text), candidate.shape_id

        return None, None

    def _guess_title_shape(self, slide):
        """slide.shapes.title 不在時の代替: 視覚的に最上部にあるタイトルらしい shape を推定.

        条件 (いずれかを満たす shape を候補化):
            A. スライド上端 (top <= 1.6cm) かつ薄い帯 (height <= 1.4cm) かつ短文 (<= 80字)
            B. スライド中央上部 (top <= 50%) かつ最大フォントサイズに近い (>=80%)
               フォント・短文（章扉スライド：中央に大きな章タイトル）
        どちらも、テンプレ装飾 / 装飾 shape は除外する。
        """
        TOP_THRESHOLD_EMU = 1_500_000
        HEIGHT_THRESHOLD_EMU = 1_400_000
        MAX_TITLE_LEN = 80

        slide_h = self._slide_height_emu or 0
        max_font = self._current_max_font_pt
        candidates = []  # (priority, top, left, shape)

        for shape in slide.shapes:
            if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                continue
            if self._is_template_placeholder(shape):
                continue
            try:
                shape_id = shape.shape_id
            except Exception:
                shape_id = None
            if shape_id is not None and shape_id in self._current_decoration_shape_ids:
                continue
            try:
                top = shape.top or 0
                left = shape.left or 0
                height = shape.height or 0
            except Exception:
                continue
            text = _safe_text(shape.text_frame.text)
            if not text or len(text) > MAX_TITLE_LEN:
                continue
            if text in self._template_texts:
                continue

            font_pt = self._extract_max_font_size_pt(shape)

            # A: 最上部・薄い帯・短文
            if top <= TOP_THRESHOLD_EMU and height <= HEIGHT_THRESHOLD_EMU:
                candidates.append((1, top, left, shape))
                continue

            # B: スライド上半分かつ最大フォントに近い大型テキスト（章扉中央タイトル等）
            if max_font and font_pt and font_pt >= max_font * 0.8 and slide_h and top <= slide_h * 0.5:
                candidates.append((2, top, left, shape))
                continue

        if not candidates:
            return None
        candidates.sort(key=lambda x: (x[0], x[1], x[2]))
        return candidates[0][3]

    @staticmethod
    def _collect_template_texts(presentation) -> set:
        """テンプレ装飾として除外するテキスト集合を構築する。

        収集対象:
            1. スライドマスタ / レイアウトに登場するテキスト
            2. 全スライドの 3 スライド以上に同テキストで出現する ASCII 短文
               (例: "Core Colors" / "Sub Colors" 等の凡例ラベル) のみ。
               日本語本文ワードや混合テキストは誤除外回避のため対象外。
        """
        texts: set = set()

        def _gather(container):
            try:
                shapes = container.shapes
            except Exception:
                return
            for shape in shapes:
                try:
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        t = _safe_text(shape.text_frame.text)
                        if t:
                            texts.add(t)
                except Exception:
                    continue

        try:
            for master in presentation.slide_masters:
                _gather(master)
                try:
                    for layout in master.slide_layouts:
                        _gather(layout)
                except Exception:
                    continue
        except Exception:
            return texts

        # ASCII 短文 (英数字 + 空白) の全スライド頻出パターンを凡例ラベルとして検出
        try:
            slides = list(presentation.slides)
        except Exception:
            return texts
        if len(slides) < 3:
            return texts

        def _iter_text_shapes(shape_iter, depth: int = 0):
            # HR-η: グループの再帰深度を MAX_GROUP_DEPTH で打切（CWE-674 対策）
            if depth > MAX_GROUP_DEPTH:
                return
            for shape in shape_iter:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        yield from _iter_text_shapes(shape.shapes, depth + 1)
                        continue
                except Exception:
                    pass
                try:
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        yield shape
                except Exception:
                    continue

        appearance: dict = {}
        for slide in slides:
            seen: set = set()
            try:
                for shape in _iter_text_shapes(slide.shapes):
                    t = _safe_text(shape.text_frame.text)
                    if not t or len(t) > 25:
                        continue
                    # ASCII 範囲のみ（半角英数字 + 記号 + 空白）に限定
                    if not all((ord(c) < 128) for c in t):
                        continue
                    # 数字のみは章番号と紛らわしいので対象外
                    if t.isdigit():
                        continue
                    seen.add(t)
            except Exception:
                continue
            for t in seen:
                appearance[t] = appearance.get(t, 0) + 1

        # 3 スライド以上で出現する ASCII 短文を装飾ラベルとみなす
        for t, cnt in appearance.items():
            if cnt >= 3:
                texts.add(t)
        return texts

    @staticmethod
    def _is_template_placeholder(shape) -> bool:
        """FOOTER / SLIDE_NUMBER / DATE プレースホルダ判定。"""
        try:
            ph = shape.placeholder_format
            if ph is None:
                return False
            ph_type = ph.type
            if ph_type is None:
                return False
            type_str = str(ph_type)
            return any(kw in type_str for kw in ("FOOTER", "SLIDE_NUMBER", "DATE"))
        except Exception:
            return False

    def _matches_template_text(self, shape) -> bool:
        """マスタ/レイアウトに同一テキストで存在する shape を装飾扱いにする。"""
        if not self._template_texts:
            return False
        try:
            if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                return False
            text = _safe_text(shape.text_frame.text)
            if not text:
                return False
            return text in self._template_texts
        except Exception:
            return False

    def _compute_slide_visual_stats(self, slide) -> None:
        """スライド内全 shape の最大/中央値フォントサイズと装飾候補 shape を識別する。

        - 装飾候補: 薄いグレー色 / 極小フォント (中央値の 70% 未満) /
          スライド端 (右下) の極小 shape / 全 shape が小さい場合は除外。
        - decoration_shape_ids を集合に格納し、後段の _collect_shape で参照する。
        """
        self._current_max_font_pt = None
        self._current_median_font_pt = None
        self._current_decoration_shape_ids = set()

        # スライド寸法（基準）
        slide_w = self._slide_width_emu or 0
        slide_h = self._slide_height_emu or 0

        shapes_metrics: List[dict] = []
        try:
            self._gather_metrics_recursive(slide.shapes, shapes_metrics)
        except Exception:
            return

        font_sizes = [m["font_size_pt"] for m in shapes_metrics if m.get("font_size_pt")]
        if font_sizes:
            self._current_max_font_pt = max(font_sizes)
            sorted_sizes = sorted(font_sizes)
            mid = len(sorted_sizes) // 2
            if len(sorted_sizes) % 2 == 0 and len(sorted_sizes) >= 2:
                self._current_median_font_pt = (sorted_sizes[mid - 1] + sorted_sizes[mid]) / 2
            else:
                self._current_median_font_pt = sorted_sizes[mid]

        # 装飾判定
        median = self._current_median_font_pt
        for meta in shapes_metrics:
            shape_id = meta["shape_id"]
            text = (meta.get("text") or "").strip()
            if not text:
                continue
            font_pt = meta.get("font_size_pt")
            is_gray = meta.get("is_grayish", False)
            char_count = meta.get("char_count", 0)
            top = meta.get("top") or 0
            left = meta.get("left") or 0
            width = meta.get("width") or 0
            height = meta.get("height") or 0

            decoration_reasons = 0

            # 1) 極小フォント（スライド中央値の 70% 未満）かつ短文
            if font_pt is not None and median is not None and font_pt < median * 0.7 and char_count <= 30:
                decoration_reasons += 1

            # 2) フォント色が薄いグレー系
            if is_gray:
                decoration_reasons += 1

            # 3) shape が極端に小さい (面積がスライドの 0.3% 未満) かつ短文
            if slide_w and slide_h:
                area = (width or 0) * (height or 0)
                slide_area = slide_w * slide_h
                if slide_area > 0 and area > 0 and (area / slide_area) < 0.003 and char_count <= 20:
                    decoration_reasons += 1

            # 4) 右下の端 (top > 75%, left > 75%) に位置する短文
            if slide_w and slide_h and (top / slide_h) > 0.75 and (left / slide_w) > 0.55 and char_count <= 25:
                decoration_reasons += 1

            if decoration_reasons >= 2:
                self._current_decoration_shape_ids.add(shape_id)

    def _gather_metrics_recursive(self, shape_iter, out: list, depth: int = 0) -> None:
        """グループを再帰展開しながら全 shape のメタデータを集める.

        HR-η: 再帰深度を MAX_GROUP_DEPTH で打切（CWE-674 対策、_walk_shape_to_dict と対称化）.
        """
        if depth > MAX_GROUP_DEPTH:
            return
        for shape in shape_iter:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._gather_metrics_recursive(shape.shapes, out, depth + 1)
                    continue
            except Exception:
                pass
            try:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    out.append(self._shape_meta(shape))
            except Exception:
                continue

    @staticmethod
    def _is_section_cover_layout(slide) -> bool:
        """章扉スライド (innerCoverA / cover 系レイアウト) を判定."""
        try:
            name = (slide.slide_layout.name or "").lower()
        except Exception:
            return False
        if not name:
            return False
        return any(kw in name for kw in ("innercover", "cover", "section", "chapter", "扉"))

    @staticmethod
    def _is_decoration_number(shape) -> bool:
        """純数字 (1〜3 桁) のテキスト shape を装飾的章番号と判定."""
        try:
            if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                return False
            text = _safe_text(shape.text_frame.text)
        except Exception:
            return False
        if not text:
            return False
        return text.isdigit() and 1 <= len(text) <= 3

    @staticmethod
    def _collect_repeated_slide_texts(slide) -> set:
        """同一スライド内で 2 回以上出現する短いテキスト (<= 30 文字) を集める."""
        counts: dict = {}
        try:
            shapes = list(slide.shapes)
        except Exception:
            return set()

        def _walk(shape_iter):
            for shape in shape_iter:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _walk(shape.shapes)
                        continue
                except Exception:
                    pass
                try:
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        text = _safe_text(shape.text_frame.text)
                        if text and len(text) <= 30:
                            counts[text] = counts.get(text, 0) + 1
                except Exception:
                    continue

        _walk(shapes)
        # 凡例ラベル等の装飾は同一スライド内で 3 回以上繰り返される傾向にある。
        # 2 回出現の場合は組織図のように内容として意図的な重複の可能性が高いため対象外。
        return {t for t, c in counts.items() if c >= 3}

    def _matches_repeated_slide_text(self, shape) -> bool:
        """同一スライド内に重複出現するテキスト shape を装飾扱いにする."""
        repeated = getattr(self, "_current_repeated_texts", None)
        if not repeated:
            return False
        try:
            if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                return False
            text = _safe_text(shape.text_frame.text)
            if not text:
                return False
            return text in repeated
        except Exception:
            return False

    # ---------- シェイプ単位 ----------

    def _collect_shape(
        self,
        shape,
        slide_no: int,
        shapes_meta: list,
        connectors: list,
        collected: list,
    ) -> None:
        if self._is_title_shape(shape):
            return

        # 当該スライドのタイトルとして使用済みの shape は重複出力を避ける
        try:
            if self._current_title_shape_id is not None and shape.shape_id == self._current_title_shape_id:
                return
        except Exception:
            pass

        # フッタ / スライド番号 / 日付プレースホルダはテンプレ装飾として除外
        if self._is_template_placeholder(shape):
            return

        # マスタ / レイアウトに同じテキストで定義された装飾要素を除外
        if self._matches_template_text(shape):
            return

        # 章扉スライドの装飾的章番号 (純数字 1〜3 桁) を除外
        if self._current_is_section_cover and self._is_decoration_number(shape):
            return

        # 視覚統計ベースの装飾判定: 極小フォント・薄いグレー・極小サイズ・隅の短文を装飾扱い
        try:
            if shape.shape_id in self._current_decoration_shape_ids:
                return
        except Exception:
            pass

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                self._collect_shape(sub, slide_no, shapes_meta, connectors, collected)
            return

        pos_top, pos_left = self._safe_position(shape)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            block = self._handle_picture(shape, slide_no)
            if block:
                collected.append({"kind": "image", "shape_id": shape.shape_id, "text": block, "top": pos_top, "left": pos_left})
            return

        if getattr(shape, "has_table", False) and shape.has_table:
            block = self._convert_table(shape.table)
            if block:
                collected.append({"kind": "table", "shape_id": shape.shape_id, "text": block, "top": pos_top, "left": pos_left})
            return

        if getattr(shape, "has_chart", False) and shape.has_chart:
            block = self._summarize_chart(shape.chart)
            collected.append({"kind": "chart", "shape_id": shape.shape_id, "text": block, "top": pos_top, "left": pos_left})
            return

        if shape.shape_type == MSO_SHAPE_TYPE.LINE or self._is_connector(shape):
            self._collect_connector(shape, connectors)
            return

        if self._is_smartart(shape):
            mermaid = None
            if not self.no_mermaid:
                mermaid = self._smartart_to_mermaid(shape)
            if mermaid:
                collected.append(
                    {"kind": "mermaid", "shape_id": shape.shape_id, "text": f"```mermaid\n{mermaid}\n```", "top": pos_top, "left": pos_left}
                )
            else:
                collected.append(
                    {"kind": "smartart", "shape_id": shape.shape_id, "text": self._smartart_fallback_text(shape), "top": pos_top, "left": pos_left}
                )
            return

        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            # Mermaid 化に必要なメタ情報は重複テキストでも収集する
            shapes_meta.append(self._shape_meta(shape))
            # 同一スライド内で繰り返し出現する短文（凡例ラベル等）は本文化を抑制
            if self._matches_repeated_slide_text(shape):
                return
            default_bullet = self._is_body_placeholder(shape)
            text = self._convert_text_frame(shape.text_frame, default_bullet=default_bullet)
            if text:
                collected.append({"kind": "text", "shape_id": shape.shape_id, "text": text, "top": pos_top, "left": pos_left})
            return

    @staticmethod
    def _merge_horizontal_text_rows(items: list) -> list:
        """top が近い連続 text 要素を 1 行 (スペース区切り) として結合する。

        対象は kind == "text" のみ。画像・テーブル・mermaid 等は独立ブロックを維持。
        top の差が 250_000 EMU (約 0.25cm) 以内のものを同一行とみなす。
        """
        if not items:
            return items
        TOP_TOLERANCE = 250_000

        merged: list = []
        current_row: list = []
        for item in items:
            if item.get("kind") != "text":
                if current_row:
                    merged.append(PPTXMarkdownConverter._fuse_row(current_row))
                    current_row = []
                merged.append(item)
                continue
            if not current_row:
                current_row = [item]
                continue
            prev_top = current_row[-1].get("top") or 0
            cur_top = item.get("top") or 0
            if abs(cur_top - prev_top) <= TOP_TOLERANCE:
                current_row.append(item)
            else:
                merged.append(PPTXMarkdownConverter._fuse_row(current_row))
                current_row = [item]
        if current_row:
            merged.append(PPTXMarkdownConverter._fuse_row(current_row))
        return merged

    @staticmethod
    def _fuse_row(row_items: list) -> dict:
        """同一行に並ぶ複数 text item を 1 ブロックに結合."""
        if len(row_items) == 1:
            return row_items[0]
        row_items_sorted = sorted(row_items, key=lambda x: (x.get("left") or 0))
        joined = "　".join(
            (it.get("text") or "").replace("\n", " ").strip()
            for it in row_items_sorted
            if (it.get("text") or "").strip()
        )
        return {
            "kind": "text",
            "shape_id": None,
            "text": joined,
            "top": row_items_sorted[0].get("top"),
            "left": row_items_sorted[0].get("left"),
        }

    @staticmethod
    def _safe_position(shape) -> tuple:
        """shape の (top, left) を取得。取得失敗時は (0, 0) を返す."""
        try:
            top = shape.top or 0
        except Exception:
            top = 0
        try:
            left = shape.left or 0
        except Exception:
            left = 0
        return top, left

    @staticmethod
    def _is_title_shape(shape) -> bool:
        try:
            placeholder = shape.placeholder_format
            if placeholder is None:
                return False
            placeholder_type = placeholder.type
            type_str = str(placeholder_type) if placeholder_type is not None else ""
            if "SUBTITLE" in type_str:
                return False
            if placeholder.idx == 0:
                return True
            return "TITLE" in type_str
        except Exception:
            return False

    @staticmethod
    def _is_body_placeholder(shape) -> bool:
        try:
            placeholder = shape.placeholder_format
            if placeholder is None:
                return False
            placeholder_type = placeholder.type
            if placeholder_type is None:
                return False
            type_str = str(placeholder_type)
            return any(kw in type_str for kw in ("BODY", "OBJECT", "CONTENT"))
        except Exception:
            return False

    @staticmethod
    def _is_connector(shape) -> bool:
        try:
            return shape.element.tag.endswith("}cxnSp")
        except Exception:
            return False

    @staticmethod
    def _is_smartart(shape) -> bool:
        try:
            for child in shape.element.iter():
                if child.tag.endswith("}graphicData"):
                    uri = child.get("uri", "")
                    if "diagram" in uri:
                        return True
        except Exception:
            return False
        return False

    @staticmethod
    def _shape_meta(shape) -> dict:
        text = ""
        try:
            if shape.has_text_frame:
                text = _safe_text(shape.text_frame.text)
        except Exception:
            text = ""
        auto = ""
        try:
            auto = str(shape.auto_shape_type) if getattr(shape, "auto_shape_type", None) else ""
        except Exception:
            auto = ""
        font_size_max = PPTXMarkdownConverter._extract_max_font_size_pt(shape)
        font_color_hex = PPTXMarkdownConverter._extract_dominant_font_color(shape)
        return {
            "shape_id": shape.shape_id,
            "name": shape.name,
            "text": text,
            "auto_shape_type": auto,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "font_size_pt": font_size_max,
            "font_color": font_color_hex,
            "is_grayish": PPTXMarkdownConverter._is_grayish_color(font_color_hex),
            "char_count": len(text),
        }

    @staticmethod
    def _extract_max_font_size_pt(shape) -> Optional[float]:
        """shape 内の最大フォントサイズ (pt) を返す。取得できない場合は None."""
        try:
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                return None
        except Exception:
            return None
        max_pt: Optional[float] = None
        try:
            for paragraph in shape.text_frame.paragraphs:
                # paragraph レベルのフォントサイズも見る
                try:
                    p_size = paragraph.font.size
                    if p_size is not None:
                        pt_val = p_size.pt
                        if max_pt is None or pt_val > max_pt:
                            max_pt = pt_val
                except Exception:
                    pass
                for run in paragraph.runs:
                    try:
                        size = run.font.size
                        if size is None:
                            continue
                        pt_val = size.pt
                        if max_pt is None or pt_val > max_pt:
                            max_pt = pt_val
                    except Exception:
                        continue
        except Exception:
            return max_pt
        return max_pt

    @staticmethod
    def _extract_dominant_font_color(shape) -> Optional[str]:
        """shape 内の最初に取得できた RGB 色を hex 文字列で返す."""
        try:
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                return None
        except Exception:
            return None
        try:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        color = run.font.color
                        if color is None:
                            continue
                        rgb = color.rgb
                        if rgb is None:
                            continue
                        # rgb は RGBColor 型。str() で "RRGGBB" を返す
                        return str(rgb).upper()
                    except Exception:
                        continue
        except Exception:
            return None
        return None

    @staticmethod
    def _is_grayish_color(hex_color: Optional[str]) -> bool:
        """フォント色が薄いグレー / 装飾色っぽいかを判定."""
        if not hex_color or len(hex_color) != 6:
            return False
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
        except Exception:
            return False
        # 3 チャネルがいずれも 0x80 (128) 以上で、互いの差が小さい → 薄いグレー系
        if r < 0x80 or g < 0x80 or b < 0x80:
            return False
        max_c = max(r, g, b)
        min_c = min(r, g, b)
        return (max_c - min_c) <= 20

    def _collect_connector(self, shape, connectors: list) -> None:
        try:
            begin = None
            end = None
            for element in shape.element.iter():
                tag = element.tag
                if tag.endswith("}stCxn") and element.get("id"):
                    begin = _safe_int_from_str(element.get("id"))
                elif tag.endswith("}endCxn") and element.get("id"):
                    end = _safe_int_from_str(element.get("id"))
            if begin is not None and end is not None:
                connectors.append({"begin": begin, "end": end})
        except Exception:
            return

    # ---------- テキストフレーム ----------

    def _convert_text_frame(self, text_frame, *, default_bullet: bool = False) -> str:
        rendered_paragraphs: List[str] = []
        code_buffer: List[str] = []

        def flush_code():
            if code_buffer:
                rendered_paragraphs.append("```\n" + "\n".join(code_buffer) + "\n```")
                code_buffer.clear()

        for paragraph in text_frame.paragraphs:
            text = self._render_paragraph_runs(paragraph)
            if not text.strip():
                flush_code()
                continue

            level = paragraph.level or 0
            if self._is_monospace_paragraph(paragraph):
                # 装飾を取り除いた素のテキストを残す
                raw = "".join((run.text or "") for run in paragraph.runs)
                code_buffer.append(raw)
                continue

            flush_code()

            bullet_state = self._bullet_state(paragraph)
            is_bulleted = bullet_state is True or (bullet_state is None and default_bullet)

            if level > 0:
                indent = "  " * level
                rendered_paragraphs.append(f"{indent}- {text}")
            elif is_bulleted:
                rendered_paragraphs.append(f"- {text}")
            else:
                rendered_paragraphs.append(text)

        flush_code()
        return "\n".join(rendered_paragraphs)

    @staticmethod
    def _bullet_state(paragraph):
        """段落の bullet 設定を判定する。

        戻り値:
            True  -- 段落に明示的な bullet 指定がある (buChar / buAutoNum)
            False -- 段落に明示的な「bullet なし」指定がある (buNone)
            None  -- 段落側に指定がなく、レイアウト側のデフォルトに従う
        """
        try:
            p_pr = paragraph._pPr
            if p_pr is None:
                return None
            for child in p_pr.iter():
                tag = child.tag
                if tag.endswith("}buChar") or tag.endswith("}buAutoNum"):
                    return True
                if tag.endswith("}buNone"):
                    return False
        except Exception:
            return None
        return None

    @staticmethod
    def _render_paragraph_runs(paragraph) -> str:
        rendered_parts: List[str] = []
        for run in paragraph.runs:
            text = run.text or ""
            if not text:
                continue
            bold = bool(run.font.bold)
            italic = bool(run.font.italic)
            strike = False
            try:
                r_pr = run._r.find(f"{{{NS_A}}}rPr")
                if r_pr is not None:
                    strike_attr = r_pr.get("strike")
                    if strike_attr and strike_attr != "noStrike":
                        strike = True
            except Exception:
                strike = False
            wrapped = text
            if strike:
                wrapped = f"~~{wrapped}~~"
            if italic:
                wrapped = f"*{wrapped}*"
            if bold:
                wrapped = f"**{wrapped}**"
            rendered_parts.append(wrapped)
        return "".join(rendered_parts)

    @staticmethod
    def _is_monospace_paragraph(paragraph) -> bool:
        runs = list(paragraph.runs)
        if not runs:
            return False
        for run in runs:
            try:
                name = run.font.name
            except Exception:
                name = None
            if not _is_monospace_font(name):
                return False
        return True

    # ---------- 表 ----------

    def _convert_table(self, table) -> str:
        rows = list(table.rows)
        if not rows:
            return ""
        header_cells = [self._cell_text(cell) for cell in rows[0].cells]
        col_count = len(header_cells)
        if col_count == 0:
            return ""
        # 全セルが空のテーブルはレイアウト装飾とみなして抑制
        all_cells_empty = True
        for row in rows:
            for cell in row.cells:
                if self._cell_text(cell).strip():
                    all_cells_empty = False
                    break
            if not all_cells_empty:
                break
        if all_cells_empty:
            return ""
        lines = [
            "| " + " | ".join(_escape_md_pipe(c) for c in header_cells) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]
        for row in rows[1:]:
            cells = [self._cell_text(cell) for cell in row.cells]
            while len(cells) < col_count:
                cells.append("")
            lines.append(
                "| " + " | ".join(_escape_md_pipe(c) for c in cells[:col_count]) + " |"
            )
        return "\n".join(lines)

    @staticmethod
    def _cell_text(cell) -> str:
        if cell is None:
            return ""
        return _safe_text(cell.text or "")

    # ---------- 画像 ----------

    def _handle_picture(self, shape, slide_no: int) -> Optional[str]:
        try:
            blob = shape.image.blob
        except Exception as exc:
            print(
                f"Warning: failed to extract image on slide {slide_no}: {exc}",
                file=sys.stderr,
            )
            return None

        size = len(blob)
        if size > self.max_image_bytes:
            print(
                f"Warning: image too large ({size} bytes > {self.max_image_bytes}); "
                f"slide {slide_no} shape '{shape.name}'",
                file=sys.stderr,
            )
            return (
                f"> 画像（サイズ超過のためバイナリ非保存）: slide={slide_no}, "
                f"name={shape.name}, size={size}B"
            )

        raw_ext = (shape.image.ext or "png").lower()
        ext = raw_ext if raw_ext in ALLOWED_IMAGE_EXTS else "bin"
        if ext != raw_ext:
            print(
                f"Warning: image extension '{raw_ext}' is not in allowlist; "
                f"normalized to 'bin' on slide {slide_no}",
                file=sys.stderr,
            )
        self._image_seq.setdefault(slide_no, 0)
        self._image_seq[slide_no] += 1
        img_no = self._image_seq[slide_no]
        filename = f"slide{slide_no}_img{img_no}.{ext}"
        # MR-6: 画像総量・枚数の上限チェック（CWE-400 / CWE-770）
        if self._image_total_count >= MAX_IMAGE_COUNT_PER_PPTX:
            raise ValueError(
                f"image count exceeds MAX_IMAGE_COUNT_PER_PPTX ({MAX_IMAGE_COUNT_PER_PPTX})"
            )
        if self._image_total_bytes + len(blob) > MAX_TOTAL_IMAGE_BYTES:
            raise ValueError(
                f"total image bytes would exceed MAX_TOTAL_IMAGE_BYTES ({MAX_TOTAL_IMAGE_BYTES})"
            )
        out_path = self.images_dir / filename
        try:
            _safe_write_bytes(out_path, blob, force=self._force)
            self._image_total_bytes += len(blob)
            self._image_total_count += 1
        except ValueError as exc:
            print(f"Warning: failed to write image (path issue): {exc}", file=sys.stderr)
            return None
        except OSError as exc:
            print(f"Warning: failed to write image: {exc}", file=sys.stderr)
            return None

        alt = _escape_md_link_text(self._image_alt(shape, img_no))
        try:
            rel = out_path.relative_to(self.output_path.parent)
        except ValueError:
            rel = out_path
        return f"![{alt}]({rel.as_posix()})"

    @staticmethod
    def _image_alt(shape, img_no: int) -> str:
        """画像 shape の alt テキスト（descr 属性）を取得.

        python-pptx の private 属性 `_element` と lxml の `xpath()` 依存を避け、
        public API (`shape.description`) または `shape.element.iter()` で再帰探索する.
        戻り値は `_safe_text` で正規化し、Bidi 制御文字・サロゲート・過長文字列を除去する
        （HR-F: 後段 convert-html での XSS / homograph 偽装対策 + 一貫したサニタイズ）.
        """
        result = None
        try:
            # python-pptx 0.6.21+ の Public API（存在すれば最優先）
            descr = getattr(shape, "description", None)
            if descr:
                result = descr
            else:
                # element 属性（new public）または _element（fallback）で XML 要素を取得
                elem = getattr(shape, "element", None) or getattr(shape, "_element", None)
                if elem is not None:
                    for sub in elem.iter():
                        v = sub.get("descr")
                        if v:
                            result = v
                            break
        except Exception:
            pass
        if not result:
            result = shape.name or f"image{img_no}"
        # _safe_text で Bidi / サロゲート / 過長文字列を除去（HR-F）
        return _safe_text(result) or f"image{img_no}"

    # ---------- チャート ----------

    @staticmethod
    def _summarize_chart(chart) -> str:
        try:
            type_name = str(chart.chart_type)
        except Exception:
            type_name = "unknown"
        series_names: List[str] = []
        try:
            for series in chart.series:
                try:
                    series_names.append(series.name)
                except Exception:
                    continue
        except Exception:
            series_names = []
        return f"> チャート: {type_name} 系列={series_names}"

    # ---------- SmartArt ----------

    def _smartart_to_mermaid(self, shape) -> Optional[str]:
        try:
            ns = {"dgm": NS_DGM, "r": NS_R, "a": NS_A}
            rel_ids = shape.element.findall(".//dgm:relIds", namespaces=ns)
            if not rel_ids:
                return None
            rid_data = rel_ids[0].get(f"{{{NS_R}}}dm")
            if not rid_data:
                return None
            slide_part = shape.part
            rel = slide_part.rels.get(rid_data)
            if rel is None:
                return None
            data_xml = etree.fromstring(rel.target_part.blob, parser=_hardened_xml_parser())
            points: dict[str, str] = {}
            for point in data_xml.findall(".//dgm:pt", namespaces=ns):
                pid = point.get("modelId")
                if not pid:
                    continue
                texts = point.findall(".//a:t", namespaces=ns)
                value = " ".join((t.text or "") for t in texts).strip()
                points[pid] = value
            edges: List[tuple] = []
            for connection in data_xml.findall(".//dgm:cxn", namespaces=ns):
                src = connection.get("srcId")
                dst = connection.get("destId")
                ctype = connection.get("type", "parOf")
                if ctype not in ("parOf", "presOf"):
                    continue
                if src in points and dst in points:
                    edges.append((src, dst))
            if not edges:
                return None

            def safe_id(pid: str) -> str:
                cleaned = re.sub(r"[^A-Za-z0-9]", "", pid)
                return f"S{cleaned[:20]}" if cleaned else f"S{abs(hash(pid)) % 100000}"

            used = {src for src, _ in edges} | {dst for _, dst in edges}
            lines = ["flowchart TD"]
            for nid in sorted(used):
                label = _mermaid_escape_label(points.get(nid, "") or nid)
                lines.append(f'    {safe_id(nid)}["{label}"]')
            for src, dst in edges:
                lines.append(f"    {safe_id(src)} --> {safe_id(dst)}")
            return "\n".join(lines)
        except Exception as exc:
            print(f"Warning: failed to parse SmartArt: {exc}", file=sys.stderr)
            return None

    @staticmethod
    def _smartart_fallback_text(shape) -> str:
        try:
            text_nodes = shape.element.findall(f".//{{{NS_A}}}t")
            lines: List[str] = []
            for node in text_nodes:
                if node.text and node.text.strip():
                    lines.append(f"- {node.text.strip()}")
            if lines:
                return "> SmartArt（テキスト抽出）:\n" + "\n".join(lines)
        except Exception:
            pass
        return "> SmartArt（解析失敗・テキスト抽出不可）"

    # ---------- スピーカーノート ----------

    @staticmethod
    def _extract_notes(slide) -> Optional[str]:
        try:
            if not slide.has_notes_slide:
                return None
            text = _safe_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            return None
        if not text:
            return None
        lines = ["> [!NOTE]"]
        for line in text.split("\n"):
            lines.append(f"> {line}")
        return "\n".join(lines)


# ----------------------------------------------------------------------------- #
# エントリポイント
# ----------------------------------------------------------------------------- #


def parse_args(argv: Iterable[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert PPTX to Markdown.")
    parser.add_argument("input", help="Input PPTX file path")
    parser.add_argument(
        "output",
        nargs="?",
        default=None,
        help="Output MD file path (default: <input>.md)",
    )
    parser.add_argument(
        "--images-dir",
        default=None,
        help="Image extraction directory (default: <basename>_images/ next to output MD)",
    )
    parser.add_argument(
        "--no-mermaid",
        action="store_true",
        help="Disable flowchart/SmartArt Mermaid conversion",
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes",
    )
    parser.add_argument(
        "--include-hidden",
        action="store_true",
        help="Include hidden slides",
    )
    parser.add_argument(
        "--no-first-slide-as-title",
        action="store_true",
        help="Treat first slide as H2 instead of H1",
    )
    parser.add_argument(
        "--max-image-size",
        type=int,
        default=DEFAULT_MAX_IMAGE_BYTES,
        help="Maximum bytes per image (default 5 MiB)",
    )
    parser.add_argument(
        "--structured-json",
        default=None,
        help=(
            "Output a machine-readable structured JSON instead of (or in addition to) Markdown. "
            "If specified alone, only JSON is written. Use this for LLM-driven semantic conversion."
        ),
    )
    parser.add_argument(
        "--json-only",
        action="store_true",
        help="When set with --structured-json, suppress Markdown output and only emit JSON.",
    )
    parser.add_argument(
        "--per-slide-json",
        default=None,
        help=(
            "Emit one JSON file per slide into the given directory (slide-NN.json) plus "
            "metadata.json. Useful for large PPTX where a monolithic JSON exceeds the LLM "
            "context window."
        ),
    )
    parser.add_argument(
        "--compact-view",
        default=None,
        help=(
            "Emit one human/LLM-friendly text view per slide into the given directory "
            "(slide-NN.txt). Each shape is rendered as one row with id/position/font/text. "
            "Designed for Phase 2 (LLM interpretation) of large PPTX."
        ),
    )
    parser.add_argument(
        "--workspace-root",
        default=None,
        help=(
            "Workspace root directory for output path validation. All secondary output "
            "paths (--structured-json / --per-slide-json / --compact-view / --images-dir) "
            "must resolve under this directory to defeat path traversal (CWE-22). "
            "Default: parent directory of the resolved output MD path."
        ),
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help=(
            "Overwrite existing output files without confirmation. Without this flag, "
            "writing to an existing file aborts with an error (CWE-377 protection)."
        ),
    )
    return parser.parse_args(list(argv))


def main(argv: Optional[Iterable[str]] = None) -> int:
    args = parse_args(argv if argv is not None else sys.argv[1:])
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr, flush=True)
        sys.stderr.flush()
        return 1
    try:
        converter = PPTXMarkdownConverter(args)
        wrote_anything = False
        # 構造化 JSON モード（参考: PPTX → 機械抽出 JSON → LLM 解釈 → MD の 2 段階処理）
        if converter.structured_json_path is not None:
            json_path = converter.export_structured_json()
            print(f"Wrote JSON: {json_path}")
            wrote_anything = True
        # スライド分割 JSON（大規模 PPTX 向け）
        if converter.per_slide_json_dir is not None:
            count = converter.export_per_slide_json()
            print(f"Wrote per-slide JSON: {converter.per_slide_json_dir} ({count} slides)")
            wrote_anything = True
        # コンパクト ビュー（人間/LLM 可読の簡潔表現）
        if converter.compact_view_dir is not None:
            count = converter.export_compact_view()
            print(f"Wrote compact view: {converter.compact_view_dir} ({count} slides)")
            wrote_anything = True
        if wrote_anything:
            print(f"Images dir: {converter.images_dir}")
            if converter.json_only:
                return 0
        # JSON 系オプション未指定 or --json-only 無し → 従来の Markdown 直接出力
        output_path = converter.convert()
        print(f"Wrote: {output_path}")
        print(f"Images dir: {converter.images_dir}")
        return 0
    except FileNotFoundError as exc:
        print(f"Error: {exc}", file=sys.stderr, flush=True)
        sys.stderr.flush()
        return 1
    except ValueError as exc:
        print(f"Error: {exc}", file=sys.stderr, flush=True)
        sys.stderr.flush()
        return 1
    except Exception as exc:  # noqa: BLE001
        print(f"Error: unexpected failure: {exc}", file=sys.stderr, flush=True)
        sys.stderr.flush()
        return 2


if __name__ == "__main__":
    # PowerShell `Start-Process -RedirectStandardOutput/Error` 経由で起動された場合、
    # 通常の `sys.exit()` 経由（atexit / 通常 GC）では Python の stdio バッファ解放が
    # 親プロセスの `WaitForExit` で検知されず「ハング」として観測される
    # 既知の Windows + Python の挙動を回避するため、終了前に明示 flush し、
    # `os._exit()` で atexit を経由せず即時終了する.
    # 本スクリプトは threading / atexit / 一時ファイルクリーンアップを行わないため、
    # `os._exit()` での即時終了は安全.
    import os as _exit_os
    try:
        _rc = main()
    except SystemExit as _se:
        _rc = int(_se.code) if isinstance(_se.code, int) else (0 if _se.code is None else 1)
    try:
        sys.stdout.flush()
        sys.stderr.flush()
    except Exception:
        pass
    _exit_os._exit(_rc)
