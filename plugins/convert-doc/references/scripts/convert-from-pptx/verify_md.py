#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""verify_md.py - PPTX -> Markdown coverage verifier (Phase 3 validation).

PPTX とそこから生成された Markdown を機械的に突き合わせて、
情報の漏れ / 誤転記 / 装飾混入の疑いを検出する。

検出する観点:
    1. text coverage  : PPTX 内テキストが MD 内に出現するか (テンプレ装飾は除外候補)
    2. table coverage : PPTX 内テーブルセルが MD 内に出現するか
    3. image coverage : PPTX 内画像が MD 内 ![](...) として参照されているか
    4. connector coverage : PPTX 内コネクタが Mermaid edge として反映されているか
    5. fabrication detection : MD 内のテキストが PPTX のどこにも存在しないか (捏造防止)

出力: JSON 形式の検証レポート。閾値未達なら exit code 1 を返す (CI 連携用)。
"""
from __future__ import annotations

import argparse
import json
import re
import sys
import unicodedata
from pathlib import Path
from typing import Optional

try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except AttributeError:
    # Python 3.6 以下では reconfigure 未実装。文字化けリスクは
    # PYTHONUTF8=1 / PYTHONIOENCODING=utf-8 で補う前提.
    pass

# XML 攻撃対策（CWE-611 / CWE-776）は convert_from_pptx.py と対称化:
# - 本スクリプトは XML を直接 parse しないが、python-pptx 経由で lxml が動くため、
#   グローバル default parser を hardened に上書きする (convert_from_pptx.py と同じ).
# - PPTX サイズ・スライド数等の上限定数も併用する.
# - 旧 `defusedxml.lxml.monkey_patch_lxml()` 呼び出しは defusedxml 0.7 で API が削除された
#   ため撤去した（呼び出すと AttributeError → fail-close で起動不能）.
try:
    from lxml import etree as _lxml_etree
    _lxml_etree.set_default_parser(_lxml_etree.XMLParser(
        resolve_entities=False,
        no_network=True,
        load_dtd=False,
        huge_tree=False,
    ))
except ImportError:
    # lxml が無い環境（python-pptx の動作前提として通常はあるが、保険として）
    pass

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:  # pragma: no cover
    print(f"Error: python-pptx not available ({exc}).", file=sys.stderr, flush=True)
    sys.stderr.flush()
    raise SystemExit(2)


CTRL_CHAR_RE = re.compile(r"[\x09-\x1f\x7f-\x9f]+")
WS_RE = re.compile(r"\s+")
PHRASE_SPLIT_RE = re.compile(r"[。｡\n\r]|(?<=\S)\.\s+")
IMAGE_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
MERMAID_BLOCK_RE = re.compile(r"```mermaid\s*\n(.*?)\n```", re.DOTALL)
EDGE_RE = re.compile(r"-{2,}>|--+|==+>|==+|-\.->")

# HR-E / HR-ζ: DoS 防御の上限定数（convert_from_pptx.py と完全対称化）
MAX_SLIDES = 1000
MAX_SHAPES_PER_SLIDE = 5000
MAX_GROUP_DEPTH = 20
MAX_TOTAL_UNCOMPRESSED_BYTES = 256 * 1024 * 1024  # 256 MiB (ZIP bomb 防御)
MAX_COMPRESSION_RATIO = 200  # ZIP 圧縮率上限
MAX_TEXT_PER_SHAPE = 1_000_000  # 1 shape あたりのテキスト長上限（CWE-400）
MAX_TOTAL_IMAGE_BYTES = 256 * 1024 * 1024  # 256 MiB (画像総量、CWE-770)
MAX_IMAGE_COUNT_PER_PPTX = 1000  # PPTX 全体での画像枚数上限


_MD_SYNTAX_RE = re.compile(r"[#>|`*~\[\]]")
_MD_LINK_RE = re.compile(r"<(https?://[^>]+)>")
_MD_BR_RE = re.compile(r"<br\s*/?>", re.IGNORECASE)


def _check_input_path(path: Path, label: str) -> None:
    """HR-D: 入力パスの安全性検証（CWE-22 / CWE-59 対策）.

    - パス文字列に `..` 含む場合は拒否（パストラバーサル）
    - シンボリックリンクの場合は拒否（TOCTOU / リンク先攻撃対策）
    """
    raw = str(path).replace("\\", "/")
    if ".." in raw.split("/"):
        print(
            f"Error: {label} contains '..' (path traversal blocked): {path}",
            file=sys.stderr,
            flush=True,
        )
        sys.stderr.flush()
        raise SystemExit(2)
    if path.is_symlink():
        print(
            f"Error: {label} is a symbolic link (refused for safety): {path}",
            file=sys.stderr,
            flush=True,
        )
        sys.stderr.flush()
        raise SystemExit(2)


def _validate_pptx_for_verify(pptx_path: Path) -> None:
    """HR-E: verify_md.py 側の PPTX 入力検証（convert_from_pptx.py の _validate_pptx と同等）.

    ZIP マジック・総展開サイズ上限・圧縮率上限を確認し、Billion Laughs / ZIP bomb 攻撃を遮断する.
    """
    import zipfile
    with open(pptx_path, "rb") as fh:
        head = fh.read(4)
    if head != b"PK\x03\x04":
        raise ValueError(f"Not a PPTX (zip) file: {pptx_path}")
    with zipfile.ZipFile(pptx_path, "r") as zf:
        total_uncompressed = 0
        for info in zf.infolist():
            total_uncompressed += info.file_size
            if total_uncompressed > MAX_TOTAL_UNCOMPRESSED_BYTES:
                raise ValueError(
                    f"Input PPTX exceeds uncompressed size limit "
                    f"({total_uncompressed} > {MAX_TOTAL_UNCOMPRESSED_BYTES}): {pptx_path}"
                )
            if info.compress_size > 0:
                ratio = info.file_size / info.compress_size
                if ratio > MAX_COMPRESSION_RATIO:
                    raise ValueError(
                        f"Suspicious compression ratio ({ratio:.1f}x) "
                        f"for entry '{info.filename}' in {pptx_path}"
                    )


def _normalize(text: str) -> str:
    """テキストを比較用に正規化する (NFKC + 制御文字除去 + 連続空白縮約 + 小文字化)."""
    if not text:
        return ""
    t = unicodedata.normalize("NFKC", text)
    t = CTRL_CHAR_RE.sub(" ", t)
    t = WS_RE.sub(" ", t)
    return t.strip().lower()


def _strip_md_syntax(text: str) -> str:
    """Markdown 構文記号 (#, >, |, *, ``, _, ~, [, ]) を除去して比較しやすくする."""
    if not text:
        return ""
    t = _MD_SYNTAX_RE.sub(" ", text)
    return _normalize(t)


def _split_into_phrases(text: str) -> list:
    if not text:
        return []
    parts = PHRASE_SPLIT_RE.split(text)
    return [_normalize(p) for p in parts if _normalize(p)]


# --------------------------------------------------------------------------- #
# PPTX 抽出
# --------------------------------------------------------------------------- #


def _is_template_placeholder(shape) -> bool:
    try:
        ph = shape.placeholder_format
        if ph is None or ph.type is None:
            return False
        type_str = str(ph.type)
        return any(kw in type_str for kw in ("FOOTER", "SLIDE_NUMBER", "DATE"))
    except Exception:
        return False


def _is_connector(shape) -> bool:
    try:
        return shape.element.tag.endswith("}cxnSp")
    except Exception:
        return False


def _connector_info(shape) -> Optional[dict]:
    try:
        begin = None
        end = None
        for element in shape.element.iter():
            tag = element.tag
            if tag.endswith("}stCxn") and element.get("id"):
                begin = int(element.get("id"))
            elif tag.endswith("}endCxn") and element.get("id"):
                end = int(element.get("id"))
        return {"begin": begin, "end": end}
    except Exception:
        return None


def _collect_template_texts(presentation) -> set:
    texts: set = set()

    def _gather(container):
        try:
            for shape in container.shapes:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    t = (shape.text_frame.text or "").strip()
                    if t:
                        texts.add(t)
        except Exception:
            pass

    try:
        for master in presentation.slide_masters:
            _gather(master)
            for layout in master.slide_layouts:
                _gather(layout)
    except Exception:
        pass
    return texts


def _is_offscreen(shape, slide_w: Optional[int], slide_h: Optional[int]) -> bool:
    """shape がスライド領域外 (top<0, left<0, top>height, left>width) なら True."""
    try:
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
    except Exception:
        return False
    if top < 0 or left < 0:
        return True
    if slide_h is not None and top >= slide_h:
        return True
    if slide_w is not None and left >= slide_w:
        return True
    return False


def _walk(shape_iter, texts: list, tables: list, images: list, connectors: list,
          slide_w=None, slide_h=None, depth: int = 0):
    """HR-E: グループネスト深度と shape 数の上限チェックを追加（CWE-400 / CWE-674）."""
    if depth > MAX_GROUP_DEPTH:
        raise ValueError(
            f"Group nesting exceeds MAX_GROUP_DEPTH ({MAX_GROUP_DEPTH})"
        )
    for shape in shape_iter:
        if len(texts) + len(tables) + len(images) >= MAX_SHAPES_PER_SLIDE:
            raise ValueError(
                f"shape count exceeds MAX_SHAPES_PER_SLIDE ({MAX_SHAPES_PER_SLIDE})"
            )
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk(shape.shapes, texts, tables, images, connectors, slide_w, slide_h, depth + 1)
                continue
        except ValueError:
            raise
        except Exception:
            pass
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.LINE or _is_connector(shape):
                ci = _connector_info(shape)
                if ci:
                    connectors.append(ci)
        except Exception:
            pass
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({"shape_id": shape.shape_id, "name": shape.name})
                continue
        except Exception:
            pass
        try:
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([(cell.text or "").strip() for cell in row.cells])
                tables.append({"shape_id": shape.shape_id, "rows": rows})
                continue
        except Exception:
            pass
        try:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    # HR-ζ: テキスト長上限による DoS 防御（CWE-400）
                    if len(t) > MAX_TEXT_PER_SHAPE:
                        t = t[:MAX_TEXT_PER_SHAPE]
                    texts.append({
                        "shape_id": shape.shape_id,
                        "text": t,
                        "is_template_placeholder": _is_template_placeholder(shape),
                        "is_offscreen": _is_offscreen(shape, slide_w, slide_h),
                    })
        except Exception:
            pass


def _extract_pptx_inventory(pptx_path: Path) -> dict:
    # HR-E: PPTX マジック + ZIP bomb / 圧縮率の事前検査（convert_from_pptx.py と同等）
    _validate_pptx_for_verify(pptx_path)

    prs = Presentation(str(pptx_path))
    template_texts = _collect_template_texts(prs)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    slides_inv = []
    for idx, slide in enumerate(prs.slides, start=1):
        # HR-E: スライド数の上限チェック（DoS 防御・CWE-400）
        if idx > MAX_SLIDES:
            raise ValueError(f"slide count exceeds MAX_SLIDES ({MAX_SLIDES})")
        texts: list = []
        tables: list = []
        images: list = []
        connectors: list = []
        _walk(slide.shapes, texts, tables, images, connectors, slide_w, slide_h)
        slides_inv.append({
            "slide_no": idx,
            "texts": texts,
            "tables": tables,
            "images": images,
            "connectors": connectors,
        })

    return {
        "slide_count": len(slides_inv),
        "template_decoration_texts": sorted(template_texts),
        "slides": slides_inv,
    }


# --------------------------------------------------------------------------- #
# Markdown 抽出
# --------------------------------------------------------------------------- #


def _extract_md_features(md_path: Path) -> dict:
    with open(md_path, "r", encoding="utf-8") as fh:
        body = fh.read()
    # 比較用: Markdown のリンク表記 <https://...> を URL に展開 → <br/> 除去 → 記号空白化
    body_for_compare = _MD_LINK_RE.sub(r" \1 ", body)
    body_for_compare = _MD_BR_RE.sub(" ", body_for_compare)
    body_for_compare = _MD_SYNTAX_RE.sub(" ", body_for_compare)
    normalized_body = _normalize(body_for_compare)

    images = [{"alt": m.group(1), "path": m.group(2)} for m in IMAGE_RE.finditer(body)]

    mermaid_blocks = MERMAID_BLOCK_RE.findall(body)
    edges = 0
    for block in mermaid_blocks:
        for line in block.splitlines():
            line = line.strip()
            if not line:
                continue
            if line.startswith(("subgraph", "end", "flowchart", "graph")):
                continue
            if EDGE_RE.search(line):
                edges += 1

    # 句点分割した phrase（Markdown 記号除去済み）
    phrases = set(_split_into_phrases(body_for_compare))

    return {
        "raw": body,
        "normalized": normalized_body,
        "images": images,
        "mermaid_blocks_count": len(mermaid_blocks),
        "mermaid_edges_count": edges,
        "phrases": phrases,
    }


# --------------------------------------------------------------------------- #
# 検証ロジック
# --------------------------------------------------------------------------- #


def _text_present_in_md(needle: str, md_norm: str, md_phrases: set) -> bool:
    needle_n = _normalize(needle)
    if not needle_n:
        return True
    if needle_n in md_phrases:
        return True
    if needle_n in md_norm:
        return True
    parts = _split_into_phrases(needle)
    if parts and all(p in md_norm for p in parts):
        return True
    return False


def _blocks_present_in_md(needle: str, md_norm: str, md_phrases: set) -> bool:
    """改行で分割した各断片が MD 内に存在するかを ALL チェックする。

    PPTX 内の段落/セルは改行で論理ブロック化されているが、Phase 2 で
    Claude が構造変換（箇条書き化・テーブル分解）すると断片が別位置に
    再配置される。各断片レベルでの存在チェックが現実的な検証。
    """
    if not needle:
        return True
    if _text_present_in_md(needle, md_norm, md_phrases):
        return True
    # 改行・VT 等で分割
    parts = re.split(r"[\n\r\x0b\x0c]+", needle)
    parts = [p.strip() for p in parts if p.strip()]
    if not parts:
        return True
    return all(_text_present_in_md(p, md_norm, md_phrases) for p in parts)


def verify(pptx_path: Path, md_path: Path, coverage_threshold: float = 0.85) -> dict:
    inv = _extract_pptx_inventory(pptx_path)
    md = _extract_md_features(md_path)
    template_set = {_normalize(t) for t in inv["template_decoration_texts"]}

    text_total = 0
    text_present = 0
    text_template_excluded = 0
    text_offscreen_excluded = 0
    text_missing: list = []
    for s in inv["slides"]:
        for t in s["texts"]:
            text_norm = _normalize(t["text"])
            if not text_norm:
                continue
            if t["is_template_placeholder"] or text_norm in template_set:
                text_template_excluded += 1
                continue
            # 画面外 shape (top<0 等) はテンプレ装飾/レイアウト補助としてスコープ外
            if t.get("is_offscreen"):
                text_offscreen_excluded += 1
                continue
            text_total += 1
            if _blocks_present_in_md(t["text"], md["normalized"], md["phrases"]):
                text_present += 1
            else:
                text_missing.append({
                    "slide_no": s["slide_no"],
                    "shape_id": t["shape_id"],
                    "text": t["text"][:120],
                })

    text_coverage = (text_present / text_total) if text_total else 1.0

    table_total = 0
    table_present = 0
    table_missing: list = []
    for s in inv["slides"]:
        for tbl in s["tables"]:
            for r_idx, row in enumerate(tbl["rows"]):
                for c_idx, cell in enumerate(row):
                    cell_n = _normalize(cell)
                    if not cell_n:
                        continue
                    table_total += 1
                    if _blocks_present_in_md(cell, md["normalized"], md["phrases"]):
                        table_present += 1
                    else:
                        table_missing.append({
                            "slide_no": s["slide_no"],
                            "shape_id": tbl["shape_id"],
                            "row": r_idx,
                            "col": c_idx,
                            "cell": cell[:120],
                        })
    table_coverage = (table_present / table_total) if table_total else 1.0

    pptx_image_total = sum(len(s["images"]) for s in inv["slides"])
    md_image_total = len(md["images"])

    pptx_conn_total = sum(
        sum(1 for c in s["connectors"] if c.get("begin") is not None and c.get("end") is not None)
        for s in inv["slides"]
    )
    mermaid_edge_total = md["mermaid_edges_count"]

    pptx_all_text_norm_parts = []
    for s in inv["slides"]:
        for t in s["texts"]:
            pptx_all_text_norm_parts.append(_normalize(t["text"]))
        for tbl in s["tables"]:
            for row in tbl["rows"]:
                for cell in row:
                    pptx_all_text_norm_parts.append(_normalize(cell))
    # HR-1: 隣接テキスト間に NUL を挟むことで境界消失による誤マッチを防ぐ.
    # 例: ["終了", "開始"] を " ".join すると "終了 開始" となり "了開" が誤ヒットしうるが、
    # NUL 挟みなら "\x00終了\x00開始\x00" となり連続マッチが起きない（cmp_phrase は通常 NUL を含まない）.
    _SENTINEL_BOUNDARY = "\x00"
    pptx_all_text_norm = _SENTINEL_BOUNDARY + _SENTINEL_BOUNDARY.join(pptx_all_text_norm_parts) + _SENTINEL_BOUNDARY

    suspicious_md_phrases: list = []
    for phrase in md["phrases"]:
        if len(phrase) < 16:
            continue
        if any(tok in phrase for tok in ("![", "```", "http", "->", "flowchart", "subgraph", "<br")):
            continue
        cmp_phrase = phrase.replace("*", "").replace("`", "").strip()
        if len(cmp_phrase) < 16:
            continue
        if cmp_phrase not in pptx_all_text_norm:
            suspicious_md_phrases.append(cmp_phrase[:160])

    report = {
        "input_pptx": str(pptx_path),
        "output_md": str(md_path),
        "summary": {
            "text_coverage": round(text_coverage, 4),
            "text_total": text_total,
            "text_present": text_present,
            "text_template_excluded": text_template_excluded,
            "text_offscreen_excluded": text_offscreen_excluded,
            "table_cell_coverage": round(table_coverage, 4),
            "table_total": table_total,
            "table_present": table_present,
            "pptx_image_total": pptx_image_total,
            "md_image_total": md_image_total,
            "pptx_connector_total": pptx_conn_total,
            "mermaid_edge_total": mermaid_edge_total,
            "suspicious_md_phrase_count": len(suspicious_md_phrases),
            "coverage_threshold": coverage_threshold,
        },
        "missing_texts": text_missing,
        "missing_table_cells": table_missing,
        "suspicious_md_phrases": suspicious_md_phrases,
    }

    failures = []
    if text_coverage < coverage_threshold:
        failures.append(f"text_coverage {text_coverage:.2%} < threshold {coverage_threshold:.0%}")
    if table_coverage < coverage_threshold:
        failures.append(f"table_cell_coverage {table_coverage:.2%} < threshold {coverage_threshold:.0%}")
    if pptx_image_total and md_image_total < pptx_image_total * 0.5:
        failures.append(f"images: pptx={pptx_image_total} md={md_image_total} (less than 50%)")
    if pptx_conn_total >= 5 and mermaid_edge_total == 0:
        failures.append(f"connectors: pptx={pptx_conn_total} but no Mermaid edges in MD")
    report["failures"] = failures
    report["passed"] = not failures
    return report


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(description="Verify PPTX -> Markdown conversion coverage.")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("md", help="Generated Markdown file")
    parser.add_argument("--report", default=None, help="Write JSON report to this path")
    parser.add_argument("--threshold", type=float, default=0.85, help="Coverage threshold (0.0-1.0)")
    parser.add_argument("--max-missing-shown", type=int, default=20, help="Truncate missing lists in console output")
    args = parser.parse_args(argv if argv is not None else sys.argv[1:])

    pptx_path = Path(args.pptx)
    md_path = Path(args.md)
    # HR-D: パストラバーサル / シンボリックリンク拒否（CWE-22 / CWE-59）
    _check_input_path(pptx_path, "pptx")
    _check_input_path(md_path, "md")
    if not pptx_path.exists():
        print(f"Error: PPTX not found: {pptx_path}", file=sys.stderr, flush=True)
        sys.stderr.flush()
        return 2
    if not md_path.exists():
        print(f"Error: Markdown not found: {md_path}", file=sys.stderr, flush=True)
        sys.stderr.flush()
        return 2

    report = verify(pptx_path, md_path, coverage_threshold=args.threshold)

    if args.report:
        report_path = Path(args.report)
        # L-12: パストラバーサル `..` を含む生パスは拒否（CWE-22 簡易対策）
        if ".." in str(report_path).replace("\\", "/").split("/"):
            print(
                f"Error: --report path contains '..' (path traversal blocked): {report_path}",
                file=sys.stderr,
                flush=True,
            )
            sys.stderr.flush()
            return 2
        # HR-ε: 既存ファイルがシンボリックリンクなら拒否（CWE-59 / CWE-367 TOCTOU 対策）
        if report_path.exists() and report_path.is_symlink():
            print(
                f"Error: --report path is a symbolic link (refused for safety): {report_path}",
                file=sys.stderr,
                flush=True,
            )
            sys.stderr.flush()
            return 2
        report_path = report_path.resolve()
        report_path.parent.mkdir(parents=True, exist_ok=True)
        with open(report_path, "w", encoding="utf-8", newline="\n") as fh:
            json.dump(report, fh, ensure_ascii=False, indent=2)
        try:
            report_path.chmod(0o600)
        except (OSError, NotImplementedError):
            pass
        print(f"Wrote report: {report_path}")

    s = report["summary"]
    print()
    print("=== Coverage Report ===")
    print(f"  Text coverage:        {s['text_coverage']:.2%} ({s['text_present']}/{s['text_total']})")
    print(f"  Table cell coverage:  {s['table_cell_coverage']:.2%} ({s['table_present']}/{s['table_total']})")
    print(f"  Images:               PPTX={s['pptx_image_total']}, MD={s['md_image_total']}")
    print(f"  Connectors / Edges:   PPTX={s['pptx_connector_total']}, MD={s['mermaid_edge_total']}")
    print(f"  Template excluded:    {s['text_template_excluded']} text shapes")
    print(f"  Offscreen excluded:   {s['text_offscreen_excluded']} text shapes")
    print(f"  Suspicious phrases:   {s['suspicious_md_phrase_count']}")
    print()

    if report["missing_texts"]:
        print(f"Missing texts (first {args.max_missing_shown}):")
        for m in report["missing_texts"][: args.max_missing_shown]:
            print(f"  slide {m['slide_no']} shape_id={m['shape_id']}: {m['text']!r}")
        print()

    if report["missing_table_cells"]:
        print(f"Missing table cells (first {args.max_missing_shown}):")
        for m in report["missing_table_cells"][: args.max_missing_shown]:
            print(f"  slide {m['slide_no']} [{m['row']},{m['col']}]: {m['cell']!r}")
        print()

    if report["suspicious_md_phrases"]:
        print(f"Suspicious MD phrases not found in PPTX (first {args.max_missing_shown}):")
        for p in report["suspicious_md_phrases"][: args.max_missing_shown]:
            print(f"  {p!r}")
        print()

    if report["passed"]:
        print("RESULT: PASSED")
        return 0
    print("RESULT: FAILED")
    for f in report["failures"]:
        print(f"  - {f}")
    return 1


if __name__ == "__main__":
    sys.exit(main())
