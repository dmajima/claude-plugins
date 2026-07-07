#!/usr/bin/env python3
"""
convert_pptx.py - Markdown to PowerPoint (PPTX) converter

Usage:
  python convert_pptx.py <input.md> [output.pptx]
    [--title TITLE] [--subtitle SUB] [--aspect 16:9|4:3]
    [--primary-color "#003879"] [--max-body-chars 2400]
    [--theme theme.json] [--dump-default-theme]

Dependencies: python-pptx, Pillow, requests
"""

import argparse
import base64
import io
import json
import math
import re
import sys
from dataclasses import dataclass, field, fields, replace
from pathlib import Path
from typing import List, Optional
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

MERMAID_PNG_ENDPOINT = "https://mermaid.ink/img/{}?type=png"

# Cap for remote fetches (mermaid PNG / linked images) to avoid memory-exhaustion
# from slow or oversized responses.
MAX_REMOTE_BYTES = 20 * 1024 * 1024


def _read_capped(resp, label: str) -> Optional[bytes]:
    """Read a streamed requests response up to MAX_REMOTE_BYTES; None if exceeded."""
    chunks = []
    total = 0
    for chunk in resp.iter_content(chunk_size=65536):
        total += len(chunk)
        if total > MAX_REMOTE_BYTES:
            print(
                f"Warning: remote resource exceeds {MAX_REMOTE_BYTES // (1024 * 1024)}MB cap, "
                f"skipped: {label}",
                file=sys.stderr,
            )
            return None
        chunks.append(chunk)
    return b"".join(chunks)


def _is_external_path_ref(src: str) -> bool:
    """True for absolute / UNC / drive-anchored paths that must never be joined
    onto base_dir (they would escape it or trigger network filesystem lookups).
    Mirrored in convert-html/convert.py — keep both copies in sync.
    """
    from pathlib import PurePosixPath, PureWindowsPath
    return bool(
        src.startswith(("\\\\", "//"))
        or PureWindowsPath(src).drive
        or PureWindowsPath(src).root
        or PurePosixPath(src).root
    )


# ===================== Composition (cover / content-header layout) =====================

# Width / height tokens resolved at draw time against the actual slide size:
#   "full" = slide width (w) / slide height (h)
#   "sym"  = slide width - x * 2 (symmetric side margins; w only)
# Slide width depends on the aspect ratio, so edge-reaching sizes must be
# written as tokens rather than absolute numbers to stay portable.

@dataclass(frozen=True)
class ShapeSpec:
    """A filled, borderless rectangle in a slide composition."""
    x: float
    y: float
    w: "float | str"        # number | "full" | "sym"
    h: "float | str"        # number | "full"
    color: str              # color token name (a `colors` section key) or hex


@dataclass(frozen=True)
class TextSpec:
    """Placement and text styling for a composition text box.

    Font face and size intentionally do NOT live here: composition describes
    geometry only, while `fonts` / `font_sizes_pt` keep describing the look
    (cover title -> title_slide_title, cover subtitle -> title_slide_subtitle,
    content-header title -> title_band).
    """
    x: float
    y: float
    w: "float | str"        # number | "sym" ("full" is used only internally)
    h: float
    color: str              # color token name or hex
    bold: bool = False
    align: str = "left"     # left | center | right
    anchor: str = "top"     # top | middle
    margin: float = 0.1     # text_frame left/right margin in inches (pptx default)


@dataclass(frozen=True)
class CoverComposition:
    """Title-slide layout: decoration shapes plus title/subtitle placement."""
    shapes: "tuple[ShapeSpec, ...]"
    title: TextSpec
    subtitle: TextSpec


@dataclass(frozen=True)
class ContentHeaderComposition:
    """Content-slide header layout plus where the body blocks start."""
    shapes: "tuple[ShapeSpec, ...]"
    title: TextSpec
    content_top: float      # body start Y (inches); also drives slide chunking


@dataclass(frozen=True)
class Composition:
    """Declarative cover / content-header layout.

    A part left as None falls back to the built-in default composition
    (parts are replaced wholesale; no deep merge within a part).
    """
    cover: "CoverComposition | None" = None
    content_header: "ContentHeaderComposition | None" = None


def build_default_composition(theme: "Theme", slide_w_in: float, slide_h_in: float) -> Composition:
    """Build the default composition (visually identical to the historical
    hard-coded layout).

    This function is the SSOT for the default composition; the reference
    listing in add-design-pptx's theme-schema.md mirrors it and is kept in
    sync by check_default_composition.py. It tracks theme values
    (title_band_height) dynamically, which is why --dump-default-theme does
    not include a `composition` section.
    """
    band_h = theme.title_band_height_in
    return Composition(
        cover=CoverComposition(
            shapes=(ShapeSpec(x=0, y=0, w=0.4, h="full", color="primary"),),
            title=TextSpec(
                x=1.0, y=2.3, w=slide_w_in - 1.5, h=2.0, color="primary", bold=True,
            ),
            subtitle=TextSpec(
                x=1.0, y=4.4, w=slide_w_in - 1.5, h=1.5, color="text",
            ),
        ),
        content_header=ContentHeaderComposition(
            shapes=(ShapeSpec(x=0, y=0, w="full", h=band_h, color="primary"),),
            # x=0 / w="full" / margin=0.35 reproduces the historical band text
            # exactly: the text frame spans the slide and the 0.35in margins
            # give the same effective text start X (0.35) as the old
            # text-on-band-shape implementation.
            title=TextSpec(
                x=0, y=0, w="full", h=band_h, color="on_primary",
                bold=True, anchor="middle", margin=0.35,
            ),
            content_top=band_h + 0.2,
        ),
    )


# ===================== Theme =====================

def _default_syntax_palette() -> "dict[str, RGBColor]":
    """Pygments syntax highlight palette (aligned with the HTML friendly theme)."""
    return {
        "keyword":   RGBColor(0x00, 0x7B, 0x83),   # Token.Keyword*
        "builtin":   RGBColor(0x5C, 0x35, 0x66),   # Token.Name.Builtin*
        "func":      RGBColor(0x00, 0x4E, 0xB0),   # Token.Name.Function / Token.Name.Class
        "tag":       RGBColor(0x00, 0x7B, 0x83),   # Token.Name.Tag
        "attr":      RGBColor(0xAA, 0x55, 0x00),   # Token.Name.Attribute
        "string":    RGBColor(0x4E, 0x95, 0x2A),   # Token.Literal.String*
        "number":    RGBColor(0xAA, 0x55, 0x00),   # Token.Literal.Number*
        "operator":  RGBColor(0x33, 0x33, 0x33),
        "comment":   RGBColor(0x80, 0x80, 0x80),   # Token.Comment*
        "error":     RGBColor(0xB7, 0x25, 0x25),
        "heading":   RGBColor(0x00, 0x4E, 0xB0),   # Token.Generic.Heading*
    }


@dataclass(frozen=True)
class Theme:
    """A complete visual design for the generated deck.

    Field defaults ARE the default design; ``--theme`` JSON files override
    them partially. Keep the defaults byte-compatible with historical output.
    """
    # --- colors ---
    primary: RGBColor = RGBColor(0x00, 0x38, 0x79)          # primary navy
    accent: RGBColor = RGBColor(0x1D, 0x6F, 0xD1)           # reserved for future use
    text: RGBColor = RGBColor(0x1F, 0x2D, 0x3D)
    on_primary: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)       # text on primary fills
    code_bg: RGBColor = RGBColor(0xF5, 0xF6, 0xF8)
    code_text: RGBColor = RGBColor(0x1F, 0x2D, 0x3D)
    code_border: RGBColor = RGBColor(0xDD, 0xE1, 0xE8)
    hr_color: RGBColor = RGBColor(0xC9, 0xD0, 0xD8)
    table_row_odd: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)
    table_row_even: RGBColor = RGBColor(0xF5, 0xF6, 0xF8)
    # --- fonts ---
    body_font: str = "Yu Gothic UI"
    heading_font: str = "Yu Gothic UI"
    code_font: str = "Consolas"
    # --- font sizes (pt) ---
    title_band_size_pt: float = 24
    title_slide_title_size_pt: float = 40
    title_slide_subtitle_size_pt: float = 18
    body_size_pt: float = 16
    heading_h3_size_pt: float = 22
    list_size_pt: float = 15
    code_size_pt: float = 11
    table_size_pt: float = 12
    # --- layout (inches) ---
    title_band_height_in: float = 0.9
    content_padding_in: float = 0.5
    mermaid_max_width_in: float = 11.5
    mermaid_max_height_in: float = 5.5
    image_max_width_in: float = 11.5
    image_max_height_in: float = 5.5
    # --- syntax highlight palette ---
    syntax_palette: "dict[str, RGBColor]" = field(default_factory=_default_syntax_palette)
    # --- composition (cover / content-header layout structure) ---
    # None = built-in default composition (which tracks title_band_height).
    composition: "Composition | None" = None


# JSON section -> (json_key, Theme field) mapping tables.
_THEME_COLOR_KEYS = {
    "primary": "primary",
    "accent": "accent",
    "text": "text",
    "on_primary": "on_primary",
    "code_bg": "code_bg",
    "code_text": "code_text",
    "code_border": "code_border",
    "hr": "hr_color",
    "table_row_odd": "table_row_odd",
    "table_row_even": "table_row_even",
}
_THEME_FONT_KEYS = {
    "body": "body_font",
    "heading": "heading_font",
    "code": "code_font",
}
_THEME_FONT_SIZE_KEYS = {
    "title_band": "title_band_size_pt",
    "title_slide_title": "title_slide_title_size_pt",
    "title_slide_subtitle": "title_slide_subtitle_size_pt",
    "body": "body_size_pt",
    "heading_h3": "heading_h3_size_pt",
    "list": "list_size_pt",
    "code": "code_size_pt",
    "table": "table_size_pt",
}
_THEME_LAYOUT_KEYS = {
    "title_band_height": "title_band_height_in",
    "content_padding": "content_padding_in",
    "mermaid_max_width": "mermaid_max_width_in",
    "mermaid_max_height": "mermaid_max_height_in",
    "image_max_width": "image_max_width_in",
    "image_max_height": "image_max_height_in",
}
_THEME_META_KEYS = ("name", "description")

# Self-check: every Theme field must be reachable from exactly one mapping
# table (plus syntax_palette / composition, handled separately). Fails at
# import time if a new Theme field is added without updating the tables, so
# the omission can never silently ship (it would otherwise be invisible in
# --dump-default-theme and impossible to override from a theme JSON).
_MAPPED_THEME_FIELDS = (
    set(_THEME_COLOR_KEYS.values()) | set(_THEME_FONT_KEYS.values())
    | set(_THEME_FONT_SIZE_KEYS.values()) | set(_THEME_LAYOUT_KEYS.values())
    | {"syntax_palette", "composition"}
)
assert _MAPPED_THEME_FIELDS == {f.name for f in fields(Theme)}, (
    "theme mapping tables are out of sync with Theme fields: "
    f"unmapped={sorted({f.name for f in fields(Theme)} - _MAPPED_THEME_FIELDS)} "
    f"unknown={sorted(_MAPPED_THEME_FIELDS - {f.name for f in fields(Theme)})}"
)


def _rgb_to_hex(color: RGBColor) -> str:
    return f"#{color}"


def _composition_to_dict(comp: Composition) -> dict:
    """Serialize a Composition back to the documented JSON shape.

    Used for round-trip completeness when a theme carries a custom
    composition. The default theme never carries one (composition is None),
    so --dump-default-theme stays composition-free: the default composition
    tracks theme values such as title_band_height dynamically, and a dumped
    static copy would silently desync from layout_in edits.
    """
    def shape_dict(s: ShapeSpec) -> dict:
        return {"x": s.x, "y": s.y, "w": s.w, "h": s.h, "color": s.color}

    def text_dict(t: TextSpec) -> dict:
        return {
            "x": t.x, "y": t.y, "w": t.w, "h": t.h, "color": t.color,
            "bold": t.bold, "align": t.align, "anchor": t.anchor, "margin": t.margin,
        }

    data = {}
    if comp.cover is not None:
        data["cover"] = {
            "shapes": [shape_dict(s) for s in comp.cover.shapes],
            "title": text_dict(comp.cover.title),
            "subtitle": text_dict(comp.cover.subtitle),
        }
    if comp.content_header is not None:
        data["content_header"] = {
            "shapes": [shape_dict(s) for s in comp.content_header.shapes],
            "title": text_dict(comp.content_header.title),
            "content_top": comp.content_header.content_top,
        }
    return data


def theme_to_json(theme: Theme) -> str:
    """Serialize a Theme to the documented JSON format (used by --dump-default-theme)."""
    data = {
        "name": "default",
        "description": "convert-pptx built-in default design",
        "colors": {k: _rgb_to_hex(getattr(theme, f)) for k, f in _THEME_COLOR_KEYS.items()},
        "fonts": {k: getattr(theme, f) for k, f in _THEME_FONT_KEYS.items()},
        "font_sizes_pt": {k: getattr(theme, f) for k, f in _THEME_FONT_SIZE_KEYS.items()},
        "layout_in": {k: getattr(theme, f) for k, f in _THEME_LAYOUT_KEYS.items()},
        "syntax_palette": {k: _rgb_to_hex(v) for k, v in theme.syntax_palette.items()},
    }
    if theme.composition is not None:
        data["composition"] = _composition_to_dict(theme.composition)
    return json.dumps(data, ensure_ascii=False, indent=2)


def _theme_section(data: dict, section: str, keys: dict, parse) -> dict:
    """Extract one JSON section into Theme field overrides. Unknown keys are errors."""
    overrides = {}
    if section not in data:
        return overrides
    raw = data[section]
    if not isinstance(raw, dict):
        # An explicit null is almost certainly a typo; require omission instead.
        raise ValueError(
            f"theme: '{section}' must be an object (omit the key entirely to use defaults)"
        )
    for key, value in raw.items():
        if key not in keys:
            raise ValueError(
                f"theme: unknown key '{section}.{key}' (allowed: {', '.join(sorted(keys))})"
            )
        try:
            overrides[keys[key]] = parse(value)
        except ValueError as e:
            raise ValueError(f"theme: '{section}.{key}': {e}") from e
    return overrides


def _parse_font_name(value) -> str:
    if not isinstance(value, str) or not value.strip():
        raise ValueError(f"expected a non-empty font name string, got {value!r}")
    return value.strip()


def _parse_positive_number(value) -> float:
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        raise ValueError(f"expected a positive finite number, got {value!r}")
    try:
        f = float(value)
    except OverflowError as e:  # arbitrarily large JSON integers
        raise ValueError(f"expected a positive finite number, got {value!r}") from e
    if not math.isfinite(f) or f <= 0:
        raise ValueError(f"expected a positive finite number, got {value!r}")
    return f


# ----- composition parsing (theme JSON -> Composition) -----

def _parse_nonneg_number(value) -> float:
    """0-or-more finite number. Composition x / y / margin legitimately use 0
    (bottom bands, left bars), so the positive-number parser must not be
    reused for them.
    """
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        raise ValueError(f"expected a non-negative finite number, got {value!r}")
    try:
        f = float(value)
    except OverflowError as e:
        raise ValueError(f"expected a non-negative finite number, got {value!r}") from e
    if not math.isfinite(f) or f < 0:
        raise ValueError(f"expected a non-negative finite number, got {value!r}")
    return f


def _parse_comp_dim(value, *, tokens) -> "float | str":
    """A positive finite number or one of the allowed size tokens."""
    if isinstance(value, str):
        if value in tokens:
            return value
        raise ValueError(
            f"expected a positive finite number or one of: {', '.join(tokens)}, got {value!r}"
        )
    return _parse_positive_number(value)


def _parse_comp_color(value) -> str:
    """Validate a composition color (token name or hex) WITHOUT resolving it.

    Tokens stay symbolic in the Composition so they track the Theme's current
    colors at draw time -- including --primary-color, which is applied to the
    Theme after load_theme() returns. Resolving here would freeze the color
    and silently ignore that CLI override.
    """
    if isinstance(value, str) and value in _THEME_COLOR_KEYS:
        return value
    try:
        _parse_hex_color(value)
    except ValueError as e:
        raise ValueError(
            f"expected a color token ({', '.join(sorted(_THEME_COLOR_KEYS))}) "
            f"or a hex string like '#RRGGBB', got {value!r}"
        ) from e
    return value


def _parse_comp_bool(value) -> bool:
    if not isinstance(value, bool):
        raise ValueError(f"expected true or false, got {value!r}")
    return value


def _parse_comp_enum(value, *, allowed) -> str:
    if value not in allowed:
        raise ValueError(f"expected one of: {', '.join(allowed)}, got {value!r}")
    return value


def _require_comp_object(raw, path: str, *, allowed, required) -> None:
    """Shared object-shape validation: dict type, unknown keys, required keys."""
    if not isinstance(raw, dict):
        raise ValueError(f"theme: '{path}' must be an object")
    for key in raw:
        if key not in allowed:
            raise ValueError(
                f"theme: unknown key '{path}.{key}' (allowed: {', '.join(sorted(allowed))})"
            )
    for req in required:
        if req not in raw:
            raise ValueError(f"theme: '{path}' is missing required key '{req}'")


def _parse_shape_spec(raw, path: str) -> ShapeSpec:
    parsers = {
        "x": _parse_nonneg_number,
        "y": _parse_nonneg_number,
        "w": lambda v: _parse_comp_dim(v, tokens=("full", "sym")),
        "h": lambda v: _parse_comp_dim(v, tokens=("full",)),
        "color": _parse_comp_color,
    }
    _require_comp_object(raw, path, allowed=parsers, required=parsers)
    values = {}
    for key, parse in parsers.items():
        try:
            values[key] = parse(raw[key])
        except ValueError as e:
            raise ValueError(f"theme: '{path}.{key}': {e}") from e
    return ShapeSpec(**values)


def _parse_shapes(raw, path: str) -> "tuple[ShapeSpec, ...]":
    """Parse a shapes array. An empty array is accepted and means the same as
    omitting the key: no decoration shapes."""
    if not isinstance(raw, list):
        raise ValueError(f"theme: '{path}' must be an array of rectangle objects")
    return tuple(_parse_shape_spec(item, f"{path}[{i}]") for i, item in enumerate(raw))


def _parse_text_spec(raw, path: str, *, default_bold: bool) -> TextSpec:
    parsers = {
        "x": _parse_nonneg_number,
        "y": _parse_nonneg_number,
        "w": lambda v: _parse_comp_dim(v, tokens=("sym",)),
        "h": _parse_positive_number,
        "color": _parse_comp_color,
        "bold": _parse_comp_bool,
        "align": lambda v: _parse_comp_enum(v, allowed=("left", "center", "right")),
        "anchor": lambda v: _parse_comp_enum(v, allowed=("top", "middle")),
        "margin": _parse_nonneg_number,
    }
    _require_comp_object(raw, path, allowed=parsers, required=("x", "y", "w", "h", "color"))
    values = {}
    for key, value in raw.items():
        try:
            values[key] = parsers[key](value)
        except ValueError as e:
            raise ValueError(f"theme: '{path}.{key}': {e}") from e
    values.setdefault("bold", default_bold)
    return TextSpec(**values)


def _parse_composition(raw) -> Composition:
    """Parse and validate the `composition` theme section.

    Each part (cover / content_header) replaces the default wholesale; there
    is no deep merge within a part. Colors are validated but kept symbolic
    (see _parse_comp_color).
    """
    if not isinstance(raw, dict):
        raise ValueError(
            "theme: 'composition' must be an object (omit the key entirely to use the default layout)"
        )
    for key in raw:
        if key not in ("cover", "content_header"):
            raise ValueError(
                f"theme: unknown key 'composition.{key}' (allowed: content_header, cover)"
            )
    if not raw:
        raise ValueError(
            "theme: 'composition' must define at least one of: cover, content_header "
            "(omit the key entirely to use the default layout)"
        )

    cover = None
    if "cover" in raw:
        c = raw["cover"]
        _require_comp_object(
            c, "composition.cover",
            allowed=("shapes", "title", "subtitle"), required=("title", "subtitle"),
        )
        cover = CoverComposition(
            shapes=_parse_shapes(c.get("shapes", []), "composition.cover.shapes"),
            title=_parse_text_spec(c["title"], "composition.cover.title", default_bold=True),
            subtitle=_parse_text_spec(c["subtitle"], "composition.cover.subtitle", default_bold=False),
        )

    content_header = None
    if "content_header" in raw:
        h = raw["content_header"]
        _require_comp_object(
            h, "composition.content_header",
            allowed=("shapes", "title", "content_top"), required=("title", "content_top"),
        )
        try:
            content_top = _parse_positive_number(h["content_top"])
        except ValueError as e:
            raise ValueError(f"theme: 'composition.content_header.content_top': {e}") from e
        content_header = ContentHeaderComposition(
            shapes=_parse_shapes(h.get("shapes", []), "composition.content_header.shapes"),
            title=_parse_text_spec(h["title"], "composition.content_header.title", default_bold=True),
            content_top=content_top,
        )

    return Composition(cover=cover, content_header=content_header)


def load_theme(path: Path) -> Theme:
    """Load a partial theme JSON and merge it over the default Theme.

    Raises ValueError with a human-readable message on any problem.
    """
    try:
        raw_text = path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError) as e:
        raise ValueError(f"theme: cannot read file {path} as UTF-8: {e}") from e
    try:
        data = json.loads(raw_text)
    except json.JSONDecodeError as e:
        raise ValueError(f"theme: invalid JSON in {path}: {e}") from e
    if not isinstance(data, dict):
        raise ValueError(f"theme: top level of {path} must be a JSON object")

    known_sections = {"colors", "fonts", "font_sizes_pt", "layout_in", "syntax_palette", "composition"}
    for key in data:
        if key not in known_sections and key not in _THEME_META_KEYS:
            raise ValueError(
                f"theme: unknown top-level key '{key}' "
                f"(allowed: {', '.join(sorted(known_sections | set(_THEME_META_KEYS)))})"
            )

    overrides = {}
    overrides.update(_theme_section(data, "colors", _THEME_COLOR_KEYS, _parse_hex_color))
    overrides.update(_theme_section(data, "fonts", _THEME_FONT_KEYS, _parse_font_name))
    overrides.update(_theme_section(data, "font_sizes_pt", _THEME_FONT_SIZE_KEYS, _parse_positive_number))
    overrides.update(_theme_section(data, "layout_in", _THEME_LAYOUT_KEYS, _parse_positive_number))

    theme = Theme(**overrides)

    if "syntax_palette" in data:
        palette_raw = data["syntax_palette"]
        if not isinstance(palette_raw, dict):
            raise ValueError(
                "theme: 'syntax_palette' must be an object (omit the key entirely to use defaults)"
            )
        palette = _default_syntax_palette()
        for key, value in palette_raw.items():
            if key not in palette:
                raise ValueError(
                    f"theme: unknown key 'syntax_palette.{key}' (allowed: {', '.join(sorted(palette))})"
                )
            try:
                palette[key] = _parse_hex_color(value)
            except ValueError as e:
                raise ValueError(f"theme: 'syntax_palette.{key}': {e}") from e
        theme = replace(theme, syntax_palette=palette)

    if "composition" in data:
        comp = _parse_composition(data["composition"])
        theme = replace(theme, composition=comp)
        # A custom content_header replaces the default band wholesale, so
        # layout_in.title_band_height stops mattering; warn (not error) when
        # both are spelled out to catch a likely misunderstanding. A
        # cover-only override keeps using it via the default content header,
        # so no warning in that case.
        layout_raw = data.get("layout_in")
        if (
            comp.content_header is not None
            and isinstance(layout_raw, dict)
            and "title_band_height" in layout_raw
        ):
            print(
                "Warning: this theme overrides composition.content_header, so "
                "layout_in.title_band_height is not used (it only affects the "
                "default composition)",
                file=sys.stderr,
            )

    return theme


# ===================== Markdown parsing =====================

@dataclass
class Block:
    """A Markdown block element."""
    kind: str
    text: str = ""
    level: int = 0
    lang: str = ""
    rows: List[List[str]] = field(default_factory=list)
    src: str = ""
    alt: str = ""


HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)(?:\s+#+)?\s*$")
IMAGE_RE = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$")
UL_RE = re.compile(r"^(\s*)[-*+]\s+(.+)$")
OL_RE = re.compile(r"^(\s*)\d+\.\s+(.+)$")
HR_RE = re.compile(r"^\s*(-{3,}|\*{3,}|_{3,})\s*$")
TABLE_DIVIDER_RE = re.compile(r"^\s*\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)*\|?\s*$")


def parse_markdown(md_text: str) -> List[Block]:
    """Parse the Markdown text into a flat list of Block objects.
    Minimal block-level parser: covers headings, paragraphs, bullet/numbered lists,
    fenced code (including mermaid), tables, images (standalone), and horizontal rules.
    """
    lines = md_text.splitlines()
    blocks: List[Block] = []
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]

        # Fenced code
        fence = re.match(r"^```(\w*)\s*$", line)
        if fence:
            lang = fence.group(1)
            i += 1
            buf = []
            while i < n and not re.match(r"^```\s*$", lines[i]):
                buf.append(lines[i])
                i += 1
            if i < n:
                i += 1  # consume closing fence
            kind = "mermaid" if lang == "mermaid" else "code"
            blocks.append(Block(kind=kind, text="\n".join(buf), lang=lang))
            continue

        # Heading
        m = HEADING_RE.match(line)
        if m:
            blocks.append(Block(kind="heading", text=m.group(2), level=len(m.group(1))))
            i += 1
            continue

        # Standalone image
        m = IMAGE_RE.match(line)
        if m:
            blocks.append(Block(kind="image", alt=m.group(1), src=m.group(2)))
            i += 1
            continue

        # Horizontal rule
        if HR_RE.match(line):
            blocks.append(Block(kind="hr"))
            i += 1
            continue

        # Table: header | divider | body rows
        if "|" in line and i + 1 < n and TABLE_DIVIDER_RE.match(lines[i + 1]):
            header = _split_table_row(line)
            i += 2
            rows = [header]
            while i < n and "|" in lines[i] and lines[i].strip():
                rows.append(_split_table_row(lines[i]))
                i += 1
            blocks.append(Block(kind="table", rows=rows))
            continue

        # Lists (group consecutive items)
        m_ul = UL_RE.match(line)
        m_ol = OL_RE.match(line)
        if m_ul or m_ol:
            kind = "ul" if m_ul else "ol"
            items = []
            while i < n:
                mm_ul = UL_RE.match(lines[i])
                mm_ol = OL_RE.match(lines[i])
                if mm_ul and kind == "ul":
                    indent = len(mm_ul.group(1))
                    items.append((indent // 2, mm_ul.group(2)))
                    i += 1
                elif mm_ol and kind == "ol":
                    indent = len(mm_ol.group(1))
                    items.append((indent // 2, mm_ol.group(2)))
                    i += 1
                else:
                    break
            blocks.append(Block(kind=kind, rows=[[str(lvl), txt] for lvl, txt in items]))
            continue

        # Blank line -> paragraph separator
        if not line.strip():
            i += 1
            continue

        # Paragraph (accumulate consecutive non-blank lines)
        para_lines = [line]
        i += 1
        while (
            i < n
            and lines[i].strip()
            and not HEADING_RE.match(lines[i])
            and not re.match(r"^```", lines[i])
            and not IMAGE_RE.match(lines[i])
            and not HR_RE.match(lines[i])
            and not UL_RE.match(lines[i])
            and not OL_RE.match(lines[i])
            and not ("|" in lines[i] and i + 1 < n and TABLE_DIVIDER_RE.match(lines[i + 1]))
        ):
            para_lines.append(lines[i])
            i += 1
        blocks.append(Block(kind="paragraph", text=" ".join(s.strip() for s in para_lines)))

    return blocks


def _split_table_row(line: str) -> List[str]:
    """Split a Markdown table row into cells."""
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


# ===================== Slide segmentation =====================

@dataclass
class SlideSpec:
    title: Optional[str]
    blocks: List[Block]


def split_into_slides(blocks: List[Block]) -> (Optional[str], Optional[str], List[SlideSpec]):
    """Produce (doc_title, doc_subtitle, slides) from a flat block list.
    - First H1 -> doc_title (becomes title slide)
    - The paragraph immediately after the first H1 (if any) -> subtitle
    - Each H2 starts a new slide; its text becomes the slide title
    - If there's no H2 at all, a single content slide holds everything after the H1
    """
    doc_title: Optional[str] = None
    doc_subtitle: Optional[str] = None
    slides: List[SlideSpec] = []

    idx = 0
    n = len(blocks)

    # Find first H1
    while idx < n and not (blocks[idx].kind == "heading" and blocks[idx].level == 1):
        idx += 1
    if idx < n:
        doc_title = blocks[idx].text
        idx += 1
        # Optional subtitle: the first paragraph before any H2
        j = idx
        while j < n:
            b = blocks[j]
            if b.kind == "heading" and b.level == 2:
                break
            if b.kind == "paragraph" and doc_subtitle is None:
                doc_subtitle = b.text
                break
            j += 1
    else:
        idx = 0  # no H1 at all; still try to build slides

    # Segment remaining blocks by H2
    current: Optional[SlideSpec] = None
    any_h2 = False
    for b in blocks[idx:]:
        if b.kind == "heading" and b.level == 2:
            any_h2 = True
            current = SlideSpec(title=b.text, blocks=[])
            slides.append(current)
            continue
        if current is None:
            # Before the first H2; skip the subtitle paragraph we already captured
            if b.kind == "paragraph" and b.text == doc_subtitle:
                continue
            current = SlideSpec(title=None, blocks=[])
            slides.append(current)
        current.blocks.append(b)

    if not any_h2 and not slides:
        # Rare case: only H1 (or nothing). Create an empty slide to avoid zero-slide decks.
        slides.append(SlideSpec(title=None, blocks=[]))

    return doc_title, doc_subtitle, slides


# ===================== Rendering helpers =====================

_HEX_COLOR_RE = re.compile(r"^#?[0-9A-Fa-f]{3}([0-9A-Fa-f]{3})?$")


def _parse_hex_color(hex_str) -> RGBColor:
    if not isinstance(hex_str, str) or not _HEX_COLOR_RE.match(hex_str):
        raise ValueError(
            f"invalid hex color: {hex_str!r} (expected '#RGB' or '#RRGGBB')"
        )
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_rgb(hex_str: str) -> RGBColor:
    """argparse type= wrapper: keeps the historical exit-code-2 behavior."""
    try:
        return _parse_hex_color(hex_str)
    except ValueError as e:
        raise argparse.ArgumentTypeError(str(e)) from e


def fetch_mermaid_png(code: str) -> Optional[bytes]:
    """Download a mermaid diagram as PNG via mermaid.ink. Returns None on failure.

    Validates Content-Type and PNG magic bytes; rejects redirects to private hosts.
    """
    import requests
    try:
        payload = {"code": code, "mermaid": {"theme": "default"}}
        encoded = base64.urlsafe_b64encode(json.dumps(payload).encode()).decode().rstrip("=")
        url = MERMAID_PNG_ENDPOINT.format(encoded)
        resp = requests.get(
            url, timeout=20, headers={"User-Agent": "convert-pptx/1.0"}, stream=True
        )
        ctype = resp.headers.get("Content-Type", "")
        if resp.status_code == 200 and "image/png" in ctype:
            content = _read_capped(resp, "mermaid.ink PNG")
            if content is not None and content.startswith(b"\x89PNG"):
                return content
        print(
            f"Warning: mermaid.ink returned status={resp.status_code} content-type={ctype!r}",
            file=sys.stderr,
        )
    except Exception as e:
        print(f"Warning: mermaid.ink fetch failed: {e}", file=sys.stderr)
    return None


def fit_size_inches(px_w: int, px_h: int, max_w_in: float, max_h_in: float) -> (Inches, Inches):
    """Fit a (pixel) size within (max_w, max_h) inches preserving aspect ratio."""
    # Assume 96 DPI for raster images coming from mermaid.ink
    w_in = px_w / 96
    h_in = px_h / 96
    if w_in <= 0 or h_in <= 0:
        return Inches(max_w_in), Inches(max_h_in)
    ratio = min(max_w_in / w_in, max_h_in / h_in, 1.0)
    return Inches(w_in * ratio), Inches(h_in * ratio)


def measure_image_bytes(data: bytes) -> (int, int):
    from PIL import Image
    with Image.open(io.BytesIO(data)) as img:
        return img.size


def _lexer_for(lang: str):
    """Return a pygments lexer for the given hint, falling back to plain text."""
    from pygments.lexers import get_lexer_by_name
    from pygments.lexers.special import TextLexer
    from pygments.util import ClassNotFound

    if lang:
        try:
            return get_lexer_by_name(lang, stripnl=False, ensurenl=False)
        except ClassNotFound:
            pass
    return TextLexer(stripnl=False, ensurenl=False)


def _token_color(tok_type, palette: "dict[str, RGBColor]"):
    """Resolve a pygments token to an RGB color by walking up the token hierarchy."""
    from pygments.token import Token

    mapping = [
        (Token.Keyword,             "keyword"),
        (Token.Name.Function,       "func"),
        (Token.Name.Class,          "func"),
        (Token.Name.Decorator,      "func"),
        (Token.Name.Builtin,        "builtin"),
        (Token.Name.Builtin.Pseudo, "builtin"),
        (Token.Name.Tag,            "tag"),
        (Token.Name.Attribute,      "attr"),
        (Token.Literal.String,      "string"),
        (Token.Literal.Number,      "number"),
        (Token.Operator,            "operator"),
        (Token.Punctuation,         "operator"),
        (Token.Comment,             "comment"),
        (Token.Error,               "error"),
        (Token.Generic.Heading,     "heading"),
        (Token.Generic.Subheading,  "heading"),
    ]
    t = tok_type
    while t is not None:
        for base, key in mapping:
            if t in base:
                return palette[key]
        parent = getattr(t, "parent", None)
        if parent is None or parent is t:
            break
        t = parent
    return None


def _lex_code_lines(code_text: str, lang: str):
    """Yield token-list-per-line for the given code string.
    Each item is a list of (token_type, text_fragment) tuples.  Empty lines
    yield an empty list so the paragraph structure matches the source.
    """
    from pygments import lex

    tokens = lex(code_text, _lexer_for(lang))
    line: list = []
    for tok, text in tokens:
        while "\n" in text:
            head, _, text = text.partition("\n")
            if head:
                line.append((tok, head))
            yield line
            line = []
        if text:
            line.append((tok, text))
    if line:
        yield line


NBSP = " "

def _normalize_leading(line_tokens):
    """Convert leading-whitespace-only runs to NBSP so PowerPoint preserves indentation.

    PowerPoint collapses/hides runs of ASCII spaces during layout, so we replace
    the spaces inside leading-whitespace-only runs with NBSP. Tabs are expanded
    to 4 spaces everywhere. Non-leading spaces are left as normal spaces so word
    wrap behaves naturally.
    """
    normalized = []
    indent_done = False
    for tok, text in line_tokens:
        if not indent_done and text and text.strip() == "":
            normalized.append((tok, text.replace("	", "    ").replace(" ", NBSP)))
        else:
            indent_done = True
            normalized.append((tok, text.replace("	", "    ")))
    return normalized
def _block_has_content(block: "Block") -> bool:
    """Return True if a block would render any visible shape on a slide."""
    if block.kind == "paragraph":
        return bool(block.text and block.text.strip())
    if block.kind in ("ul", "ol"):
        return bool(block.rows) and any((r[1] or "").strip() for r in block.rows)
    if block.kind == "heading":
        return bool(block.text and block.text.strip())
    if block.kind == "table":
        return bool(block.rows) and any(any((c or "").strip() for c in row) for row in block.rows)
    if block.kind == "code":
        return bool(block.text and block.text.strip())
    if block.kind == "mermaid":
        return bool(block.text and block.text.strip())
    if block.kind == "image":
        return bool(block.src)
    if block.kind == "hr":
        return False  # horizontal rule on its own does not justify a slide
    return False


def strip_inline_markdown(text: str) -> str:
    """Minimal inline cleanup for display in text frames."""
    # GFM task list markers -> unicode ballot boxes. Do this *before* we strip
    # links so the brackets are not confused with link syntax.
    text = re.sub(r"^\s*\[[ xX]\]\s+", lambda m: ("☑ " if "x" in m.group(0).lower() else "☐ "), text)
    # Links: [text](url) -> text
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
    # Inline code
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # Bold / italic markers (we keep content; runs are not split here)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    # Strikethrough
    text = re.sub(r"~~([^~]+)~~", r"\1", text)
    return text


# ===================== PPTX builder =====================

class Deck:
    def __init__(self, *, aspect: str, theme: Theme):
        self.theme = theme
        self.prs = Presentation()
        # Keep the inch sizes as literals (not derived from EMU) so that
        # composition math (e.g. "sym" widths) reproduces the historical
        # EMU-exact geometry without float round-trip drift.
        if aspect == "4:3":
            self.slide_w_in, self.slide_h_in = 10.0, 7.5
        else:  # 16:9
            self.slide_w_in, self.slide_h_in = 13.333, 7.5
        self.prs.slide_width = Inches(self.slide_w_in)
        self.prs.slide_height = Inches(self.slide_h_in)
        self.blank = self.prs.slide_layouts[6]
        # Effective composition, resolved per part: a theme may override
        # cover / content_header independently; a part left as None falls
        # back to the built-in default composition.
        default_comp = build_default_composition(theme, self.slide_w_in, self.slide_h_in)
        theme_comp = theme.composition or Composition()
        self.composition = Composition(
            cover=theme_comp.cover or default_comp.cover,
            content_header=theme_comp.content_header or default_comp.content_header,
        )

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self.blank)
        return slide

    # ----- composition-driven drawing -----

    def _resolve_comp_w(self, x: float, w) -> float:
        """Resolve a composition width (number | 'full' | 'sym') to inches."""
        if w == "full":
            return self.slide_w_in
        if w == "sym":
            return self.slide_w_in - x * 2
        return float(w)

    def _resolve_comp_h(self, h) -> float:
        """Resolve a composition height (number | 'full') to inches."""
        if h == "full":
            return self.slide_h_in
        return float(h)

    def _resolve_comp_color(self, color: str) -> RGBColor:
        """Resolve a composition color token / hex string to an RGBColor.

        Resolution happens at draw time (not at theme load) on purpose:
        --primary-color is applied to the Theme *after* load_theme(), so an
        early resolution would not see CLI overrides in "primary" tokens.
        """
        field_name = _THEME_COLOR_KEYS.get(color)
        if field_name is not None:
            return getattr(self.theme, field_name)
        return _parse_hex_color(color)

    def _draw_comp_shape(self, slide, spec: ShapeSpec) -> None:
        """Draw one composition rectangle (solid fill, no outline)."""
        w_in = self._resolve_comp_w(spec.x, spec.w)
        h_in = self._resolve_comp_h(spec.h)
        if w_in <= 0 or h_in <= 0:
            # "sym" can go negative for large x; the validator cannot know the
            # slide width, so this draw-time skip is the safety net.
            print(
                f"Warning: composition shape at x={spec.x}, y={spec.y} resolves to a "
                f"non-positive size ({w_in:.3g} x {h_in:.3g} in) and was skipped",
                file=sys.stderr,
            )
            return
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(spec.x), Inches(spec.y), Inches(w_in), Inches(h_in),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._resolve_comp_color(spec.color)
        shape.line.fill.background()

    def _draw_comp_text(
        self, slide, spec: TextSpec, text: str, font_name: str, size_pt: float,
    ) -> None:
        """Draw one composition text box (title / subtitle / header title)."""
        w_in = self._resolve_comp_w(spec.x, spec.w)
        if w_in <= 0 or spec.h <= 0:
            print(
                f"Warning: composition text at x={spec.x}, y={spec.y} resolves to a "
                f"non-positive size ({w_in:.3g} x {spec.h:.3g} in) and was skipped",
                file=sys.stderr,
            )
            return
        tb = slide.shapes.add_textbox(
            Inches(spec.x), Inches(spec.y), Inches(w_in), Inches(spec.h),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(spec.margin)
        tf.margin_right = Inches(spec.margin)
        if spec.anchor == "middle":
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        if spec.align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif spec.align == "right":
            p.alignment = PP_ALIGN.RIGHT
        p.font.name = font_name
        p.font.size = Pt(size_pt)
        p.font.bold = spec.bold
        p.font.color.rgb = self._resolve_comp_color(spec.color)

    def _add_content_header(self, slide, text: str):
        """Draw the content-slide header: shapes in array order, then the
        title text on top (z-order: decorations below, text frontmost)."""
        header = self.composition.content_header
        for spec in header.shapes:
            self._draw_comp_shape(slide, spec)
        self._draw_comp_text(
            slide, header.title, text,
            self.theme.heading_font, self.theme.title_band_size_pt,
        )

    def add_title_slide(self, title: str, subtitle: Optional[str]):
        slide = self._new_slide()
        cover = self.composition.cover
        for spec in cover.shapes:
            self._draw_comp_shape(slide, spec)
        self._draw_comp_text(
            slide, cover.title, title,
            self.theme.heading_font, self.theme.title_slide_title_size_pt,
        )
        if subtitle:
            self._draw_comp_text(
                slide, cover.subtitle, strip_inline_markdown(subtitle),
                self.theme.body_font, self.theme.title_slide_subtitle_size_pt,
            )

    def add_content_slide(
        self,
        spec: SlideSpec,
        base_dir: Path,
        max_body_chars: int,
    ):
        # Drop blocks that would render as nothing (empty paragraphs, blank hr
        # lists produced by upstream quirks) so we don't emit empty slides.
        nonempty_blocks = [b for b in spec.blocks if _block_has_content(b)]
        if not nonempty_blocks and not spec.title:
            # Truly empty slide -> skip entirely.
            return
        chunks = self._chunk_blocks(nonempty_blocks, max_body_chars)
        # Drop empty chunks defensively. _chunk_blocks returns [[]] for empty
        # input; a slide whose only purpose was the title band is still useful,
        # but a chunk with no blocks AND no title should not create a slide.
        chunks = [c for c in chunks if c or spec.title]
        if not chunks:
            return
        for idx, chunk in enumerate(chunks):
            slide = self._new_slide()
            if spec.title:
                title = spec.title if idx == 0 else f"{spec.title} ({idx + 1})"
                self._add_content_header(slide, title)
                top_cursor = Inches(self.composition.content_header.content_top)
            else:
                # Physical slides without a title (e.g. H2-less continuation
                # pages) keep the historical header-less layout.
                top_cursor = Inches(0.3)
            self._render_blocks(slide, chunk, top_cursor, base_dir)

    def _estimate_block_height_in(self, b: Block) -> float:
        """Estimate the rendered height of a block in inches.

        The numbers intentionally mirror the actual geometry used by the
        _render_* methods plus the 0.1in gap added after each block, so the
        chunker's budget stays in sync with what _render_blocks will actually
        consume. A small over-estimate is preferable to under-estimate because
        it avoids overflow cutoffs that silently drop shapes.
        """
        gap = 0.1
        if b.kind == "heading":
            return 0.45 + gap
        if b.kind == "paragraph":
            approx_lines = max(1, len(b.text) // 60 + 1)
            return 0.35 * approx_lines + gap
        if b.kind in ("ul", "ol"):
            n = max(1, len(b.rows))
            return 0.4 * n + 0.1 + gap
        if b.kind == "code":
            n_lines = max(1, b.text.count("\n") + 1)
            return max(0.5, 0.28 * n_lines + 0.2) + gap
        if b.kind == "table":
            n_rows = max(1, len(b.rows))
            return 0.4 * n_rows + gap
        if b.kind == "mermaid":
            return self.theme.mermaid_max_height_in + gap
        if b.kind == "image":
            return self.theme.image_max_height_in + gap
        if b.kind == "hr":
            return 0.15 + gap
        return 0.3 + gap

    def _chunk_blocks(self, blocks: List[Block], max_chars: int) -> List[List[Block]]:
        """Chunk the block list into slide-sized groups based on estimated height.

        The `max_chars` parameter is preserved for API compatibility but is now
        used only as a safety cap: the primary budget is the slide's remaining
        vertical space (slide_height - title band - top/bottom padding - a
        small safety margin). This keeps the chunker's splits aligned with the
        renderer's vertical layout, so shapes aren't silently dropped by the
        bottom-limit cutoff in _render_blocks.

        After packing, an orphan-heading pass moves trailing headings
        (e.g. `### 7.1 ...`) to the next chunk so a heading never sits alone at
        the tail of a slide while its content starts on the next one.
        """
        # Vertical budget per slide (inches), anchored on the composition's
        # content_top (where body blocks actually start). K derivation: the
        # historical budget was
        #   slide_height - title_band_height - content_padding - 0.6
        # and content started at title_band_height + 0.2, so rewriting the
        # budget in content_top terms gives
        #   slide_height - (content_top - 0.2) - content_padding - 0.6
        #     = slide_height - content_top - (content_padding + 0.4)
        # i.e. K = content_padding + 0.4 (NOT a fixed 0.9: content_padding is
        # documented to influence chunking, so it must keep scaling here).
        # Defaults (content_top=1.1, content_padding=0.5) yield the historical
        # 7.5 - 1.1 - 0.9 = 5.5in budget.
        budget_in = (
            self.slide_h_in
            - self.composition.content_header.content_top
            - (self.theme.content_padding_in + 0.4)
        )

        chunks: List[List[Block]] = []
        current: List[Block] = []
        height_used = 0.0
        chars_used = 0

        for b in blocks:
            h = self._estimate_block_height_in(b)
            # rows are [level, text] pairs for lists and cell lists for tables;
            # single-cell table rows have no x[1], fall back to the first cell.
            chars = max(1, len(b.text or "")) + sum(
                len((x[1] if len(x) > 1 else x[0]) or "") + 4 for x in (b.rows or [])
            )

            if h > budget_in:
                # A single block taller than one slide cannot be split by the
                # chunker; it will overflow past the slide edge when rendered.
                print(
                    f"Warning: a single '{b.kind}' block (estimated {h:.1f}in) exceeds "
                    f"the slide budget ({budget_in:.1f}in) and will overflow; "
                    "consider splitting it in the source Markdown",
                    file=sys.stderr,
                )

            would_overflow_height = current and (height_used + h > budget_in)
            would_overflow_chars = current and (chars_used + chars > max_chars)
            if would_overflow_height or would_overflow_chars:
                chunks.append(current)
                current = []
                height_used = 0.0
                chars_used = 0

            current.append(b)
            height_used += h
            chars_used += chars

        if current:
            chunks.append(current)

        # Orphan-heading fix: a heading at the very end of a chunk (with its
        # actual content spilling into the next chunk) reads as a dangling
        # title. Move it forward to sit with the content it introduces.
        for i in range(len(chunks) - 1):
            moved: List[Block] = []
            while chunks[i] and chunks[i][-1].kind == "heading":
                moved.append(chunks[i].pop())
            if moved:
                chunks[i + 1] = list(reversed(moved)) + chunks[i + 1]

        # Drop any chunks that became empty after the orphan migration.
        chunks = [c for c in chunks if c]
        if not chunks:
            chunks.append([])
        return chunks

    def _render_blocks(self, slide, blocks: List[Block], start_top: Emu, base_dir: Path):
        left = Inches(self.theme.content_padding_in)
        width = self.prs.slide_width - Inches(self.theme.content_padding_in * 2)
        bottom_limit = self.prs.slide_height - Inches(0.3)
        cursor = start_top
        skipped = 0

        for b in blocks:
            if cursor >= bottom_limit:
                skipped += 1
                continue

            if b.kind == "paragraph":
                height = self._render_text(
                    slide, left, cursor, width, strip_inline_markdown(b.text),
                    Pt(self.theme.body_size_pt),
                )
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "heading":
                # H3+ rendered as bold accent line
                size = max(14, self.theme.heading_h3_size_pt - (b.level - 3) * 2)
                height = self._render_text(
                    slide, left, cursor, width,
                    strip_inline_markdown(b.text),
                    Pt(size), bold=True, color=self.theme.primary,
                )
                cursor = cursor + height + Inches(0.1)

            elif b.kind in ("ul", "ol"):
                items = b.rows
                tb = slide.shapes.add_textbox(left, cursor, width, Inches(0.4 * len(items) + 0.1))
                tf = tb.text_frame
                tf.word_wrap = True
                for i, (lvl_str, txt) in enumerate(items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    bullet = "- " if b.kind == "ul" else f"{i + 1}. "
                    p.text = f"{'  ' * int(lvl_str)}{bullet}{strip_inline_markdown(txt)}"
                    p.font.name = self.theme.body_font
                    p.font.size = Pt(self.theme.list_size_pt)
                    p.font.color.rgb = self.theme.text
                cursor = cursor + tb.height + Inches(0.1)

            elif b.kind == "code":
                height = self._render_code(slide, left, cursor, width, b.text, b.lang)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "mermaid":
                height = self._render_mermaid(slide, left, cursor, width, b.text)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "image":
                height = self._render_image(slide, left, cursor, width, b.src, b.alt, base_dir)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "table":
                height = self._render_table(slide, left, cursor, width, b.rows)
                cursor = cursor + height + Inches(0.1)

            elif b.kind == "hr":
                ln = slide.shapes.add_connector(1, left, cursor, left + width, cursor)
                ln.line.color.rgb = self.theme.hr_color
                cursor = cursor + Inches(0.15)

        if skipped:
            print(
                f"Warning: {skipped} block(s) did not fit on a slide and were not rendered "
                "(check layout_in / font_sizes_pt theme values)",
                file=sys.stderr,
            )

    def _render_text(
        self, slide, left, top, width, text, size: Pt, *, bold=False,
        color: Optional[RGBColor] = None,
    ):
        if color is None:
            color = self.theme.text
        # Rough height estimate based on char count
        approx_lines = max(1, int(len(text) / 60) + 1)
        height = Inches(0.35 * approx_lines)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.theme.body_font
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        return height

    def _render_code(self, slide, left, top, width, code_text: str, lang: str = ""):
        """Render a fenced code block with Pygments syntax highlight.

        Each line becomes one paragraph; each Pygments token becomes a run so
        the whole block keeps a uniform monospace font while coloring keywords,
        strings, comments, etc.  Leading whitespace on each line is converted
        to NBSP so PowerPoint preserves indentation.
        """
        line_groups = list(_lex_code_lines(code_text, lang))
        if not line_groups:
            line_groups = [[]]

        height = Inches(max(0.5, 0.28 * len(line_groups) + 0.2))
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self.theme.code_bg
        box.line.color.rgb = self.theme.code_border
        tf = box.text_frame
        # Code should not reflow: it preserves the author's line breaks.
        tf.word_wrap = False
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)

        for i, line_tokens in enumerate(line_groups):
            line_tokens = _normalize_leading(line_tokens)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Baseline paragraph formatting (applies if a run doesn't override).
            p.font.name = self.theme.code_font
            p.font.size = Pt(self.theme.code_size_pt)
            p.font.color.rgb = self.theme.code_text

            if not line_tokens:
                # Preserve blank lines with a NBSP run so the line height stays.
                run = p.add_run()
                run.text = NBSP
                run.font.name = self.theme.code_font
                run.font.size = Pt(self.theme.code_size_pt)
                run.font.color.rgb = self.theme.code_text
                continue

            for tok, txt in line_tokens:
                if not txt:
                    continue
                run = p.add_run()
                run.text = txt
                run.font.name = self.theme.code_font
                run.font.size = Pt(self.theme.code_size_pt)
                color = _token_color(tok, self.theme.syntax_palette) or self.theme.code_text
                run.font.color.rgb = color
        return height

    def _render_mermaid(self, slide, left, top, width, code: str):
        png = fetch_mermaid_png(code)
        if not png:
            return self._render_code(slide, left, top, width, code)
        max_w = self.theme.mermaid_max_width_in
        max_h = self.theme.mermaid_max_height_in
        try:
            w, h = measure_image_bytes(png)
            img_w, img_h = fit_size_inches(w, h, max_w, max_h)
        except Exception as e:
            print(f"Warning: cannot measure mermaid PNG size ({e}); using max size", file=sys.stderr)
            img_w, img_h = Inches(max_w), Inches(max_h)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=img_w, height=img_h)
        return img_h

    def _render_image(self, slide, left, top, width, src: str, alt: str, base_dir: Path):
        data = self._load_image_bytes(src, base_dir)
        if data is None:
            return self._render_text(
                slide, left, top, width, f"[画像が見つかりません: {alt or src}]", Pt(13)
            )
        max_w = self.theme.image_max_width_in
        max_h = self.theme.image_max_height_in
        try:
            w, h = measure_image_bytes(data)
            img_w, img_h = fit_size_inches(w, h, max_w, max_h)
        except Exception as e:
            print(f"Warning: cannot measure image size for {src} ({e}); using max size", file=sys.stderr)
            img_w, img_h = Inches(max_w), Inches(max_h)
        slide.shapes.add_picture(io.BytesIO(data), left, top, width=img_w, height=img_h)
        return img_h

    @staticmethod
    def _is_within(parent: Path, child: Path) -> bool:
        try:
            parent_r = parent.resolve()
            child_r = child.resolve()
        except OSError:
            return False
        return child_r == parent_r or parent_r in child_r.parents

    @staticmethod
    def _is_public_host(host: str) -> bool:
        """Reject loopback / link-local / private hosts to mitigate SSRF.

        IP literals are checked directly. Hostnames are resolved and every
        resolved address must be public (fail-close on resolution failure).
        A TOCTOU DNS-rebinding window remains because requests re-resolves;
        full mitigation would require pinning the connection to the checked IP.
        """
        import ipaddress
        import socket
        if not host:
            return False

        def _blocked(ip) -> bool:
            return ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved

        try:
            return not _blocked(ipaddress.ip_address(host))
        except ValueError:
            pass
        # Hostname (not an IP literal): resolve and check every address.
        try:
            infos = socket.getaddrinfo(host, None)
        except OSError:
            return False
        addrs = {info[4][0].split("%")[0] for info in infos}
        if not addrs:
            return False
        for addr in addrs:
            try:
                if _blocked(ipaddress.ip_address(addr)):
                    return False
            except ValueError:
                return False
        return True

    @classmethod
    def _load_image_bytes(cls, src: str, base_dir: Path) -> Optional[bytes]:
        parsed = urlparse(src)
        if parsed.scheme in ("http", "https"):
            if not cls._is_public_host(parsed.hostname or ""):
                print(f"Warning: image URL host blocked (SSRF guard): {parsed.hostname}", file=sys.stderr)
                return None
            try:
                import requests
                r = requests.get(
                    src,
                    timeout=20,
                    headers={"User-Agent": "convert-pptx/1.0"},
                    allow_redirects=False,
                    stream=True,
                )
                if r.status_code == 200:
                    return _read_capped(r, src)
            except Exception:
                return None
            return None
        # local path: reject absolute/UNC references BEFORE any filesystem access
        # (a UNC src would otherwise trigger an SMB lookup during .is_file()).
        if _is_external_path_ref(src):
            print(f"Warning: absolute/UNC image path rejected: {src}", file=sys.stderr)
            return None
        # must resolve inside base_dir to prevent traversal
        candidate = base_dir / src
        if candidate.is_file() and cls._is_within(base_dir, candidate):
            try:
                return candidate.read_bytes()
            except OSError as e:
                print(f"Warning: cannot read image {candidate}: {e}", file=sys.stderr)
                return None
        return None

    def _render_table(self, slide, left, top, width, rows: List[List[str]]):
        if not rows:
            return Inches(0.3)
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        height = Inches(0.4 * n_rows)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table
        for r, row in enumerate(rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.text = strip_inline_markdown(row[c]) if c < len(row) else ""
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.theme.body_font
                        run.font.size = Pt(self.theme.table_size_pt)
                        if r == 0:
                            run.font.bold = True
                            run.font.color.rgb = self.theme.on_primary
                        else:
                            run.font.color.rgb = self.theme.text
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme.primary
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        self.theme.table_row_odd if r % 2 == 1 else self.theme.table_row_even
                    )
        return height

    def save(self, path: Path):
        self.prs.save(str(path))


# ===================== Entry point =====================

def main() -> None:
    parser = argparse.ArgumentParser(description="Convert Markdown to PowerPoint (PPTX)")
    parser.add_argument("input", nargs="?", help="Input .md file path")
    parser.add_argument("output", nargs="?", help="Output .pptx file path")
    parser.add_argument("--title", help="Title slide heading (default: first H1 in the Markdown)")
    parser.add_argument("--subtitle", help="Title slide subtitle")
    parser.add_argument("--aspect", default="16:9", choices=["16:9", "4:3"])
    parser.add_argument("--primary-color", type=hex_to_rgb, default=None,
                        help="Primary color #RRGGBB (overrides the theme's primary)")
    def _positive_int(value: str) -> int:
        n = int(value)
        if n <= 0:
            raise argparse.ArgumentTypeError(f"expected a positive integer, got {value!r}")
        return n

    parser.add_argument("--max-body-chars", type=_positive_int, default=2400,
                        help="Soft budget for a single slide's body text")
    parser.add_argument("--theme", default=None,
                        help="Path to a theme JSON file (partial override of the default design)")
    parser.add_argument("--dump-default-theme", action="store_true",
                        help="Print the built-in default theme as JSON and exit")
    args = parser.parse_args()

    if args.dump_default_theme:
        print(theme_to_json(Theme()))
        return

    if not args.input:
        parser.error("the following arguments are required: input")

    try:
        theme = load_theme(Path(args.theme)) if args.theme else Theme()
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
    if args.primary_color is not None:
        theme = replace(theme, primary=args.primary_color)

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output) if args.output else input_path.with_suffix(".pptx")
    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    md_text = input_path.read_text(encoding="utf-8")
    blocks = parse_markdown(md_text)
    doc_title, doc_subtitle, slides = split_into_slides(blocks)

    if args.title:
        doc_title = args.title
    if args.subtitle:
        doc_subtitle = args.subtitle

    deck = Deck(aspect=args.aspect, theme=theme)
    if doc_title:
        deck.add_title_slide(doc_title, doc_subtitle)
    for spec in slides:
        deck.add_content_slide(spec, base_dir=input_path.parent, max_body_chars=args.max_body_chars)

    deck.save(output_path)
    print(f"Generated: {output_path}")


if __name__ == "__main__":
    main()
